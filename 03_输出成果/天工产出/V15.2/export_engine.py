"""
产物导出引擎 - 多格式导出模块
V15.2 天工产出

支持格式：Markdown、Word、PPT、Excel
"""

from typing import Optional, Dict, Any, List
from pathlib import Path
import os


class ExportError(Exception):
    """导出错误基类"""
    pass


class StyleConfig:
    """样式配置类"""
    
    def __init__(
        self,
        font_name: str = "微软雅黑",
        font_size: int = 12,
        title_font_size: int = 18,
        heading_font_size: int = 14,
        bold_headings: bool = True,
        line_spacing: float = 1.5,
        margin_top: float = 2.54,
        margin_bottom: float = 2.54,
        margin_left: float = 3.17,
        margin_right: float = 3.17,
    ):
        """
        初始化样式配置
        
        Args:
            font_name: 字体名称
            font_size: 正文字体大小（磅）
            title_font_size: 标题字体大小（磅）
            heading_font_size: 小标题字体大小（磅）
            bold_headings: 标题是否加粗
            line_spacing: 行距倍数
            margin_top: 上边距（厘米）
            margin_bottom: 下边距（厘米）
            margin_left: 左边距（厘米）
            margin_right: 右边距（厘米）
        """
        self.font_name = font_name
        self.font_size = font_size
        self.title_font_size = title_font_size
        self.heading_font_size = heading_font_size
        self.bold_headings = bold_headings
        self.line_spacing = line_spacing
        self.margin_top = margin_top
        self.margin_bottom = margin_bottom
        self.margin_left = margin_left
        self.margin_right = margin_right


class ProductExporter:
    """产物导出器
    
    支持将内容导出为多种格式：
    - Markdown (.md)
    - Word文档 (.docx)
    - PowerPoint演示文稿 (.pptx)
    - Excel表格 (.xlsx)
    
    使用示例：
        exporter = ProductExporter()
        exporter.to_markdown(content, "output.md")
        exporter.to_docx(content, "output.docx", style=StyleConfig(font_size=14))
    """
    
    def __init__(self, default_style: Optional[StyleConfig] = None):
        """
        初始化导出器
        
        Args:
            default_style: 默认样式配置，为None时使用默认值
        """
        self.default_style = default_style or StyleConfig()
    
    def _ensure_directory(self, output_path: str) -> None:
        """确保输出目录存在"""
        dir_path = os.path.dirname(output_path)
        if dir_path:
            os.makedirs(dir_path, exist_ok=True)
    
    def _validate_content(self, content: str) -> None:
        """验证内容有效性"""
        if content is None:
            raise ExportError("内容不能为None")
        if not isinstance(content, str):
            raise ExportError(f"内容类型错误，期望str，实际为{type(content).__name__}")
    
    def to_markdown(
        self,
        content: str,
        output_path: str,
        encoding: str = "utf-8",
        **kwargs
    ) -> str:
        """
        导出为Markdown文件
        
        Args:
            content: 要导出的内容
            output_path: 输出文件路径
            encoding: 文件编码，默认utf-8
            
        Returns:
            导出文件的绝对路径
            
        Raises:
            ExportError: 导出失败时抛出
        """
        try:
            self._validate_content(content)
            self._ensure_directory(output_path)
            
            with open(output_path, "w", encoding=encoding) as f:
                f.write(content)
            
            return os.path.abspath(output_path)
            
        except IOError as e:
            raise ExportError(f"Markdown导出失败: {e}")
        except Exception as e:
            raise ExportError(f"未知错误: {e}")
    
    def to_docx(
        self,
        content: str,
        output_path: str,
        style: Optional[StyleConfig] = None,
        **kwargs
    ) -> str:
        """
        导出为Word文档（使用python-docx）
        
        Args:
            content: 要导出的内容
            output_path: 输出文件路径
            style: 样式配置，为None时使用默认样式
            
        Returns:
            导出文件的绝对路径
            
        Raises:
            ExportError: 导出失败时抛出
        """
        try:
            from docx import Document
            from docx.shared import Pt, Cm
            from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
        except ImportError:
            raise ExportError("缺少依赖库 python-docx，请执行: pip install python-docx")
        
        try:
            self._validate_content(content)
            self._ensure_directory(output_path)
            
            config = style or self.default_style
            doc = Document()
            
            # 设置页边距
            sections = doc.sections
            for section in sections:
                section.top_margin = Cm(config.margin_top)
                section.bottom_margin = Cm(config.margin_bottom)
                section.left_margin = Cm(config.margin_left)
                section.right_margin = Cm(config.margin_right)
            
            # 解析并添加内容
            paragraphs = self._parse_content_to_paragraphs(content)
            
            for para_info in paragraphs:
                para = doc.add_paragraph()
                
                if para_info["type"] == "title":
                    run = para.add_run(para_info["text"])
                    run.font.name = config.font_name
                    run.font.size = Pt(config.title_font_size)
                    run.bold = config.bold_headings
                    para.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
                    
                elif para_info["type"] == "heading":
                    run = para.add_run(para_info["text"])
                    run.font.name = config.font_name
                    run.font.size = Pt(config.heading_font_size)
                    run.bold = config.bold_headings
                    
                else:
                    run = para.add_run(para_info["text"])
                    run.font.name = config.font_name
                    run.font.size = Pt(config.font_size)
                
                # 设置行距
                para.paragraph_format.line_spacing = config.line_spacing
            
            doc.save(output_path)
            return os.path.abspath(output_path)
            
        except Exception as e:
            raise ExportError(f"Word文档导出失败: {e}")
    
    def to_pptx(
        self,
        content: str,
        output_path: str,
        style: Optional[StyleConfig] = None,
        title: Optional[str] = None,
        **kwargs
    ) -> str:
        """
        导出为PPT（使用python-pptx）
        
        Args:
            content: 要导出的内容
            output_path: 输出文件路径
            style: 样式配置
            title: 演示文稿标题，用于首页
            
        Returns:
            导出文件的绝对路径
            
        Raises:
            ExportError: 导出失败时抛出
        """
        try:
            from pptx import Presentation
            from pptx.util import Pt, Inches
            from pptx.enum.text import PP_ALIGN
        except ImportError:
            raise ExportError("缺少依赖库 python-pptx，请执行: pip install python-pptx")
        
        try:
            self._validate_content(content)
            self._ensure_directory(output_path)
            
            config = style or self.default_style
            prs = Presentation()
            
            # 设置幻灯片尺寸
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
            
            # 添加标题页
            if title:
                slide_layout = prs.slide_layouts[6]  # 空白布局
                slide = prs.slides.add_slide(slide_layout)
                
                left = Inches(0.5)
                top = Inches(2.5)
                width = Inches(12.333)
                height = Inches(2)
                
                title_box = slide.shapes.add_textbox(left, top, width, height)
                tf = title_box.text_frame
                p = tf.paragraphs[0]
                p.text = title
                p.font.name = config.font_name
                p.font.size = Pt(config.title_font_size)
                p.font.bold = config.bold_headings
                p.alignment = PP_ALIGN.CENTER
            
            # 解析内容并创建幻灯片
            slides_content = self._parse_content_to_slides(content)
            
            for slide_content in slides_content:
                slide_layout = prs.slide_layouts[6]  # 空白布局
                slide = prs.slides.add_slide(slide_layout)
                
                # 标题文本框
                if slide_content["title"]:
                    title_box = slide.shapes.add_textbox(
                        Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8)
                    )
                    tf = title_box.text_frame
                    p = tf.paragraphs[0]
                    p.text = slide_content["title"]
                    p.font.name = config.font_name
                    p.font.size = Pt(config.heading_font_size)
                    p.font.bold = config.bold_headings
                
                # 内容文本框
                content_box = slide.shapes.add_textbox(
                    Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5)
                )
                tf = content_box.text_frame
                tf.word_wrap = True
                
                for i, line in enumerate(slide_content["content"]):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    
                    p.text = line
                    p.font.name = config.font_name
                    p.font.size = Pt(config.font_size)
                    p.level = 0
            
            prs.save(output_path)
            return os.path.abspath(output_path)
            
        except Exception as e:
            raise ExportError(f"PPT导出失败: {e}")
    
    def to_xlsx(
        self,
        content: str,
        output_path: str,
        style: Optional[StyleConfig] = None,
        sheet_name: str = "Sheet1",
        **kwargs
    ) -> str:
        """
        导出为Excel（使用openpyxl）
        
        Args:
            content: 要导出的内容，支持以下格式：
                - 纯文本：按行分割
                - CSV格式：逗号分隔
                - 表格格式：使用制表符或逗号分隔列
            output_path: 输出文件路径
            style: 样式配置
            sheet_name: 工作表名称
            
        Returns:
            导出文件的绝对路径
            
        Raises:
            ExportError: 导出失败时抛出
        """
        try:
            from openpyxl import Workbook
            from openpyxl.styles import Font, Alignment, Border, Side, PatternFill
        except ImportError:
            raise ExportError("缺少依赖库 openpyxl，请执行: pip install openpyxl")
        
        try:
            self._validate_content(content)
            self._ensure_directory(output_path)
            
            config = style or self.default_style
            wb = Workbook()
            ws = wb.active
            ws.title = sheet_name
            
            # 定义样式
            header_font = Font(
                name=config.font_name,
                size=config.font_size,
                bold=config.bold_headings
            )
            content_font = Font(
                name=config.font_name,
                size=config.font_size
            )
            center_alignment = Alignment(
                horizontal="center",
                vertical="center",
                wrap_text=True
            )
            thin_border = Border(
                left=Side(style="thin"),
                right=Side(style="thin"),
                top=Side(style="thin"),
                bottom=Side(style="thin")
            )
            header_fill = PatternFill(
                start_color="E0E0E0",
                end_color="E0E0E0",
                fill_type="solid"
            )
            
            # 解析内容
            rows = self._parse_content_to_rows(content)
            
            for row_idx, row_data in enumerate(rows, start=1):
                for col_idx, cell_value in enumerate(row_data, start=1):
                    cell = ws.cell(row=row_idx, column=col_idx, value=cell_value)
                    
                    # 第一行作为表头
                    if row_idx == 1:
                        cell.font = header_font
                        cell.fill = header_fill
                    else:
                        cell.font = content_font
                    
                    cell.alignment = center_alignment
                    cell.border = thin_border
            
            # 自动调整列宽
            for col in ws.columns:
                max_length = 0
                column = col[0].column_letter
                
                for cell in col:
                    if cell.value:
                        cell_length = len(str(cell.value))
                        # 中文字符宽度修正
                        cell_length = sum(2 if ord(c) > 127 else 1 for c in str(cell.value))
                        max_length = max(max_length, cell_length)
                
                adjusted_width = min(max_length + 2, 50)
                ws.column_dimensions[column].width = adjusted_width
            
            wb.save(output_path)
            return os.path.abspath(output_path)
            
        except Exception as e:
            raise ExportError(f"Excel导出失败: {e}")
    
    def _parse_content_to_paragraphs(self, content: str) -> List[Dict[str, Any]]:
        """
        解析内容为段落列表
        
        识别标题（#开头）和小标题（##开头）
        """
        paragraphs = []
        lines = content.split("\n")
        
        for line in lines:
            line = line.rstrip()
            
            if line.startswith("### "):
                paragraphs.append({
                    "type": "heading",
                    "text": line[4:],
                    "level": 3
                })
            elif line.startswith("## "):
                paragraphs.append({
                    "type": "heading",
                    "text": line[3:],
                    "level": 2
                })
            elif line.startswith("# "):
                paragraphs.append({
                    "type": "title",
                    "text": line[2:],
                    "level": 1
                })
            elif line.strip():
                paragraphs.append({
                    "type": "paragraph",
                    "text": line
                })
        
        return paragraphs
    
    def _parse_content_to_slides(self, content: str) -> List[Dict[str, Any]]:
        """
        解析内容为幻灯片列表
        
        使用 --- 或 === 作为幻灯片分隔符
        """
        slides = []
        current_slide = {"title": "", "content": []}
        
        lines = content.split("\n")
        
        for line in lines:
            stripped = line.strip()
            
            # 幻灯片分隔符
            if stripped in ["---", "===", "***"]:
                if current_slide["content"] or current_slide["title"]:
                    slides.append(current_slide)
                current_slide = {"title": "", "content": []}
                continue
            
            # 标题
            if stripped.startswith("# "):
                current_slide["title"] = stripped[2:]
                continue
            
            # 内容
            if stripped:
                current_slide["content"].append(stripped)
        
        # 添加最后一张幻灯片
        if current_slide["content"] or current_slide["title"]:
            slides.append(current_slide)
        
        return slides if slides else [{"title": "内容", "content": [content]}]
    
    def _parse_content_to_rows(self, content: str) -> List[List[str]]:
        """
        解析内容为表格行
        
        支持制表符分隔、逗号分隔和按行分割
        """
        rows = []
        lines = content.strip().split("\n")
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
            
            # 尝试检测分隔符
            if "\t" in line:
                cells = line.split("\t")
            elif "," in line and "|" not in line:
                cells = [cell.strip() for cell in line.split(",")]
            elif "|" in line:
                # Markdown表格格式
                cells = [cell.strip() for cell in line.split("|") if cell.strip()]
                # 跳过分隔行（如 |---|---|）
                if all(set(cell) <= {"-", ":"} for cell in cells):
                    continue
            else:
                cells = [line]
            
            rows.append(cells)
        
        return rows if rows else [[content]]


# 便捷函数
def export_to_markdown(content: str, output_path: str, **kwargs) -> str:
    """便捷函数：导出为Markdown"""
    exporter = ProductExporter()
    return exporter.to_markdown(content, output_path, **kwargs)


def export_to_docx(content: str, output_path: str, style: Optional[StyleConfig] = None, **kwargs) -> str:
    """便捷函数：导出为Word文档"""
    exporter = ProductExporter()
    return exporter.to_docx(content, output_path, style=style, **kwargs)


def export_to_pptx(content: str, output_path: str, style: Optional[StyleConfig] = None, **kwargs) -> str:
    """便捷函数：导出为PPT"""
    exporter = ProductExporter()
    return exporter.to_pptx(content, output_path, style=style, **kwargs)


def export_to_xlsx(content: str, output_path: str, style: Optional[StyleConfig] = None, **kwargs) -> str:
    """便捷函数：导出为Excel"""
    exporter = ProductExporter()
    return exporter.to_xlsx(content, output_path, style=style, **kwargs)


if __name__ == "__main__":
    # 测试代码
    test_content = """# 项目报告

## 一、项目概述
这是项目概述内容，介绍项目背景和目标。

## 二、主要内容
1. 第一点内容
2. 第二点内容
3. 第三点内容

### 详细说明
这里是详细的说明文字。

---
第二页内容
这是第二页的详细内容。
"""
    
    exporter = ProductExporter()
    
    # 测试Markdown导出
    md_path = exporter.to_markdown(test_content, "/tmp/test.md")
    print(f"Markdown导出成功: {md_path}")
    
    # 测试Word导出
    docx_path = exporter.to_docx(test_content, "/tmp/test.docx")
    print(f"Word导出成功: {docx_path}")
    
    # 测试PPT导出
    pptx_path = exporter.to_pptx(test_content, "/tmp/test.pptx", title="项目报告")
    print(f"PPT导出成功: {pptx_path}")
    
    # 测试Excel导出
    table_content = """项目,状态,负责人,进度
V15.2开发,进行中,天工,80%
导出功能,待测试,天工,50%
文档编写,未开始,天工,0%"""
    
    xlsx_path = exporter.to_xlsx(table_content, "/tmp/test.xlsx")
    print(f"Excel导出成功: {xlsx_path}")

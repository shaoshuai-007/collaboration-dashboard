#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
湖北电信AI智能配案 - 麦肯锡风格PPT
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

# 配色方案
COLOR_PRIMARY = RGBColor(26, 54, 93)      # 深蓝灰 #1A365D
COLOR_GRAY = RGBColor(74, 85, 104)        # 灰色 #4A5568
COLOR_ACCENT = RGBColor(49, 130, 206)     # 强调蓝 #3182CE
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_LIGHT_GRAY = RGBColor(240, 242, 245)

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 定义幻灯片布局
    title_slide_layout = prs.slide_layouts[6]  # 空白布局
    
    return prs

def add_title_slide(prs, title, subtitle):
    """封面页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景色块 - 左侧深蓝灰矩形
    left_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(6), Inches(7.5)
    )
    left_shape.fill.solid()
    left_shape.fill.fore_color.rgb = COLOR_PRIMARY
    left_shape.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(5), Inches(1.5))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = COLOR_WHITE
    
    # 副标题
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(5), Inches(1))
    sub_frame = sub_box.text_frame
    sub_para = sub_frame.paragraphs[0]
    sub_para.text = subtitle
    sub_para.font.size = Pt(20)
    sub_para.font.color.rgb = COLOR_ACCENT
    
    # 右侧装饰 - 强调蓝圆形
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(8), Inches(1), Inches(2.5), Inches(2.5)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = COLOR_ACCENT
    circle.line.fill.background()
    
    # 右侧文字
    right_text = slide.shapes.add_textbox(Inches(7.5), Inches(4), Inches(5), Inches(2))
    tf = right_text.text_frame
    p = tf.paragraphs[0]
    p.text = "九星智囊团"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    
    p2 = tf.add_paragraph()
    p2.text = "诚邀合作"
    p2.font.size = Pt(18)
    p2.font.color.rgb = COLOR_GRAY

def add_content_slide(prs, title, content_left="", content_right="", chart_data=None, chart_type=None):
    """通用内容页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 顶部标题区域背景
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = COLOR_PRIMARY
    header.line.fill.background()
    
    # 标题文字
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    
    # 左侧内容
    if content_left:
        left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(5.5), Inches(5))
        tf = left_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = content_left
        p.font.size = Pt(16)
        p.font.color.rgb = COLOR_GRAY
        p.line_spacing = 1.5
    
    # 右侧内容或图表
    if chart_data:
        # 图表区域
        chart_box = slide.shapes.add_chart(
            chart_type, Inches(6), Inches(1.6), Inches(6.5), Inches(5)
        )
        chart = chart_box.chart
        chart.chart_data = chart_data
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.plots[0].has_data_labels = True
    elif content_right:
        right_box = slide.shapes.add_textbox(Inches(6.5), Inches(1.6), Inches(6), Inches(5))
        tf = right_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = content_right
        p.font.size = Pt(16)
        p.font.color.rgb = COLOR_GRAY

def add_chart_slide(prs, title, chart_data, chart_type=XL_CHART_TYPE.COLUMN_CLUSTERED):
    """纯图表页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 顶部强调蓝条
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), Inches(7.5)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = COLOR_ACCENT
    header.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    
    # 图表
    chart_box = slide.shapes.add_chart(
        chart_type, Inches(1), Inches(1.5), Inches(11), Inches(5), chart_data
    )
    chart.has_legend = False
    chart.plots[0].has_data_labels = True

def add_three_column_slide(prs, title, items):
    """三列内容页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 顶部标题背景
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = COLOR_PRIMARY
    header.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    
    # 三列
    for i, item in enumerate(items):
        left = 0.5 + i * 4.3
        # 卡片背景
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(1.8), Inches(4), Inches(5)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_LIGHT_GRAY
        card.line.color.rgb = COLOR_ACCENT
        
        # 序号
        num_box = slide.shapes.add_textbox(Inches(left + 0.3), Inches(2.1), Inches(0.6), Inches(0.6))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = str(i + 1)
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = COLOR_ACCENT
        
        # 标题
        title_text = slide.shapes.add_textbox(Inches(left + 0.3), Inches(2.7), Inches(3.4), Inches(0.8))
        tf = title_text.text_frame
        p = tf.paragraphs[0]
        p.text = item['title']
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY
        
        # 内容
        content_text = slide.shapes.add_textbox(Inches(left + 0.3), Inches(3.5), Inches(3.4), Inches(3))
        tf = content_text.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = item['content']
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_GRAY

def add_six_agents_slide(prs):
    """6大智能体介绍"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = COLOR_PRIMARY
    header.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "6大智能体矩阵"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    
    agents = [
        {"name": "客群智脑", "desc": "用户画像\n精准分群"},
        {"name": "需求探针", "desc": "需求洞察\n预测分析"},
        {"name": "方案生成", "desc": "智能配案\n自动推荐"},
        {"name": "谈判助手", "desc": "话术优化\n实时辅导"},
        {"name": "服务管家", "desc": "全生命周期\n管理"},
        {"name": "决策智眼", "desc": "数据驾驶舱\n决策支持"}
    ]
    
    for i, agent in enumerate(agents):
        row = i // 3
        col = i % 3
        left = 0.8 + col * 4.3
        top = 1.6 + row * 2.8
        
        # 卡片
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(4), Inches(2.4)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_WHITE
        card.line.color.rgb = COLOR_ACCENT
        card.line.width = Pt(2)
        
        # 圆形图标
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(left + 1.5), Inches(top + 0.3), Inches(1), Inches(1)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = COLOR_ACCENT
        circle.line.fill.background()
        
        # 名称
        name_box = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 1.5), Inches(3.6), Inches(0.5))
        tf = name_box.text_frame
        p = tf.paragraphs[0]
        p.text = agent['name']
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY
        p.alignment = PP_ALIGN.CENTER
        
        # 描述
        desc_box = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 2), Inches(3.6), Inches(0.4))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = agent['desc']
        p.font.size = Pt(12)
        p.font.color.rgb = COLOR_GRAY
        p.alignment = PP_ALIGN.CENTER

def add_roadmap_slide(prs):
    """实施路线图"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = COLOR_PRIMARY
    header.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "实施路线图"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    
    phases = [
        {"phase": "M1-M2", "title": "基础建设", "items": "数据治理\n系统对接\n智能体定制"},
        {"phase": "M3-M4", "title": "试点运行", "items": "孝感/宜昌试点\n效果验证\n优化迭代"},
        {"phase": "M5-M6", "title": "规模推广", "items": "全省覆盖\n培训落地\n运营固化"}
    ]
    
    # 时间轴线
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(4.3), Inches(10), Inches(0.08)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_ACCENT
    line.line.fill.background()
    
    for i, phase in enumerate(phases):
        left = 0.8 + i * 4.3
        
        # 节点圆圈
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(left + 1.7), Inches(4.1), Inches(0.5), Inches(0.5)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = COLOR_ACCENT
        circle.line.fill.background()
        
        # 阶段标签
        phase_box = slide.shapes.add_textbox(Inches(left), Inches(1.8), Inches(4), Inches(0.5))
        tf = phase_box.text_frame
        p = tf.paragraphs[0]
        p.text = phase['phase']
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY
        p.alignment = PP_ALIGN.CENTER
        
        # 标题
        title_box = slide.shapes.add_textbox(Inches(left), Inches(2.3), Inches(4), Inches(0.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = phase['title']
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLOR_ACCENT
        p.alignment = PP_ALIGN.CENTER
        
        # 内容
        content_box = slide.shapes.add_textbox(Inches(left), Inches(2.9), Inches(4), Inches(1.2))
        tf = content_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = phase['items']
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_GRAY
        p.alignment = PP_ALIGN.CENTER

def add_team_slide(prs):
    """团队页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = COLOR_PRIMARY
    header.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "九星智囊团"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    
    # 九宫格团队展示
    team = [
        ("CEO", "战略总控"),
        ("CTO", "技术架构"),
        ("CDO", "数据治理"),
        ("产品VP", "需求规划"),
        ("架构师", "系统设计"),
        ("AI科学家", "模型训练"),
        ("运营总监", "效果运营"),
        ("客户成功", "服务保障"),
        ("财务总监", "投控管理")
    ]
    
    for i, (role, desc) in enumerate(team):
        row = i // 3
        col = i % 3
        left = 0.8 + col * 4.3
        top = 1.6 + row * 1.9
        
        # 卡片
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(4), Inches(1.6)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_WHITE
        card.line.color.rgb = COLOR_GRAY
        
        # 头像占位
        avatar = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(left + 0.3), Inches(top + 0.3), Inches(1), Inches(1)
        )
        avatar.fill.solid()
        avatar.fill.fore_color.rgb = COLOR_LIGHT_GRAY
        avatar.line.color.rgb = COLOR_ACCENT
        
        # 角色
        role_box = slide.shapes.add_textbox(Inches(left + 1.5), Inches(top + 0.4), Inches(2.2), Inches(0.5))
        tf = role_box.text_frame
        p = tf.paragraphs[0]
        p.text = role
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY
        
        # 描述
        desc_box = slide.shapes.add_textbox(Inches(left + 1.5), Inches(top + 0.9), Inches(2.2), Inches(0.4))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(12)
        p.font.color.rgb = COLOR_GRAY

def add_pricing_slide(prs):
    """报价页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = COLOR_PRIMARY
    header.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "项目预算"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    
    # 预算表格
    rows = 5
    cols = 3
    left = Inches(1.5)
    top = Inches(1.8)
    width = Inches(10)
    height = Inches(0.6)
    
    # 表头
    headers = ["项目", "内容", "金额(万元)"]
    for i, h in enumerate(headers):
        cell = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + i*width/cols, top, width/cols, height)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_PRIMARY
        cell.line.fill.background()
        
        text = slide.shapes.add_textbox(left + i*width/cols + 0.2, top + 0.15, width/cols - 0.4, 0.4)
        tf = text.text_frame
        p = tf.paragraphs[0]
        p.text = h
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
    
    # 数据行
    data = [
        ("软件开发", "6大智能体开发+平台建设", "180"),
        ("实施服务", "调研+实施+培训", "80"),
        ("运维保障", "1年运维+持续优化", "40")
    ]
    
    for row_idx, row_data in enumerate(data):
        y = top + (row_idx + 1) * height
        for col_idx, cell_data in enumerate(row_data):
            x = left + col_idx * width / cols
            cell = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, width/cols, height)
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_LIGHT_GRAY
            cell.line.color.rgb = COLOR_GRAY
            
            text = slide.shapes.add_textbox(x + 0.2, y + 0.15, width/cols - 0.4, 0.4)
            tf = text.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = cell_data
            p.font.size = Pt(12)
            p.font.color.rgb = COLOR_PRIMARY if col_idx == 2 else COLOR_GRAY
    
    # 总计
    y = top + 4 * height
    total_cell = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, y, width, height)
    total_cell.fill.solid()
    total_cell.fill.fore_color.rgb = COLOR_ACCENT
    total_cell.line.fill.background()
    
    total_text = slide.shapes.add_textbox(left + 0.2, y + 0.15, width - 0.4, 0.4)
    tf = total_text.text_frame
    p = tf.paragraphs[0]
    p.text = "项目总计：300万元"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE

def add_closing_slide(prs):
    """封底页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 全屏背景
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_PRIMARY
    bg.line.fill.background()
    
    # 中心圆
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(5.4), Inches(1.5), Inches(2.5), Inches(2.5)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = COLOR_ACCENT
    circle.line.fill.background()
    
    # 主标题
    title_box = slide.shapes.add_textbox(Inches(2), Inches(4.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "下一步行动"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # 行动项
    actions = ["① 安排高管对接会议  ② 开展现场需求调研  ③ 确认试点区域"]
    action_box = slide.shapes.add_textbox(Inches(2), Inches(5.3), Inches(9), Inches(1))
    tf = action_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = actions[0]
    p.font.size = Pt(18)
    p.font.color.rgb = COLOR_ACCENT
    p.alignment = PP_ALIGN.CENTER
    
    # 联系方式
    contact_box = slide.shapes.add_textbox(Inches(4), Inches(6.5), Inches(5), Inches(0.6))
    tf = contact_box.text_frame
    p = tf.paragraphs[0]
    p.text = "联系人：九星智囊团 | contact@jiuxing.ai"
    p.font.size = Pt(14)
    p.font.color.rgb = COLOR_LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER


def main():
    prs = create_presentation()
    
    # 1. 封面
    add_title_slide(prs, 
                    "湖北电信AI智能配案",
                    "一键生成个性化营销方案\n赋能一线营收增长")
    
    # 2. 现状 - 营销痛点数据
    chart_data = CategoryChartData()
    chart_data.categories = ['转化率', '客单价', '响应速度', '推荐准确率']
    chart_data.add_series('当前水平', (12, 180, 24, 25))
    chart_data.add_series('行业标杆', (28, 320, 4, 65))
    add_chart_slide(prs, "营销效率显著落后", chart_data, XL_CHART_TYPE.COLUMN_CLUSTERED)
    
    # 3. 三大挑战
    add_three_column_slide(prs, "三大核心挑战", [
        {"title": "数据孤岛", "content": "CRM、BOSS、渠道系统独立\n用户画像不完整"},
        {"title": "能力断层", "content": "优秀经验难复制\n新人上手周期3-6个月"},
        {"title": "效率瓶颈", "content": "方案配置依赖人工\n日均处理量<50单"}
    ])
    
    # 4. 机遇 - AI窗口期
    chart_data = CategoryChartData()
    chart_data.categories = ['2023', '2024', '2025', '2026(E)']
    chart_data.add_series('AI渗透率', (8, 22, 45, 72))
    add_chart_slide(prs, "AI+电信进入黄金窗口期", chart_data, XL_CHART_TYPE.LINE)
    
    # 5. 方案架构
    add_content_slide(prs, "AI智能配案架构",
                      content_left="四层架构\n• 数据层：全域数据汇聚\n• 能力层：6大智能体\n• 应用层：营销全流程\n• 展现层：多端触达",
                      content_right="核心技术\n• 大语言模型\n• 知识图谱\n• 推荐算法\n• 实时计算")
    
    # 6. 6大智能体
    add_six_agents_slide(prs)
    
    # 7. ROI收益
    chart_data = CategoryChartData()
    chart_data.categories = ['转化率↑', '客单价↑', '处理效率↑', '人力成本↓']
    chart_data.add_series('预期提升', (135, 78, 400, -30))
    add_chart_slide(prs, "量化收益ROI", chart_data, XL_CHART_TYPE.BAR_CLUSTERED)
    
    # 8. Demo效果
    add_content_slide(prs, "Demo演示效果",
                      content_left="智能配案只需3步\n\n1. 输入客户基本特征\n2. AI自动分析需求\n3. 一键生成方案\n\n客户经理反馈：\n\"原来需要2小时的配案\n现在3分钟就完成了\"",
                      content_right="")
    
    # 9. 实施路线图
    add_roadmap_slide(prs)
    
    # 10. 团队
    add_team_slide(prs)
    
    # 11. 报价
    add_pricing_slide(prs)
    
    # 12. 封底
    add_closing_slide(prs)
    
    # 保存
    output_path = "/root/.openclaw/workspace/03_输出成果/湖北电信AI智能配案_麦肯锡风格.pptx"
    prs.save(output_path)
    print(f"PPT已保存到: {output_path}")
    print(f"总页数: {len(prs.slides)}")

if __name__ == "__main__":
    main()

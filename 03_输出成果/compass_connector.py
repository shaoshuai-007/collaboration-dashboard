#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
指南针工程对接模块
- 对接采薇（需求文档）
- 对接织锦（思维导图）
- 对接筑台（方案举措）
- 对接呈彩（方案PPT）
- 对接工尺（详细设计）
- 对接玉衡（项目管控）

Author: 南乔
Date: 2026-03-14
"""

from typing import List, Dict, Optional
from dataclasses import dataclass
from datetime import datetime
import json
import os
import subprocess


# ==================== 对接配置 ====================
COMPASS_SKILLS = {
    'caiwei': {
        'name': '采薇',
        'skill': 'compass-needdoc',
        'description': '需求文档生成',
        'input_type': 'discussion_result',
        'output_type': 'srs_document',
        'priority': 1
    },
    'zhijin': {
        'name': '织锦',
        'skill': 'compass-mindmap',
        'description': '思维导图生成',
        'input_type': 'requirement_document',
        'output_type': 'mindmap_html',
        'priority': 2
    },
    'zhutai': {
        'name': '筑台',
        'skill': 'compass-solution',
        'description': '方案举措生成',
        'input_type': 'mindmap',
        'output_type': 'solution_excel',
        'priority': 3
    },
    'chengcai': {
        'name': '呈彩',
        'skill': 'compass-ppt',
        'description': '方案PPT生成',
        'input_type': 'solution_excel',
        'output_type': 'pptx',
        'priority': 4
    },
    'gongchi': {
        'name': '工尺',
        'skill': 'compass-design',
        'description': '详细设计生成',
        'input_type': 'pptx',
        'output_type': 'design_document',
        'priority': 5
    },
    'yuheng': {
        'name': '玉衡',
        'skill': 'compass-project',
        'description': '项目管控生成',
        'input_type': 'design_document',
        'output_type': 'project_excel',
        'priority': 6
    }
}


# ==================== 对接基类 ====================
class CompassSkillConnector:
    """指南针技能对接基类"""
    
    def __init__(self, skill_id: str):
        self.skill_id = skill_id
        self.config = COMPASS_SKILLS.get(skill_id, {})
        self.skill_path = f"/root/.openclaw/skills/{self.config.get('skill', '')}"
    
    def check_skill_available(self) -> bool:
        """检查技能是否可用"""
        return os.path.exists(self.skill_path)
    
    def prepare_input(self, data: Dict) -> str:
        """准备输入数据"""
        raise NotImplementedError
    
    def execute(self, input_data: str) -> Dict:
        """执行技能"""
        raise NotImplementedError
    
    def parse_output(self, output: str) -> Dict:
        """解析输出"""
        raise NotImplementedError


# ==================== 采薇对接（需求文档） ====================
class CaiweiConnector(CompassSkillConnector):
    """采薇对接 - 需求文档生成"""
    
    def __init__(self):
        super().__init__('caiwei')
    
    def prepare_input(self, discussion_data: Dict) -> str:
        """准备需求文档输入"""
        return json.dumps({
            'task': discussion_data.get('task', ''),
            'turns': discussion_data.get('turns', []),
            'key_points': discussion_data.get('key_points', []),
            'risks': discussion_data.get('risks', []),
            'consensus_level': discussion_data.get('consensus_level', 0)
        }, ensure_ascii=False)
    
    def execute(self, input_data: str) -> Dict:
        """
        执行需求文档生成
        
        Returns:
            {
                'success': bool,
                'output_path': str,
                'srs_data': dict
            }
        """
        data = json.loads(input_data)
        
        # 生成SRS数据结构
        srs_data = self._generate_srs_structure(data)
        
        # 生成Word文档
        output_path = self._generate_word_document(srs_data)
        
        return {
            'success': True,
            'output_path': output_path,
            'srs_data': srs_data
        }
    
    def _generate_srs_structure(self, data: Dict) -> Dict:
        """生成SRS数据结构"""
        task = data.get('task', '未命名项目')
        key_points = data.get('key_points', [])
        risks = data.get('risks', [])
        
        # 提取功能需求
        functional_reqs = []
        for i, point in enumerate(key_points[:10]):
            functional_reqs.append({
                'req_id': f'FR-{i+1:03d}',
                'name': point[:50],
                'priority': self._detect_priority(point),
                'user_story': f'作为用户，我希望{point[:30]}',
                'acceptance_criteria': [f'完成{point[:20]}的实现']
            })
        
        # 提取非功能需求
        non_functional_reqs = self._extract_nfr(data.get('turns', []))
        
        return {
            'document_info': {
                'project_name': task.split('需求')[0].strip() if '需求' in task else task,
                'version': 'V1.0',
                'author': '智能体协作平台',
                'date': datetime.now().strftime('%Y-%m-%d')
            },
            'overview': {
                'description': f'本文档基于团队讨论自动生成。任务：{task}',
                'scope': '；'.join(key_points[:5]) if key_points else '详见需求列表'
            },
            'functional_requirements': functional_reqs,
            'non_functional_requirements': non_functional_reqs,
            'constraints': [f'风险约束：{r}' for r in risks]
        }
    
    def _detect_priority(self, text: str) -> str:
        """检测优先级"""
        high_keywords = ['必须', '核心', '关键', '重要', '基础']
        if any(kw in text for kw in high_keywords):
            return '高'
        return '中'
    
    def _extract_nfr(self, turns: List[Dict]) -> List[Dict]:
        """提取非功能需求"""
        nfrs = []
        nfr_keywords = {
            '性能': ['性能', '响应', '并发', '吞吐'],
            '安全': ['安全', '加密', '权限', '认证'],
            '可用性': ['可用', '稳定', '可靠'],
            '可扩展': ['扩展', '升级', '迭代']
        }
        
        for turn in turns:
            content = turn.get('content', '')
            for nfr_type, keywords in nfr_keywords.items():
                if any(kw in content for kw in keywords):
                    nfrs.append({
                        'type': nfr_type,
                        'description': content[:100]
                    })
                    break
        
        return nfrs
    
    def _generate_word_document(self, srs_data: Dict) -> str:
        """生成Word文档"""
        from docx import Document
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        
        doc = Document()
        
        # 标题
        title = doc.add_heading('需求规格说明书（SRS）', 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # 项目信息
        doc.add_heading('一、项目信息', level=1)
        info = srs_data['document_info']
        doc.add_paragraph(f"项目名称：{info['project_name']}")
        doc.add_paragraph(f"版本：{info['version']}")
        doc.add_paragraph(f"编制人：{info['author']}")
        doc.add_paragraph(f"日期：{info['date']}")
        
        # 概述
        doc.add_heading('二、项目概述', level=1)
        doc.add_paragraph(srs_data['overview']['description'])
        doc.add_paragraph(f"范围：{srs_data['overview']['scope']}")
        
        # 功能需求
        doc.add_heading('三、功能需求', level=1)
        for req in srs_data['functional_requirements']:
            doc.add_heading(f"{req['req_id']} {req['name']}", level=2)
            doc.add_paragraph(f"优先级：{req['priority']}")
            doc.add_paragraph(f"用户故事：{req['user_story']}")
        
        # 非功能需求
        if srs_data['non_functional_requirements']:
            doc.add_heading('四、非功能需求', level=1)
            for nfr in srs_data['non_functional_requirements']:
                doc.add_paragraph(f"• {nfr['type']}：{nfr['description']}")
        
        # 约束条件
        if srs_data['constraints']:
            doc.add_heading('五、约束条件', level=1)
            for constraint in srs_data['constraints']:
                doc.add_paragraph(f"• {constraint}")
        
        # 保存
        output_dir = '/root/.openclaw/workspace/03_输出成果/指南针输出'
        os.makedirs(output_dir, exist_ok=True)
        output_path = f"{output_dir}/SRS_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        doc.save(output_path)
        
        return output_path


# ==================== 织锦对接（思维导图） ====================
class ZhijinConnector(CompassSkillConnector):
    """织锦对接 - 思维导图生成"""
    
    def __init__(self):
        super().__init__('zhijin')
    
    def prepare_input(self, requirement_data: Dict) -> str:
        """准备思维导图输入"""
        return json.dumps(requirement_data, ensure_ascii=False)
    
    def execute(self, input_data: str) -> Dict:
        """执行思维导图生成"""
        data = json.loads(input_data)
        
        # 生成思维导图数据结构
        mindmap_data = self._generate_mindmap_structure(data)
        
        # 生成HTML
        output_path = self._generate_html_mindmap(mindmap_data)
        
        return {
            'success': True,
            'output_path': output_path,
            'mindmap_data': mindmap_data
        }
    
    def _generate_mindmap_structure(self, data: Dict) -> Dict:
        """生成思维导图结构"""
        project_name = data.get('document_info', {}).get('project_name', '未命名项目')
        func_reqs = data.get('functional_requirements', [])
        
        # 构建思维导图节点
        children = []
        for req in func_reqs[:8]:  # 最多8个分支
            children.append({
                'name': req['name'][:20],
                'children': [
                    {'name': f"优先级：{req['priority']}"},
                    {'name': req['user_story'][:30]}
                ]
            })
        
        return {
            'name': project_name,
            'children': [
                {'name': '功能需求', 'children': children},
                {'name': '非功能需求', 'children': [
                    {'name': nfr['type']} for nfr in data.get('non_functional_requirements', [])[:4]
                ]},
                {'name': '约束条件', 'children': [
                    {'name': c[:20]} for c in data.get('constraints', [])[:3]
                ]}
            ]
        }
    
    def _generate_html_mindmap(self, mindmap_data: Dict) -> str:
        """生成HTML思维导图"""
        output_dir = '/root/.openclaw/workspace/03_输出成果/指南针输出'
        os.makedirs(output_dir, exist_ok=True)
        output_path = f"{output_dir}/思维导图_{datetime.now().strftime('%Y%m%d_%H%M%S')}.html"
        
        # 生成简单HTML
        html_content = f"""<!DOCTYPE html>
<html>
<head>
    <meta charset="UTF-8">
    <title>思维导图 - {mindmap_data['name']}</title>
    <style>
        body {{ font-family: Arial, sans-serif; margin: 20px; }}
        .mindmap {{ padding: 20px; }}
        .node {{ margin: 10px 0; padding: 10px; background: #f0f0f0; border-radius: 5px; }}
        .level-1 {{ font-size: 24px; font-weight: bold; color: #C93832; }}
        .level-2 {{ font-size: 18px; margin-left: 20px; color: #006EBD; }}
        .level-3 {{ font-size: 14px; margin-left: 40px; color: #595959; }}
    </style>
</head>
<body>
    <h1>思维导图 - {mindmap_data['name']}</h1>
    <div class="mindmap">
        {self._render_nodes(mindmap_data, 1)}
    </div>
</body>
</html>"""
        
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(html_content)
        
        return output_path
    
    def _render_nodes(self, node: Dict, level: int) -> str:
        """递归渲染节点"""
        html = f'<div class="node level-{level}">{node["name"]}</div>\n'
        for child in node.get('children', []):
            html += self._render_nodes(child, level + 1)
        return html


# ==================== 筑台对接（方案举措） ====================
class ZhutaiConnector(CompassSkillConnector):
    """筑台对接 - 方案举措生成"""
    
    def __init__(self):
        super().__init__('zhutai')
    
    def prepare_input(self, mindmap_data: Dict) -> str:
        return json.dumps(mindmap_data, ensure_ascii=False)
    
    def execute(self, input_data: str) -> Dict:
        data = json.loads(input_data)
        output_path = self._generate_solution_excel(data)
        
        return {
            'success': True,
            'output_path': output_path
        }
    
    def _generate_solution_excel(self, data: Dict) -> str:
        """生成方案举措Excel"""
        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill
        
        wb = Workbook()
        ws = wb.active
        ws.title = '方案举措'
        
        # 表头
        headers = ['序号', '举措名称', '目标', '措施', '责任人', '完成时间', '状态']
        for col, header in enumerate(headers, 1):
            cell = ws.cell(row=1, column=col, value=header)
            cell.font = Font(bold=True)
            cell.fill = PatternFill(start_color='C93832', end_color='C93832', fill_type='solid')
            cell.font = Font(bold=True, color='FFFFFF')
        
        # 从思维导图提取举措
        row = 2
        for i, child in enumerate(data.get('children', [])[:8], 1):
            ws.cell(row=row, column=1, value=i)
            ws.cell(row=row, column=2, value=child.get('name', '')[:30])
            ws.cell(row=row, column=3, value='待定义')
            ws.cell(row=row, column=4, value='待细化')
            ws.cell(row=row, column=5, value='待分配')
            ws.cell(row=row, column=6, value='')
            ws.cell(row=row, column=7, value='待开始')
            row += 1
        
        output_dir = '/root/.openclaw/workspace/03_输出成果/指南针输出'
        os.makedirs(output_dir, exist_ok=True)
        output_path = f"{output_dir}/方案举措_{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx"
        wb.save(output_path)
        
        return output_path


# ==================== 呈彩对接（方案PPT） ====================
class ChengcaiConnector(CompassSkillConnector):
    """呈彩对接 - 方案PPT生成"""
    
    def __init__(self):
        super().__init__('chengcai')
    
    def prepare_input(self, solution_data: Dict) -> str:
        return json.dumps(solution_data, ensure_ascii=False)
    
    def execute(self, input_data: str) -> Dict:
        data = json.loads(input_data)
        output_path = self._generate_ppt(data)
        
        return {
            'success': True,
            'output_path': output_path
        }
    
    def _generate_ppt(self, data: Dict) -> str:
        """生成PPT"""
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        
        prs = Presentation()
        
        # 封面
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide.shapes.title
        title.text = data.get('name', '方案汇报')
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(201, 56, 50)  # 电信红
        
        # 目录页
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = "目录"
        
        # 内容页（从children生成）
        for child in data.get('children', [])[:5]:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = slide.shapes.title
            title.text = child.get('name', '')
            title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 110, 189)  # 深蓝
        
        output_dir = '/root/.openclaw/workspace/03_输出成果/指南针输出'
        os.makedirs(output_dir, exist_ok=True)
        output_path = f"{output_dir}/方案PPT_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        prs.save(output_path)
        
        return output_path


# ==================== 工尺对接（详细设计） ====================
class GongchiConnector(CompassSkillConnector):
    """工尺对接 - 详细设计生成"""
    
    def __init__(self):
        super().__init__('gongchi')
    
    def prepare_input(self, ppt_data: Dict) -> str:
        return json.dumps(ppt_data, ensure_ascii=False)
    
    def execute(self, input_data: str) -> Dict:
        data = json.loads(input_data)
        output_path = self._generate_design_doc(data)
        
        return {
            'success': True,
            'output_path': output_path
        }
    
    def _generate_design_doc(self, data: Dict) -> str:
        """生成详细设计文档"""
        from docx import Document
        
        doc = Document()
        doc.add_heading('详细设计文档', 0)
        
        doc.add_heading('一、系统架构设计', level=1)
        doc.add_paragraph('（待补充）')
        
        doc.add_heading('二、接口设计', level=1)
        doc.add_paragraph('（待补充）')
        
        doc.add_heading('三、数据库设计', level=1)
        doc.add_paragraph('（待补充）')
        
        doc.add_heading('四、安全设计', level=1)
        doc.add_paragraph('（待补充）')
        
        output_dir = '/root/.openclaw/workspace/03_输出成果/指南针输出'
        os.makedirs(output_dir, exist_ok=True)
        output_path = f"{output_dir}/详细设计_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        doc.save(output_path)
        
        return output_path


# ==================== 玉衡对接（项目管控） ====================
class YuhengConnector(CompassSkillConnector):
    """玉衡对接 - 项目管控生成"""
    
    def __init__(self):
        super().__init__('yuheng')
    
    def prepare_input(self, design_data: Dict) -> str:
        return json.dumps(design_data, ensure_ascii=False)
    
    def execute(self, input_data: str) -> Dict:
        data = json.loads(input_data)
        output_path = self._generate_project_excel(data)
        
        return {
            'success': True,
            'output_path': output_path
        }
    
    def _generate_project_excel(self, data: Dict) -> str:
        """生成项目管控Excel"""
        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill
        
        wb = Workbook()
        
        # 任务清单
        ws = wb.active
        ws.title = '任务清单'
        headers = ['任务ID', '任务名称', '负责人', '开始时间', '结束时间', '状态']
        for col, header in enumerate(headers, 1):
            cell = ws.cell(row=1, column=col, value=header)
            cell.font = Font(bold=True)
        
        # 风险清单
        ws2 = wb.create_sheet('风险清单')
        headers = ['风险ID', '风险描述', '影响程度', '应对措施', '责任人']
        for col, header in enumerate(headers, 1):
            cell = ws2.cell(row=1, column=col, value=header)
            cell.font = Font(bold=True)
        
        output_dir = '/root/.openclaw/workspace/03_输出成果/指南针输出'
        os.makedirs(output_dir, exist_ok=True)
        output_path = f"{output_dir}/项目管控_{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx"
        wb.save(output_path)
        
        return output_path


# ==================== 统一对接API ====================
class CompassConnector:
    """指南针工程统一对接入口"""
    
    def __init__(self):
        self.connectors = {
            'caiwei': CaiweiConnector(),
            'zhijin': ZhijinConnector(),
            'zhutai': ZhutaiConnector(),
            'chengcai': ChengcaiConnector(),
            'gongchi': GongchiConnector(),
            'yuheng': YuhengConnector()
        }
    
    def execute_workflow(self, discussion_data: Dict, 
                        skills: List[str] = None) -> Dict:
        """
        执行工作流
        
        Args:
            discussion_data: 讨论数据
            skills: 要执行的技能列表，默认全部
        
        Returns:
            执行结果
        """
        if skills is None:
            skills = ['caiwei', 'zhijin', 'zhutai', 'chengcai', 'gongchi', 'yuheng']
        
        results = {}
        current_data = discussion_data
        
        for skill_id in skills:
            if skill_id not in self.connectors:
                continue
            
            connector = self.connectors[skill_id]
            
            # 准备输入
            input_data = connector.prepare_input(current_data)
            
            # 执行
            result = connector.execute(input_data)
            results[skill_id] = result
            
            # 传递给下一阶段
            if result['success']:
                current_data = result.get('srs_data') or result.get('mindmap_data') or result
        
        return results
    
    def execute_single(self, skill_id: str, data: Dict) -> Dict:
        """执行单个技能"""
        if skill_id not in self.connectors:
            return {'success': False, 'error': f'未知技能: {skill_id}'}
        
        connector = self.connectors[skill_id]
        input_data = connector.prepare_input(data)
        return connector.execute(input_data)


# ==================== 测试 ====================
if __name__ == '__main__':
    print("=" * 60)
    print("指南针工程对接模块测试")
    print("=" * 60)
    
    # 测试数据
    test_discussion = {
        'task': '智能客服系统需求分析',
        'turns': [
            {'speaker_name': '采薇', 'content': '【需求分析】核心功能：智能问答、多轮对话、知识库管理。'},
            {'speaker_name': '织锦', 'content': '【架构设计】推荐微服务架构，需考虑高并发性能。'},
            {'speaker_name': '筑台', 'content': '【成本质疑】微服务架构成本较高，需评估必要性。'},
        ],
        'key_points': ['智能问答功能', '多轮对话能力', '知识库管理系统', '高并发性能支持'],
        'risks': ['微服务架构成本较高'],
        'consensus_level': 75
    }
    
    connector = CompassConnector()
    
    # 测试工作流
    print("\n1. 测试完整工作流...")
    results = connector.execute_workflow(test_discussion)
    
    for skill_id, result in results.items():
        skill_name = COMPASS_SKILLS[skill_id]['name']
        if result['success']:
            print(f"   ✅ {skill_name}: {result['output_path']}")
        else:
            print(f"   ❌ {skill_name}: {result.get('error', '执行失败')}")
    
    print("\n✅ 指南针工程对接测试完成")

#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""麦肯锡风格PPT生成器 - 简化版"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# 麦肯锡配色
MCKINSEY = {
    'primary': RGBColor(26, 54, 93),      # 深蓝 #1A365D
    'secondary': RGBColor(74, 85, 104),   # 灰色 #4A5568
    'accent': RGBColor(49, 130, 206),     # 强调蓝 #3182CE
    'text': RGBColor(26, 32, 44),         # 近黑 #1A202C
    'white': RGBColor(255, 255, 255),
    'lightgray': RGBColor(240, 240, 240),
}

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_slide_title(slide, text, fontsize=36):
    """添加页面标题"""
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(fontsize)
    p.font.color.rgb = MCKINSEY['primary']
    p.font.bold = True

def add_bullet(slide, text, top=2):
    """添加要点"""
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(12), Inches(0.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "• " + text
    p.font.size = Pt(20)
    p.font.color.rgb = MCKINSEY['text']
    p.space_after = Pt(10)
    return top + 0.5

def add_box(slide, title, top=2, height=4):
    """添加内容框"""
    # 标题
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(12), Inches(0.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.color.rgb = MCKINSEY['secondary']
    p.font.bold = True
    
    # 占位框
    shape = slide.shapes.add_shape(1, Inches(0.5), Inches(top+0.6), Inches(12.333), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = MCKINSEY['lightgray']
    shape.line.color.rgb = MCKINSEY['secondary']
    
    # 标签
    label = slide.shapes.add_textbox(Inches(5), Inches(top+height/2), Inches(3), Inches(0.5))
    tf = label.text_frame
    p = tf.paragraphs[0]
    p.text = "[图表占位]"
    p.font.size = Pt(16)
    p.font.color.rgb = MCKINSEY['secondary']
    p.alignment = PP_ALIGN.CENTER

# === 1. 封面 ===
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
title.text = "湖北电信AI智能配案系统"
title.text_frame.paragraphs[0].font.size = Pt(44)
title.text_frame.paragraphs[0].font.color.rgb = MCKINSEY['primary']
title.text_frame.paragraphs[0].font.bold = True
title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

subtitle = slide.placeholders[1]
subtitle.text = "AI赋能电信营销，实现智能推荐与精准配案"
subtitle.text_frame.paragraphs[0].font.size = Pt(24)
subtitle.text_frame.paragraphs[0].font.color.rgb = MCKINSEY['secondary']
subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# === 2. 现状 ===
slide = prs.slides.add_slide(prs.slide_layouts[5])
add_slide_title(slide, "现状：营销效率瓶颈")
add_box(slide, "关键数据洞察", top=1.5, height=4.5)
add_bullet(slide, "传统配案依赖人工，效率低", top=2)
add_bullet(slide, "客户经理需花费大量时间筛选产品", top=2.5)
add_bullet(slide, "配案成功率仅35%，客户满意度待提升", top=3)

# === 3. 挑战 ===
slide = prs.slides.add_slide(prs.slide_layouts[5])
add_slide_title(slide, "挑战：三大核心痛点")
challenges = [
    "产品组合复杂：200+产品，组合规则难掌握",
    "客户需求多样：需快速匹配个性化方案",
    "依赖经验积累：新人上手周期长达3个月",
]
top = 1.8
for c in challenges:
    top = add_bullet(slide, c, top)

# === 4. 机遇 ===
slide = prs.slides.add_slide(prs.slide_layouts[5])
add_slide_title(slide, "机遇：AI+电信窗口期")
opportunities = [
    "政策支持：AI项目无上限，电信行业优先",
    "技术成熟：大模型能力突破，业务落地可期",
    "竞争态势：友商尚未大规模布局，存在先机",
]
top = 1.8
for o in opportunities:
    top = add_bullet(slide, o, top)

# === 5. 方案 ===
slide = prs.slides.add_slide(prs.slide_layouts[5])
add_slide_title(slide, "方案：AI智能配案架构")
add_box(slide, "五层系统架构", top=1.5, height=4.5)

# === 6. 能力 ===
slide = prs.slides.add_slide(prs.slide_layouts[5])
add_slide_title(slide, "能力：6大智能体")
abilities = [
    "需求分析智能体：精准识别客户需求",
    "产品推荐智能体：智能匹配最优产品",
    "方案生成智能体：自动生成配案方案",
    "价值测算智能体：实时计算ROI",
    "风险预警智能体：智能识别潜在风险",
    "知识问答智能体：实时解答业务问题",
]
top = 1.8
for a in abilities:
    top = add_bullet(slide, a, top)

# === 7. 价值 ===
slide = prs.slides.add_slide(prs.slide_layouts[5])
add_slide_title(slide, "价值：量化收益")
add_box(slide, "ROI测算", top=1.5, height=2.5)
values = [
    "配案效率提升80%：从30分钟→6分钟",
    "转化率提升57%：35%→55%",
    "人效提升：客户经理服务能力翻倍",
    "年度收益：预计节省人力成本500万+",
]
top = 4.5
for v in values:
    top = add_bullet(slide, v, top)

# === 8. 计划 ===
slide = prs.slides.add_slide(prs.slide_layouts[5])
add_slide_title(slide, "计划：实施路线图")
add_box(slide, "12周甘特图", top=1.5, height=4.5)

# === 9. 团队 ===
slide = prs.slides.add_slide(prs.slide_layouts[5])
add_slide_title(slide, "团队：九星智囊团")
team = [
    "🌸 采薇  需求分析  |  🧵 织锦  架构设计  |  🎨 呈彩  PPT制作",
    "📐 工尺  详细设计  |  ⚖️ 玉衡  项目管控  |  🏗️ 筑台  售前支持",
    "📚 折桂  知识管理  |  💻 天工  开发实现  |  📊 知微  数据分析",
]
top = 1.8
for t in team:
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(12), Inches(0.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = t
    p.font.size = Pt(18)
    p.font.color.rgb = MCKINSEY['text']
    top += 0.7

# === 10. 封底 ===
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
title.text = "下一步行动"
title.text_frame.paragraphs[0].font.size = Pt(40)
title.text_frame.paragraphs[0].font.color.rgb = MCKINSEY['primary']
title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

subtitle = slide.placeholders[1]
subtitle.text = "1. 立项审批  →  2. 需求调研  →  3. 原型开发\n\n联系人：少帅 | 湖北电信电渠中心"
subtitle.text_frame.paragraphs[0].font.size = Pt(20)
subtitle.text_frame.paragraphs[0].font.color.rgb = MCKINSEY['secondary']
subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# 保存
output_path = "/root/.openclaw/workspace/03_输出成果/湖北电信AI智能配案_麦肯锡风格.pptx"
prs.save(output_path)
print(f"✅ 已生成: {output_path}")
print(f"📊 文件大小: {os.path.getsize(output_path)} bytes")
print(f"📑 共 {len(prs.slides)} 页")

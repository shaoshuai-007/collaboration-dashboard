#!/usr/bin/env python3
"""
NVIDIA GTC 2026 黄仁勋主题演讲总结 PPT生成器 V3
内容丰富、详细、专业的版本
"""

import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor as RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
import os

# 配色方案
COLORS = {
    'primary': RgbColor(0xC9, 0x38, 0x32),      # 电信红
    'secondary': RgbColor(0x00, 0x6E, 0xBD),    # 深蓝
    'neutral': RgbColor(0x59, 0x59, 0x59),      # 灰色
    'white': RgbColor(0xFF, 0xFF, 0xFF),
    'light_gray': RgbColor(0xF5, 0xF5, 0xF5),
    'dark_gray': RgbColor(0x33, 0x33, 0x33),
    'accent': RgbColor(0x00, 0x7A, 0xCC),       # 亮蓝
}

def create_presentation():
    """创建PPT演示文稿"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs

def add_cover_slide(prs, title, subtitle):
    """添加封面页"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景 - 渐变效果（使用深蓝色背景）
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5)
    )
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['secondary']
    background.line.fill.background()
    
    # 装饰条 - 电信红
    accent_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(2.5), Inches(13.333), Inches(0.1)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = COLORS['primary']
    accent_bar.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12.333), Inches(1))
    tf = subtitle_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    
    # 底部信息
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.5))
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.text = "九星智囊团 · 2026年3月"
    p.font.size = Pt(14)
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_toc_slide(prs, sections):
    """添加目录页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5)
    )
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['white']
    background.line.fill.background()
    
    # 标题区域 - 电信红背景
    title_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.2)
    )
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = COLORS['primary']
    title_bg.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "目录"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    
    # 目录内容 - 分栏显示
    col_width = Inches(6)
    start_y = Inches(1.6)
    item_height = Inches(0.4)
    
    for i, section in enumerate(sections):
        col = 0 if i < 3 else 1
        row = i % 3
        x = Inches(0.5) + col * Inches(6.5)
        y = start_y + row * Inches(1.8)
        
        # 部分标题
        section_box = slide.shapes.add_textbox(x, y, col_width, Inches(0.4))
        tf = section_box.text_frame
        p = tf.paragraphs[0]
        p.text = section['name']
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLORS['primary']
        
        # 子项目
        for j, item in enumerate(section['items']):
            item_box = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(0.5) + j * Inches(0.35), col_width - Inches(0.3), Inches(0.35))
            tf = item_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"• {item}"
            p.font.size = Pt(14)
            p.font.color.rgb = COLORS['neutral']
    
    return slide

def add_content_slide_grid(prs, title, sections):
    """添加内容页 - 网格布局"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5)
    )
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['white']
    background.line.fill.background()
    
    # 标题区域
    title_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.2)
    )
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = COLORS['primary']
    title_bg.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    
    # 内容区域 - 根据section数量决定布局
    num_sections = len(sections)
    if num_sections <= 2:
        # 1行2列
        cols, rows = 2, 1
    elif num_sections <= 4:
        # 2行2列
        cols, rows = 2, 2
    else:
        # 2行3列
        cols, rows = 3, 2
    
    card_width = (Inches(12.333) - Inches(0.5) * (cols - 1)) / cols
    card_height = (Inches(5.8) - Inches(0.3) * (rows - 1)) / rows
    
    for i, section in enumerate(sections):
        col = i % cols
        row = i // cols
        x = Inches(0.5) + col * (card_width + Inches(0.5))
        y = Inches(1.5) + row * (card_height + Inches(0.3))
        
        # 卡片背景
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_width, card_height
        )
        card.fill.solid()
        card.fill.fore_color.rgb = COLORS['light_gray']
        card.line.color.rgb = COLORS['secondary']
        card.line.width = Pt(1)
        
        # 卡片标题
        card_title = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.15), card_width - Inches(0.4), Inches(0.4))
        tf = card_title.text_frame
        p = tf.paragraphs[0]
        p.text = section.get('title', '')
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLORS['primary']
        
        # 卡片内容
        content = section.get('content', '')
        if content:
            content_box = slide.shapes.add_textbox(
                x + Inches(0.2), y + Inches(0.55), 
                card_width - Inches(0.4), card_height - Inches(0.7)
            )
            tf = content_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = content
            p.font.size = Pt(11)
            p.font.color.rgb = COLORS['dark_gray']
            p.line_spacing = 1.2
        
        # 指标数据
        metrics = section.get('metrics', [])
        if metrics:
            metrics_y = y + Inches(0.55)
            for metric in metrics:
                metric_box = slide.shapes.add_textbox(
                    x + Inches(0.2), metrics_y,
                    card_width - Inches(0.4), Inches(0.35)
                )
                tf = metric_box.text_frame
                p = tf.paragraphs[0]
                # 处理不同格式的metrics
                if isinstance(metric, dict):
                    p.text = f"• {metric.get('label', '')}: {metric.get('value', '')}"
                else:
                    p.text = f"• {metric}"
                p.font.size = Pt(11)
                p.font.color.rgb = COLORS['secondary']
                p.font.bold = True
                metrics_y += Inches(0.35)
    
    return slide

def add_content_slide_left_right(prs, title, left_sections, right_sections):
    """添加内容页 - 左右布局"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5)
    )
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['white']
    background.line.fill.background()
    
    # 标题区域
    title_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.2)
    )
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = COLORS['primary']
    title_bg.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    
    # 左侧内容区域
    left_y = Inches(1.5)
    for section in left_sections:
        # 小标题
        section_title = slide.shapes.add_textbox(Inches(0.5), left_y, Inches(6), Inches(0.35))
        tf = section_title.text_frame
        p = tf.paragraphs[0]
        p.text = section.get('title', '')
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLORS['secondary']
        
        left_y += Inches(0.4)
        
        # 内容
        content = section.get('content', '')
        if content:
            content_box = slide.shapes.add_textbox(Inches(0.5), left_y, Inches(6), Inches(0.8))
            tf = content_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = content
            p.font.size = Pt(12)
            p.font.color.rgb = COLORS['dark_gray']
            left_y += Inches(0.9)
        
        # 要点列表
        points = section.get('points', [])
        for point in points:
            point_box = slide.shapes.add_textbox(Inches(0.7), left_y, Inches(5.8), Inches(0.3))
            tf = point_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"• {point}"
            p.font.size = Pt(12)
            p.font.color.rgb = COLORS['dark_gray']
            left_y += Inches(0.35)
        
        left_y += Inches(0.2)
    
    # 右侧数据卡片区域
    right_start_x = Inches(7)
    card_height = (Inches(5.8) - Inches(0.2) * (len(right_sections) - 1)) / len(right_sections)
    
    for i, item in enumerate(right_sections):
        y = Inches(1.5) + i * (card_height + Inches(0.2))
        
        # 卡片背景
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, right_start_x, y, Inches(5.8), card_height
        )
        card.fill.solid()
        card.fill.fore_color.rgb = COLORS['light_gray']
        card.line.color.rgb = COLORS['secondary']
        
        # 卡片名称
        name_box = slide.shapes.add_textbox(right_start_x + Inches(0.2), y + Inches(0.1), Inches(5.4), Inches(0.35))
        tf = name_box.text_frame
        p = tf.paragraphs[0]
        p.text = item.get('name', '')
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS['neutral']
        
        # 数值
        value_box = slide.shapes.add_textbox(right_start_x + Inches(0.2), y + Inches(0.4), Inches(5.4), Inches(0.5))
        tf = value_box.text_frame
        p = tf.paragraphs[0]
        p.text = str(item.get('value', ''))
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = COLORS['primary']
        
        # 描述
        desc = item.get('desc', '') or item.get('trend', '')
        if desc:
            desc_box = slide.shapes.add_textbox(right_start_x + Inches(0.2), y + Inches(0.85), Inches(5.4), Inches(0.3))
            tf = desc_box.text_frame
            p = tf.paragraphs[0]
            p.text = desc
            p.font.size = Pt(11)
            p.font.color.rgb = COLORS['neutral']
    
    return slide

def add_content_slide_top_bottom(prs, title, top_sections, bottom_content):
    """添加内容页 - 上下布局"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5)
    )
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['white']
    background.line.fill.background()
    
    # 标题区域
    title_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.2)
    )
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = COLORS['primary']
    title_bg.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    
    # 上半部分内容区域
    num_cols = len(top_sections)
    col_width = (Inches(12.333) - Inches(0.5) * (num_cols - 1)) / num_cols
    
    for i, section in enumerate(top_sections):
        x = Inches(0.5) + i * (col_width + Inches(0.5))
        y = Inches(1.5)
        
        # 卡片背景
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, col_width, Inches(2.8)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = COLORS['light_gray']
        card.line.color.rgb = COLORS['secondary']
        
        # 小标题
        section_title = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.1), col_width - Inches(0.3), Inches(0.35))
        tf = section_title.text_frame
        p = tf.paragraphs[0]
        p.text = section.get('title', '')
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLORS['primary']
        
        # 内容
        content = section.get('content', '')
        if content:
            content_box = slide.shapes.add_textbox(
                x + Inches(0.15), y + Inches(0.45),
                col_width - Inches(0.3), Inches(2.2)
            )
            tf = content_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = content
            p.font.size = Pt(11)
            p.font.color.rgb = COLORS['dark_gray']
            p.line_spacing = 1.15
    
    # 下半部分 - 表格或要点
    if bottom_content:
        bottom_y = Inches(4.5)
        
        # 底部标题
        bottom_title = slide.shapes.add_textbox(Inches(0.5), bottom_y, Inches(12.333), Inches(0.4))
        tf = bottom_title.text_frame
        p = tf.paragraphs[0]
        p.text = bottom_content.get('title', '')
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLORS['secondary']
        
        # 表格数据
        rows_data = bottom_content.get('rows', [])
        if rows_data:
            table = slide.shapes.add_table(
                len(rows_data) + 1, len(rows_data[0]) if rows_data else 3,
                Inches(0.5), bottom_y + Inches(0.5),
                Inches(12.333), Inches(0.35) * (len(rows_data) + 1)
            ).table
            
            # 设置列宽
            col_widths = [Inches(3), Inches(4), Inches(5.333)]
            for i, width in enumerate(col_widths):
                if i < len(table.columns):
                    table.columns[i].width = width
            
            # 表头
            headers = bottom_content.get('columns', ['指标', '数值', '说明'])
            for i, header in enumerate(headers):
                cell = table.cell(0, i)
                cell.text = header
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS['secondary']
                p = cell.text_frame.paragraphs[0]
                p.font.size = Pt(12)
                p.font.bold = True
                p.font.color.rgb = COLORS['white']
                p.alignment = PP_ALIGN.CENTER
            
            # 数据行
            for row_idx, row_data in enumerate(rows_data):
                for col_idx, cell_text in enumerate(row_data):
                    cell = table.cell(row_idx + 1, col_idx)
                    cell.text = str(cell_text)
                    p = cell.text_frame.paragraphs[0]
                    p.font.size = Pt(11)
                    p.font.color.rgb = COLORS['dark_gray']
                    if col_idx == 0:
                        p.alignment = PP_ALIGN.LEFT
                    else:
                        p.alignment = PP_ALIGN.CENTER
        
        # 要点列表
        points = bottom_content.get('points', [])
        if points:
            points_y = bottom_y + Inches(0.5)
            for point in points:
                point_box = slide.shapes.add_textbox(Inches(0.7), points_y, Inches(12), Inches(0.35))
                tf = point_box.text_frame
                p = tf.paragraphs[0]
                p.text = f"• {point}"
                p.font.size = Pt(12)
                p.font.color.rgb = COLORS['dark_gray']
                points_y += Inches(0.35)
    
    return slide

def add_metrics_slide(prs, title, metrics):
    """添加指标对比页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5)
    )
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['white']
    background.line.fill.background()
    
    # 标题区域
    title_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.2)
    )
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = COLORS['primary']
    title_bg.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    
    # 指标卡片
    num_metrics = len(metrics)
    card_width = (Inches(12.333) - Inches(0.3) * (num_metrics - 1)) / min(num_metrics, 5)
    
    for i, metric in enumerate(metrics):
        col = i % 5
        row = i // 5
        x = Inches(0.5) + col * (card_width + Inches(0.3))
        y = Inches(1.6) + row * Inches(2.9)
        
        # 卡片背景
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_width, Inches(2.7)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = COLORS['light_gray']
        card.line.color.rgb = COLORS['primary']
        card.line.width = Pt(2)
        
        # 指标名称
        name_box = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.1), card_width - Inches(0.3), Inches(0.4))
        tf = name_box.text_frame
        p = tf.paragraphs[0]
        p.text = metric.get('name', '')
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLORS['primary']
        p.alignment = PP_ALIGN.CENTER
        
        # 当前值
        current_box = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.5), card_width - Inches(0.3), Inches(0.35))
        tf = current_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"当前: {metric.get('current', '')}"
        p.font.size = Pt(11)
        p.font.color.rgb = COLORS['neutral']
        p.alignment = PP_ALIGN.CENTER
        
        # 目标值
        target_box = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.85), card_width - Inches(0.3), Inches(0.35))
        tf = target_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"目标: {metric.get('target', '')}"
        p.font.size = Pt(11)
        p.font.color.rgb = COLORS['secondary']
        p.alignment = PP_ALIGN.CENTER
        
        # 提升幅度
        improve_box = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(1.3), card_width - Inches(0.3), Inches(0.5))
        tf = improve_box.text_frame
        p = tf.paragraphs[0]
        p.text = metric.get('improve', '')
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLORS['primary']
        p.alignment = PP_ALIGN.CENTER
        
        # 箭头指示
        arrow_box = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(1.9), card_width - Inches(0.3), Inches(0.5))
        tf = arrow_box.text_frame
        p = tf.paragraphs[0]
        p.text = "↑"
        p.font.size = Pt(24)
        p.font.color.rgb = COLORS['primary']
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_summary_slide(prs, title, key_points):
    """添加总结页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5)
    )
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['secondary']
    background.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    
    # 核心观点
    y = Inches(1.5)
    for i, point in enumerate(key_points):
        # 观点标题
        point_title = slide.shapes.add_textbox(Inches(0.5), y, Inches(12.333), Inches(0.45))
        tf = point_title.text_frame
        p = tf.paragraphs[0]
        p.text = f"{i+1}. {point.get('point', '')}"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        
        # 观点详情
        detail = point.get('detail', '')
        if detail:
            detail_box = slide.shapes.add_textbox(Inches(0.8), y + Inches(0.45), Inches(11.5), Inches(0.5))
            tf = detail_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = detail
            p.font.size = Pt(14)
            p.font.color.rgb = RgbColor(0xE0, 0xE0, 0xE0)
        
        y += Inches(1.1)
    
    return slide

def add_final_slide(prs, title, subtitle, links, footer):
    """添加结束页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5)
    )
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['secondary']
    background.line.fill.background()
    
    # 装饰条
    accent_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(2.5), Inches(13.333), Inches(0.1)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = COLORS['primary']
    accent_bar.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.333), Inches(0.6))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(20)
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    
    # 链接
    y = Inches(5)
    for link in links:
        link_box = slide.shapes.add_textbox(Inches(0.5), y, Inches(12.333), Inches(0.35))
        tf = link_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{link['name']}: {link['url']}"
        p.font.size = Pt(12)
        p.font.color.rgb = RgbColor(0xE0, 0xE0, 0xE0)
        p.alignment = PP_ALIGN.CENTER
        y += Inches(0.35)
    
    # 底部信息
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.5))
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.text = footer
    p.font.size = Pt(14)
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def generate_ppt(outline_path, output_path):
    """生成PPT"""
    # 读取大纲
    with open(outline_path, 'r', encoding='utf-8') as f:
        outline = json.load(f)
    
    # 创建演示文稿
    prs = create_presentation()
    
    # 遍历页面
    for page in outline['pages']:
        page_type = page.get('type', 'content')
        
        if page_type == 'cover':
            # 检查是否有links（结束页）
            if 'links' in page:
                add_final_slide(
                    prs,
                    page.get('title', ''),
                    page.get('subtitle', ''),
                    page.get('links', []),
                    page.get('footer', '')
                )
            else:
                add_cover_slide(
                    prs,
                    page.get('title', ''),
                    page.get('subtitle', '')
                )
        
        elif page_type == 'toc':
            add_toc_slide(prs, page.get('sections', []))
        
        elif page_type == 'content':
            layout = page.get('layout', 'grid')
            
            if layout == 'grid':
                add_content_slide_grid(
                    prs,
                    page.get('title', ''),
                    page.get('sections', [])
                )
            
            elif layout == 'left-text-right-image':
                add_content_slide_left_right(
                    prs,
                    page.get('title', ''),
                    page.get('left_sections', []),
                    page.get('right_sections', [])
                )
            
            elif layout == 'top-text-bottom-image':
                add_content_slide_top_bottom(
                    prs,
                    page.get('title', ''),
                    page.get('top_sections', []),
                    page.get('bottom_content', {})
                )
        
        elif page_type == 'metrics':
            add_metrics_slide(
                prs,
                page.get('title', ''),
                page.get('metrics', [])
            )
        
        elif page_type == 'summary':
            add_summary_slide(
                prs,
                page.get('title', ''),
                page.get('key_points', [])
            )
    
    # 保存
    prs.save(output_path)
    print(f"PPT已生成: {output_path}")
    return output_path

if __name__ == '__main__':
    outline_path = '/root/.openclaw/workspace/03_输出成果/nvidia_gtc_v3_outline.json'
    output_path = '/root/.openclaw/workspace/03_输出成果/NVIDIA_GTC_2026_视频总结_V3.pptx'
    generate_ppt(outline_path, output_path)

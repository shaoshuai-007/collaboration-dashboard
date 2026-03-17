#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPT生成脚本 - 基于美学工坊模板规范
项目：湖北电信AI智能配案系统
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor as RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# 配色方案 - 中国电信
COLORS = {
    'primary': RgbColor(0xC9, 0x38, 0x32),      # 电信红
    'secondary': RgbColor(0xAC, 0x00, 0x00),    # 深红
    'accent': RgbColor(0xED, 0x7D, 0x31),       # 橙色
    'blue': RgbColor(0x3B, 0x82, 0xF6),         # 蓝色
    'dark_blue': RgbColor(0x1E, 0x40, 0xAF),    # 深蓝
    'green': RgbColor(0x10, 0xB9, 0x81),        # 绿色
    'orange': RgbColor(0xF5, 0x9E, 0x0B),       # 橙黄
    'purple': RgbColor(0x8B, 0x5C, 0xF6),       # 紫色
    'text_dark': RgbColor(0x37, 0x41, 0x51),    # 深灰
    'text_gray': RgbColor(0x6B, 0x72, 0x80),    # 灰色
    'bg_light': RgbColor(0xF5, 0xF7, 0xFA),     # 浅灰背景
    'white': RgbColor(0xFF, 0xFF, 0xFF),        # 白色
}

def create_presentation():
    """创建演示文稿"""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    return prs

def add_slide(prs, layout_index=6):
    """添加幻灯片"""
    slide_layout = prs.slide_layouts[layout_index]
    return prs.slides.add_slide(slide_layout)

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=None, align=PP_ALIGN.LEFT):
    """添加文本框"""
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.alignment = align
    if color:
        p.font.color.rgb = color
    return shape

def add_rectangle(slide, left, top, width, height, fill_color=None):
    """添加矩形"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_rounded_rectangle(slide, left, top, width, height, fill_color=None):
    """添加圆角矩形"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

# ==================== 第1页：封面 ====================
def create_cover_slide(prs):
    """封面页"""
    slide = add_slide(prs)
    
    # 背景
    add_rectangle(slide, 0, 0, 13.33, 7.5, COLORS['primary'])
    
    # 主标题
    add_text_box(slide, 1, 2.5, 11.33, 1.2,
                 "湖北电信AI智能配案系统",
                 font_size=48, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)
    
    # 副标题
    add_text_box(slide, 1, 3.8, 11.33, 0.6,
                 "建设方案",
                 font_size=28, color=COLORS['white'], align=PP_ALIGN.CENTER)
    
    # 价值标签
    add_text_box(slide, 1, 4.8, 11.33, 0.5,
                 "集团考核达标  |  一线效率提升  |  客户体验改善",
                 font_size=20, color=COLORS['white'], align=PP_ALIGN.CENTER)
    
    # 单位
    add_text_box(slide, 1, 6.5, 11.33, 0.4,
                 "中国电信湖北分公司",
                 font_size=16, color=COLORS['white'], align=PP_ALIGN.CENTER)
    
    # 日期
    add_text_box(slide, 1, 6.9, 11.33, 0.3,
                 "2026年3月",
                 font_size=14, color=COLORS['white'], align=PP_ALIGN.CENTER)
    
    return slide

# ==================== 第2页：目录 ====================
def create_toc_slide(prs):
    """目录页"""
    slide = add_slide(prs)
    
    # 标题
    add_text_box(slide, 0.5, 0.3, 12.33, 0.6,
                 "目录",
                 font_size=32, bold=True, color=COLORS['dark_blue'], align=PP_ALIGN.LEFT)
    
    # 四大章节
    chapters = [
        ("01", "方案背景", "集团考核要求 + 一线痛点调研"),
        ("02", "方案目标", "四大核心能力 + 关键指标"),
        ("03", "解决举措", "四大方案 + 技术架构 + 实施路径"),
        ("04", "预期效果", "效率提升 + 体验改善 + 价值量化"),
    ]
    
    for i, (num, title, desc) in enumerate(chapters):
        y = 1.3 + i * 1.5
        
        # 编号圆圈
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(1), Inches(y), Inches(0.8), Inches(0.8)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = COLORS['primary']
        circle.line.fill.background()
        
        # 编号文字
        add_text_box(slide, 1, y + 0.15, 0.8, 0.5,
                     num, font_size=24, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)
        
        # 章节标题
        add_text_box(slide, 2.2, y + 0.1, 4, 0.5,
                     title, font_size=24, bold=True, color=COLORS['text_dark'])
        
        # 描述
        add_text_box(slide, 2.2, y + 0.55, 8, 0.4,
                     desc, font_size=14, color=COLORS['text_gray'])
    
    return slide

# ==================== 第3页：项目背景 ====================
def create_background_slide(prs):
    """项目背景页"""
    slide = add_slide(prs)
    
    # 标题
    add_text_box(slide, 0.5, 0.3, 12.33, 0.6,
                 "方案背景：满足考核解决痛点",
                 font_size=28, bold=True, color=COLORS['dark_blue'])
    
    # 左侧内容区
    add_rectangle(slide, 0.5, 1.2, 6, 5.8, COLORS['bg_light'])
    
    # 目标
    add_text_box(slide, 0.8, 1.5, 5.4, 0.4,
                 "【目标】满足集团考核要求",
                 font_size=18, bold=True, color=COLORS['primary'])
    
    add_text_box(slide, 0.8, 2.0, 5.4, 0.6,
                 "云运〔2025〕4号文件要求：AI嵌入10个关键场景",
                 font_size=14, color=COLORS['text_dark'])
    
    # 举措
    add_text_box(slide, 0.8, 2.8, 5.4, 0.4,
                 "【举措】调研12分公司40+场景",
                 font_size=18, bold=True, color=COLORS['primary'])
    
    # 痛点列表
    pains = [
        "配案时间长：平均30分钟/单",
        "规则查询难：分散在多个系统",
        "错误率高：手工填单易出错",
        "知识获取难：缺乏统一入口"
    ]
    
    for i, pain in enumerate(pains):
        add_text_box(slide, 1.0, 3.3 + i * 0.5, 5, 0.4,
                     f"• {pain}", font_size=14, color=COLORS['text_dark'])
    
    # 右侧图示区
    add_rectangle(slide, 6.8, 1.2, 6, 5.8, COLORS['white'])
    
    # 仪表盘示意
    add_text_box(slide, 7.3, 1.5, 5, 0.5,
                 "集团考核压力", font_size=16, bold=True, color=COLORS['dark_blue'], align=PP_ALIGN.CENTER)
    
    # 三个指标卡片
    metrics = [
        ("10", "AI场景覆盖率", COLORS['primary']),
        ("70", "满意度指标", COLORS['blue']),
        ("31", "省公司排名", COLORS['green']),
    ]
    
    for i, (value, label, color) in enumerate(metrics):
        x = 7.3 + i * 1.8
        add_rounded_rectangle(slide, x, 2.2, 1.5, 1.8, color)
        add_text_box(slide, x, 2.4, 1.5, 0.8,
                     value, font_size=36, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)
        add_text_box(slide, x, 3.2, 1.5, 0.5,
                     label, font_size=11, color=COLORS['white'], align=PP_ALIGN.CENTER)
    
    # 痛点分布
    add_text_box(slide, 7.3, 4.3, 5, 0.4,
                 "一线调研痛点分布", font_size=14, bold=True, color=COLORS['dark_blue'], align=PP_ALIGN.CENTER)
    
    pain_points = [
        ("配案时间长", "35%"),
        ("规则查询难", "28%"),
        ("错误率高", "22%"),
        ("知识获取难", "15%"),
    ]
    
    for i, (name, pct) in enumerate(pain_points):
        y = 4.8 + i * 0.5
        add_rectangle(slide, 7.5, y, 4.5, 0.35, COLORS['bg_light'])
        add_rectangle(slide, 7.5, y, float(pct.replace('%', '')) / 100 * 4.5, 0.35, COLORS['primary'])
        add_text_box(slide, 7.6, y, 3, 0.35,
                     f"{name} {pct}", font_size=11, color=COLORS['white'])
    
    return slide

# ==================== 第4页：问题分析 ====================
def create_problem_slide(prs):
    """问题分析页"""
    slide = add_slide(prs)
    
    # 标题
    add_text_box(slide, 0.5, 0.3, 12.33, 0.6,
                 "问题分析：表象背后的根因",
                 font_size=28, bold=True, color=COLORS['dark_blue'])
    
    # 左侧内容
    add_rectangle(slide, 0.5, 1.2, 6, 5.8, COLORS['bg_light'])
    
    add_text_box(slide, 0.8, 1.5, 5.4, 0.4,
                 "【目标】解决核心根因问题",
                 font_size=18, bold=True, color=COLORS['primary'])
    
    add_text_box(slide, 0.8, 2.0, 5.4, 0.4,
                 "【举措】建设AI智能配案系统",
                 font_size=18, bold=True, color=COLORS['primary'])
    
    add_text_box(slide, 0.8, 2.6, 5.4, 1,
                 "一线营业员无AI助手辅助，手工操作效率低、错误率高、知识获取难",
                 font_size=14, color=COLORS['text_dark'])
    
    # 右侧冰山模型
    add_text_box(slide, 7, 1.5, 5.5, 0.4,
                 "冰山模型分析", font_size=16, bold=True, color=COLORS['dark_blue'], align=PP_ALIGN.CENTER)
    
    # 水面线
    add_rectangle(slide, 6.8, 3.5, 6, 0.05, COLORS['blue'])
    add_text_box(slide, 11.5, 3.3, 1.3, 0.3,
                 "水面线", font_size=10, color=COLORS['blue'])
    
    # 表象层（水面以上）
    add_rounded_rectangle(slide, 7.3, 2.0, 5, 1.3, COLORS['orange'])
    add_text_box(slide, 7.5, 2.1, 4.6, 0.3,
                 "表象（看得见）", font_size=12, bold=True, color=COLORS['white'])
    add_text_box(slide, 7.5, 2.5, 4.6, 0.7,
                 "• 配案时间长\n• 错误率高\n• 客户满意度低",
                 font_size=11, color=COLORS['white'])
    
    # 根因层（水面以下）
    add_rounded_rectangle(slide, 7.3, 3.8, 5, 2.8, COLORS['primary'])
    add_text_box(slide, 7.5, 3.9, 4.6, 0.3,
                 "根因（看不见）", font_size=12, bold=True, color=COLORS['white'])
    add_text_box(slide, 7.5, 4.3, 4.6, 2,
                 "• 缺乏智能化配案工具\n• 规则分散无统一入口\n• 无AI知识问答能力\n• 无智能校验机制\n• 缺少配方案例库",
                 font_size=11, color=COLORS['white'])
    
    return slide

# ==================== 第5页：建设目标 ====================
def create_goal_slide(prs):
    """建设目标页"""
    slide = add_slide(prs)
    
    # 标题
    add_text_box(slide, 0.5, 0.3, 12.33, 0.6,
                 "建设目标：四大核心能力",
                 font_size=28, bold=True, color=COLORS['dark_blue'])
    
    # 四个能力卡片
    capabilities = [
        ("AI知识库", "规则覆盖率\n0%→100%", "📚", COLORS['primary']),
        ("智能推荐", "准确率\n60%→80%", "🎯", COLORS['blue']),
        ("一体化配案", "配案时间\n30分钟→5分钟", "📝", COLORS['green']),
        ("智能校验", "准确率\n85%→98%", "✅", COLORS['orange']),
    ]
    
    for i, (name, metric, icon, color) in enumerate(capabilities):
        x = 0.8 + i * 3.1
        
        # 卡片背景
        add_rounded_rectangle(slide, x, 1.3, 2.8, 4.5, color)
        
        # 图标
        add_text_box(slide, x, 1.5, 2.8, 0.8,
                     icon, font_size=40, color=COLORS['white'], align=PP_ALIGN.CENTER)
        
        # 名称
        add_text_box(slide, x, 2.4, 2.8, 0.5,
                     name, font_size=18, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)
        
        # 分隔线
        add_rectangle(slide, x + 0.3, 3.0, 2.2, 0.02, COLORS['white'])
        
        # 指标
        add_text_box(slide, x, 3.2, 2.8, 1.2,
                     metric, font_size=14, color=COLORS['white'], align=PP_ALIGN.CENTER)
    
    # 关键指标汇总
    add_rectangle(slide, 0.5, 6.0, 12.33, 1.2, COLORS['bg_light'])
    add_text_box(slide, 0.8, 6.2, 11.73, 0.4,
                 "关键指标：配案时间缩短83% | 规则覆盖率100% | 校验准确率98% | 用户满意度90%",
                 font_size=14, bold=True, color=COLORS['dark_blue'], align=PP_ALIGN.CENTER)
    
    return slide

# ==================== 第6页：整体架构 ====================
def create_architecture_slide(prs):
    """整体架构页"""
    slide = add_slide(prs)
    
    # 标题
    add_text_box(slide, 0.5, 0.3, 12.33, 0.6,
                 "整体架构：六层分层设计",
                 font_size=28, bold=True, color=COLORS['dark_blue'])
    
    # 左侧内容
    add_rectangle(slide, 0.5, 1.2, 4.5, 5.8, COLORS['bg_light'])
    
    add_text_box(slide, 0.8, 1.5, 4, 0.4,
                 "【目标】架构清晰可落地",
                 font_size=16, bold=True, color=COLORS['primary'])
    
    add_text_box(slide, 0.8, 2.0, 4, 0.4,
                 "【举措】六层分层设计",
                 font_size=16, bold=True, color=COLORS['primary'])
    
    add_text_box(slide, 0.8, 2.5, 4, 1.5,
                 "便于开发、测试、运维\n支持横向扩展\n标准化接口设计",
                 font_size=13, color=COLORS['text_dark'])
    
    # 右侧六层架构
    layers = [
        ("交互层", "PC端 / 移动端 / 小程序", COLORS['primary']),
        ("智能体层", "问答智能体 / 推荐智能体 / 配案智能体 / 校验智能体", COLORS['blue']),
        ("能力层", "NLP / 知识检索 / 推荐算法 / 规则校验", COLORS['green']),
        ("模型层", "大模型 / 向量模型 / 深度学习模型", COLORS['orange']),
        ("数据层", "MySQL / Milvus / Neo4j / ES / Redis", COLORS['purple']),
        ("集成层", "CRM / 计费 / 订单 / 客户中心 / 随翼选", RgbColor(0xEC, 0x48, 0x99)),
    ]
    
    for i, (name, content, color) in enumerate(layers):
        y = 1.1 + i * 1.0
        
        # 层名
        add_rounded_rectangle(slide, 5.3, y, 2, 0.85, color)
        add_text_box(slide, 5.3, y + 0.25, 2, 0.4,
                     name, font_size=13, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)
        
        # 内容
        add_rectangle(slide, 7.5, y, 5.3, 0.85, COLORS['bg_light'])
        add_text_box(slide, 7.7, y + 0.25, 5, 0.4,
                     content, font_size=11, color=COLORS['text_dark'])
    
    return slide

# ==================== 第7页：方案1-AI知识库 ====================
def create_solution1_slide(prs):
    """方案1-AI知识库"""
    slide = add_slide(prs)
    
    # 标题
    add_text_box(slide, 0.5, 0.3, 12.33, 0.6,
                 "方案1：AI知识库建设",
                 font_size=28, bold=True, color=COLORS['dark_blue'])
    
    # 左侧内容
    add_rectangle(slide, 0.5, 1.2, 5.5, 5.8, COLORS['bg_light'])
    
    add_text_box(slide, 0.8, 1.5, 5, 0.4,
                 "【目标】规则覆盖率0%→100%",
                 font_size=16, bold=True, color=COLORS['primary'])
    
    # 举措列表
    measures = [
        "知识图谱构建",
        "规则录入系统",
        "智能搜索引擎",
        "向量检索能力",
    ]
    
    for i, m in enumerate(measures):
        add_text_box(slide, 0.8, 2.1 + i * 0.5, 5, 0.4,
                     f"• {m}", font_size=14, color=COLORS['text_dark'])
    
    # 右侧流程图
    add_text_box(slide, 6.5, 1.5, 6, 0.4,
                 "技术流程", font_size=16, bold=True, color=COLORS['dark_blue'], align=PP_ALIGN.CENTER)
    
    flow = ["用户提问", "意图识别", "知识检索", "重排序", "大模型生成", "答案输出"]
    
    for i, step in enumerate(flow):
        x = 6.5 + i * 1.1
        add_rounded_rectangle(slide, x, 2.2, 1, 0.6, COLORS['primary'])
        add_text_box(slide, x, 2.35, 1, 0.4,
                     step, font_size=10, color=COLORS['white'], align=PP_ALIGN.CENTER)
        
        if i < len(flow) - 1:
            add_text_box(slide, x + 1, 2.4, 0.1, 0.3,
                         "→", font_size=14, color=COLORS['text_gray'])
    
    # 预期效果
    add_rectangle(slide, 6.5, 3.2, 6, 3.5, COLORS['white'])
    add_text_box(slide, 6.8, 3.4, 5.5, 0.4,
                 "预期效果", font_size=14, bold=True, color=COLORS['dark_blue'])
    
    effects = [
        ("规则覆盖率", "0% → 100%"),
        ("查询响应时间", "10分钟 → 30秒"),
        ("知识准确率", "≥95%"),
    ]
    
    for i, (name, value) in enumerate(effects):
        add_text_box(slide, 6.8, 3.9 + i * 0.7, 3, 0.4,
                     name, font_size=12, color=COLORS['text_dark'])
        add_text_box(slide, 9.8, 3.9 + i * 0.7, 2.5, 0.4,
                     value, font_size=12, bold=True, color=COLORS['primary'])
    
    return slide

# ==================== 第8页：方案2-智能推荐 ====================
def create_solution2_slide(prs):
    """方案2-智能推荐"""
    slide = add_slide(prs)
    
    # 标题
    add_text_box(slide, 0.5, 0.3, 12.33, 0.6,
                 "方案2：智能推荐引擎",
                 font_size=28, bold=True, color=COLORS['dark_blue'])
    
    # 左侧内容
    add_rectangle(slide, 0.5, 1.2, 5.5, 5.8, COLORS['bg_light'])
    
    add_text_box(slide, 0.8, 1.5, 5, 0.4,
                 "【目标】推荐准确率60%→80%",
                 font_size=16, bold=True, color=COLORS['blue'])
    
    measures = [
        "客户画像系统",
        "推荐算法开发",
        "优惠匹配引擎",
        "推荐理由生成",
    ]
    
    for i, m in enumerate(measures):
        add_text_box(slide, 0.8, 2.1 + i * 0.5, 5, 0.4,
                     f"• {m}", font_size=14, color=COLORS['text_dark'])
    
    # 右侧
    add_rectangle(slide, 6.5, 1.2, 6, 5.8, COLORS['white'])
    add_text_box(slide, 6.8, 1.5, 5.5, 0.4,
                 "推荐流程", font_size=14, bold=True, color=COLORS['dark_blue'])
    
    # 推荐流程示意
    steps = [
        ("客户画像", "行为数据分析"),
        ("需求识别", "意图理解"),
        ("套餐匹配", "智能推荐"),
        ("优惠计算", "最优组合"),
    ]
    
    for i, (name, desc) in enumerate(steps):
        y = 2.1 + i * 1.2
        add_rounded_rectangle(slide, 7, y, 5, 0.9, COLORS['blue'])
        add_text_box(slide, 7.2, y + 0.15, 4.6, 0.3,
                     name, font_size=13, bold=True, color=COLORS['white'])
        add_text_box(slide, 7.2, y + 0.45, 4.6, 0.3,
                     desc, font_size=11, color=COLORS['white'])
    
    return slide

# ==================== 第9页：方案3-一体化配案 ====================
def create_solution3_slide(prs):
    """方案3-一体化配案"""
    slide = add_slide(prs)
    
    # 标题
    add_text_box(slide, 0.5, 0.3, 12.33, 0.6,
                 "方案3：一体化配案平台",
                 font_size=28, bold=True, color=COLORS['dark_blue'])
    
    # 左侧
    add_rectangle(slide, 0.5, 1.2, 5.5, 5.8, COLORS['bg_light'])
    
    add_text_box(slide, 0.8, 1.5, 5, 0.4,
                 "【目标】配案时间30分钟→5分钟",
                 font_size=16, bold=True, color=COLORS['green'])
    
    measures = [
        "配案流程引擎",
        "智能填单功能",
        "配方案例库",
        "移动端适配",
    ]
    
    for i, m in enumerate(measures):
        add_text_box(slide, 0.8, 2.1 + i * 0.5, 5, 0.4,
                     f"• {m}", font_size=14, color=COLORS['text_dark'])
    
    # 右侧配案流程
    add_rectangle(slide, 6.5, 1.2, 6, 5.8, COLORS['white'])
    add_text_box(slide, 6.8, 1.5, 5.5, 0.4,
                 "配案流程优化", font_size=14, bold=True, color=COLORS['dark_blue'])
    
    # 对比：优化前后
    add_text_box(slide, 6.8, 2.0, 2.5, 0.3,
                 "优化前", font_size=12, bold=True, color=COLORS['primary'])
    add_text_box(slide, 9.5, 2.0, 2.5, 0.3,
                 "优化后", font_size=12, bold=True, color=COLORS['green'])
    
    comparisons = [
        ("人工查询规则", "AI智能问答"),
        ("手工填写表单", "智能自动填单"),
        ("人工校验合规", "智能规则校验"),
        ("多系统切换", "一体化平台"),
    ]
    
    for i, (before, after) in enumerate(comparisons):
        y = 2.5 + i * 0.8
        add_rectangle(slide, 6.8, y, 2.5, 0.5, COLORS['bg_light'])
        add_text_box(slide, 6.9, y + 0.1, 2.3, 0.3,
                     before, font_size=11, color=COLORS['text_gray'])
        
        add_text_box(slide, 9.4, y + 0.1, 0.3, 0.3,
                     "→", font_size=14, color=COLORS['green'])
        
        add_rectangle(slide, 9.8, y, 2.5, 0.5, COLORS['green'])
        add_text_box(slide, 9.9, y + 0.1, 2.3, 0.3,
                     after, font_size=11, color=COLORS['white'])
    
    return slide

# ==================== 第10页：方案4-智能校验 ====================
def create_solution4_slide(prs):
    """方案4-智能校验"""
    slide = add_slide(prs)
    
    # 标题
    add_text_box(slide, 0.5, 0.3, 12.33, 0.6,
                 "方案4：智能校验系统",
                 font_size=28, bold=True, color=COLORS['dark_blue'])
    
    # 左侧
    add_rectangle(slide, 0.5, 1.2, 5.5, 5.8, COLORS['bg_light'])
    
    add_text_box(slide, 0.8, 1.5, 5, 0.4,
                 "【目标】校验准确率85%→98%",
                 font_size=16, bold=True, color=COLORS['orange'])
    
    measures = [
        "规则校验引擎",
        "风险预警系统",
        "校验日志分析",
        "规则优化建议",
    ]
    
    for i, m in enumerate(measures):
        add_text_box(slide, 0.8, 2.1 + i * 0.5, 5, 0.4,
                     f"• {m}", font_size=14, color=COLORS['text_dark'])
    
    # 右侧决策树示意
    add_rectangle(slide, 6.5, 1.2, 6, 5.8, COLORS['white'])
    add_text_box(slide, 6.8, 1.5, 5.5, 0.4,
                 "校验决策流程", font_size=14, bold=True, color=COLORS['dark_blue'])
    
    # 决策节点
    add_rounded_rectangle(slide, 8.5, 2.0, 2, 0.6, COLORS['orange'])
    add_text_box(slide, 8.5, 2.15, 2, 0.4,
                 "配案数据输入", font_size=11, color=COLORS['white'], align=PP_ALIGN.CENTER)
    
    # 分支
    add_rounded_rectangle(slide, 7, 3.0, 2, 0.6, COLORS['green'])
    add_text_box(slide, 7, 3.15, 2, 0.4,
                 "合规校验", font_size=11, color=COLORS['white'], align=PP_ALIGN.CENTER)
    
    add_rounded_rectangle(slide, 10, 3.0, 2, 0.6, COLORS['primary'])
    add_text_box(slide, 10, 3.15, 2, 0.4,
                 "风险预警", font_size=11, color=COLORS['white'], align=PP_ALIGN.CENTER)
    
    # 结果
    add_rounded_rectangle(slide, 7, 4.2, 2, 0.6, COLORS['blue'])
    add_text_box(slide, 7, 4.35, 2, 0.4,
                 "通过", font_size=11, color=COLORS['white'], align=PP_ALIGN.CENTER)
    
    add_rounded_rectangle(slide, 10, 4.2, 2, 0.6, COLORS['primary'])
    add_text_box(slide, 10, 4.35, 2, 0.4,
                 "拦截+提示", font_size=11, color=COLORS['white'], align=PP_ALIGN.CENTER)
    
    return slide

# ==================== 第11页：系统集成 ====================
def create_integration_slide(prs):
    """系统集成页"""
    slide = add_slide(prs)
    
    # 标题
    add_text_box(slide, 0.5, 0.3, 12.33, 0.6,
                 "系统集成：八大系统互联互通",
                 font_size=28, bold=True, color=COLORS['dark_blue'])
    
    # 中心节点
    add_rounded_rectangle(slide, 5.5, 3, 2.5, 1.2, COLORS['primary'])
    add_text_box(slide, 5.5, 3.3, 2.5, 0.6,
                 "AI智能配案\n系统", font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)
    
    # 周边系统
    systems = [
        ("CRM系统", 2, 1.5),
        ("客户中心", 5.5, 1.5),
        ("随翼选", 9, 1.5),
        ("订单中心", 2, 5),
        ("计费中心", 5.5, 5),
        ("云哨平台", 9, 5),
    ]
    
    for name, x, y in systems:
        add_rounded_rectangle(slide, x, y, 2, 0.8, COLORS['blue'])
        add_text_box(slide, x, y + 0.2, 2, 0.4,
                     name, font_size=12, color=COLORS['white'], align=PP_ALIGN.CENTER)
    
    # API网关说明
    add_rectangle(slide, 0.5, 6.5, 12.33, 0.8, COLORS['bg_light'])
    add_text_box(slide, 0.8, 6.7, 11.73, 0.4,
                 "通过API网关统一集成，标准化接口，降低集成复杂度",
                 font_size=13, color=COLORS['text_dark'], align=PP_ALIGN.CENTER)
    
    return slide

# ==================== 第12页：实施路径 ====================
def create_roadmap_slide(prs):
    """实施路径页"""
    slide = add_slide(prs)
    
    # 标题
    add_text_box(slide, 0.5, 0.3, 12.33, 0.6,
                 "实施路径：四阶段稳步推进",
                 font_size=28, bold=True, color=COLORS['dark_blue'])
    
    # 四个阶段
    phases = [
        ("Q1", "基础能力建设", "AI知识库搭建\n智能问答上线\n规则库完善", COLORS['primary']),
        ("Q2", "核心功能开发", "智能推荐引擎\n智能配案平台\n移动端适配", COLORS['blue']),
        ("Q3", "智能能力增强", "智能校验系统\n配方案例库\n优化迭代", COLORS['green']),
        ("Q4", "全面推广运营", "全省推广\n运营优化\n持续迭代", COLORS['orange']),
    ]
    
    for i, (quarter, name, content, color) in enumerate(phases):
        x = 0.8 + i * 3.1
        
        # 时间节点
        add_rounded_rectangle(slide, x, 1.3, 0.8, 0.8, color)
        add_text_box(slide, x, 1.5, 0.8, 0.5,
                     quarter, font_size=16, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)
        
        # 阶段卡片
        add_rounded_rectangle(slide, x, 2.3, 2.8, 3.8, color)
        add_text_box(slide, x, 2.5, 2.8, 0.5,
                     name, font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)
        add_text_box(slide, x + 0.2, 3.1, 2.4, 2.5,
                     content, font_size=12, color=COLORS['white'])
    
    # 时间轴
    add_rectangle(slide, 1, 6.3, 11.33, 0.05, COLORS['text_gray'])
    add_text_box(slide, 1, 6.5, 11.33, 0.4,
                 "2026年1月                                               2026年12月",
                 font_size=11, color=COLORS['text_gray'], align=PP_ALIGN.CENTER)
    
    return slide

# ==================== 第13页：关键指标 ====================
def create_metrics_slide(prs):
    """关键指标页"""
    slide = add_slide(prs)
    
    # 标题
    add_text_box(slide, 0.5, 0.3, 12.33, 0.6,
                 "关键指标：效果可量化",
                 font_size=28, bold=True, color=COLORS['dark_blue'])
    
    # 指标卡片
    metrics = [
        ("配案时间", "30分钟→5分钟", "↓83%", COLORS['primary']),
        ("规则覆盖率", "0%→100%", "+100%", COLORS['blue']),
        ("校验准确率", "85%→98%", "+13%", COLORS['green']),
        ("用户满意度", "70%→90%", "+20%", COLORS['orange']),
    ]
    
    for i, (name, change, pct, color) in enumerate(metrics):
        x = 0.8 + i * 3.1
        
        add_rounded_rectangle(slide, x, 1.3, 2.8, 2.5, color)
        add_text_box(slide, x, 1.5, 2.8, 0.5,
                     name, font_size=14, color=COLORS['white'], align=PP_ALIGN.CENTER)
        add_text_box(slide, x, 2.1, 2.8, 0.6,
                     change, font_size=12, color=COLORS['white'], align=PP_ALIGN.CENTER)
        add_text_box(slide, x, 2.8, 2.8, 0.6,
                     pct, font_size=20, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)
    
    # 效果汇总表
    add_rectangle(slide, 0.5, 4.2, 12.33, 2.8, COLORS['bg_light'])
    add_text_box(slide, 0.8, 4.4, 11.73, 0.4,
                 "效果汇总", font_size=14, bold=True, color=COLORS['dark_blue'])
    
    # 表头
    headers = ["指标名称", "现状", "目标", "提升幅度"]
    for i, h in enumerate(headers):
        x = 0.8 + i * 2.9
        add_text_box(slide, x, 4.9, 2.7, 0.4,
                     h, font_size=12, bold=True, color=COLORS['text_dark'], align=PP_ALIGN.CENTER)
    
    # 数据行
    rows = [
        ("配案时间", "30分钟", "5分钟", "↓83%"),
        ("规则覆盖率", "0%", "100%", "+100%"),
        ("校验准确率", "85%", "98%", "+13%"),
        ("用户满意度", "70%", "90%", "+20%"),
    ]
    
    for i, row in enumerate(rows):
        y = 5.4 + i * 0.4
        for j, cell in enumerate(row):
            x = 0.8 + j * 2.9
            add_text_box(slide, x, y, 2.7, 0.4,
                         cell, font_size=11, color=COLORS['text_dark'], align=PP_ALIGN.CENTER)
    
    return slide

# ==================== 第14页：核心价值 ====================
def create_value_slide(prs):
    """核心价值页"""
    slide = add_slide(prs)
    
    # 标题
    add_text_box(slide, 0.5, 0.3, 12.33, 0.6,
                 "核心价值：四大价值维度",
                 font_size=28, bold=True, color=COLORS['dark_blue'])
    
    # 价值链
    values = [
        ("考核达标", "满足集团AI场景覆盖率要求", COLORS['primary']),
        ("效率提升", "配案时间缩短83%", COLORS['blue']),
        ("体验改善", "用户满意度提升20%", COLORS['green']),
        ("人力释放", "人力成本降低47%", COLORS['orange']),
    ]
    
    for i, (name, desc, color) in enumerate(values):
        x = 0.8 + i * 3.1
        
        # 价值卡片
        add_rounded_rectangle(slide, x, 1.5, 2.8, 3.5, color)
        add_text_box(slide, x, 1.8, 2.8, 0.6,
                     name, font_size=18, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)
        add_text_box(slide, x + 0.2, 2.6, 2.4, 2,
                     desc, font_size=13, color=COLORS['white'], align=PP_ALIGN.CENTER)
        
        # 连接箭头
        if i < len(values) - 1:
            add_text_box(slide, x + 2.8, 2.8, 0.3, 0.5,
                         "→", font_size=24, color=COLORS['text_gray'])
    
    # 底部总结
    add_rectangle(slide, 0.5, 5.5, 12.33, 1.5, COLORS['bg_light'])
    add_text_box(slide, 0.8, 5.8, 11.73, 0.8,
                 "价值链：满足考核 → 提升效率 → 改善体验 → 释放人力\n多方共赢，全面赋能数字化转型",
                 font_size=14, color=COLORS['text_dark'], align=PP_ALIGN.CENTER)
    
    return slide

# ==================== 第15页：创新亮点 ====================
def create_innovation_slide(prs):
    """创新亮点页"""
    slide = add_slide(prs)
    
    # 标题
    add_text_box(slide, 0.5, 0.3, 12.33, 0.6,
                 "创新亮点：五大技术创新",
                 font_size=28, bold=True, color=COLORS['dark_blue'])
    
    # 创新点
    innovations = [
        ("智能体集群", "四大智能体协同工作", "🤖"),
        ("知识图谱", "结构化知识体系", "📊"),
        ("主动推荐", "预测用户需求", "🎯"),
        ("学习型系统", "持续优化迭代", "📚"),
        ("场景化卡片", "移动端快捷操作", "📱"),
    ]
    
    for i, (name, desc, icon) in enumerate(innovations):
        x = 0.8 + (i % 3) * 4
        y = 1.3 + (i // 3) * 2.5
        
        add_rounded_rectangle(slide, x, y, 3.6, 2, COLORS['purple'])
        add_text_box(slide, x, y + 0.3, 3.6, 0.5,
                     f"{icon} {name}", font_size=16, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)
        add_text_box(slide, x, y + 1, 3.6, 0.6,
                     desc, font_size=12, color=COLORS['white'], align=PP_ALIGN.CENTER)
    
    # 引领行业
    add_rectangle(slide, 0.5, 6.0, 12.33, 1, COLORS['primary'])
    add_text_box(slide, 0.8, 6.2, 11.73, 0.6,
                 "技术领先，打造行业标杆",
                 font_size=16, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)
    
    return slide

# ==================== 第16页：成功案例 ====================
def create_case_slide(prs):
    """成功案例页"""
    slide = add_slide(prs)
    
    # 标题
    add_text_box(slide, 0.5, 0.3, 12.33, 0.6,
                 "成功案例：三地成功实践",
                 font_size=28, bold=True, color=COLORS['dark_blue'])
    
    # 三个案例
    cases = [
        ("陕西电信", "办理时长缩短40%", COLORS['primary']),
        ("重庆电信", "满足考核要求", COLORS['blue']),
        ("安徽电信", "31省标杆", COLORS['green']),
    ]
    
    for i, (name, result, color) in enumerate(cases):
        x = 0.8 + i * 4.1
        
        add_rounded_rectangle(slide, x, 1.5, 3.8, 3, color)
        add_text_box(slide, x, 1.8, 3.8, 0.6,
                     name, font_size=20, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)
        add_text_box(slide, x, 2.6, 3.8, 0.8,
                     result, font_size=14, color=COLORS['white'], align=PP_ALIGN.CENTER)
    
    # 复制经验
    add_rectangle(slide, 0.5, 5, 12.33, 2, COLORS['bg_light'])
    add_text_box(slide, 0.8, 5.3, 11.73, 0.4,
                 "复制成熟经验，降低风险", font_size=14, bold=True, color=COLORS['dark_blue'])
    add_text_box(slide, 0.8, 5.8, 11.73, 0.8,
                 "• 技术方案成熟，可直接复用\n• 实施经验丰富，风险可控\n• 效果可预期，投资回报明确",
                 font_size=12, color=COLORS['text_dark'])
    
    return slide

# ==================== 第17页：总结展望 ====================
def create_summary_slide(prs):
    """总结展望页"""
    slide = add_slide(prs)
    
    # 标题
    add_text_box(slide, 0.5, 0.3, 12.33, 0.6,
                 "总结与展望",
                 font_size=28, bold=True, color=COLORS['dark_blue'])
    
    # 四大关键点
    points = [
        ("方案成熟", "四步方案，层层递进"),
        ("技术可行", "六层架构，清晰落地"),
        ("价值明确", "四大价值，多方共赢"),
        ("风险可控", "分阶段实施，稳步推进"),
    ]
    
    for i, (name, desc) in enumerate(points):
        x = 0.8 + (i % 2) * 6.3
        y = 1.3 + (i // 2) * 2.5
        
        add_rounded_rectangle(slide, x, y, 5.8, 2, COLORS['blue'])
        add_text_box(slide, x + 0.3, y + 0.3, 5.2, 0.5,
                     f"{i+1}. {name}", font_size=18, bold=True, color=COLORS['white'])
        add_text_box(slide, x + 0.3, y + 1, 5.2, 0.6,
                     desc, font_size=14, color=COLORS['white'])
    
    # 下一步
    add_rectangle(slide, 0.5, 6, 12.33, 1, COLORS['primary'])
    add_text_box(slide, 0.8, 6.2, 11.73, 0.6,
                 "下一步：启动项目立项，分阶段实施，持续优化迭代",
                 font_size=14, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)
    
    return slide

# ==================== 第18页：结束页 ====================
def create_end_slide(prs):
    """结束页"""
    slide = add_slide(prs)
    
    # 背景
    add_rectangle(slide, 0, 0, 13.33, 7.5, COLORS['primary'])
    
    # 感谢
    add_text_box(slide, 1, 2.5, 11.33, 1,
                 "谢谢！",
                 font_size=56, bold=True, color=COLORS['white'], align=PP_ALIGN.CENTER)
    
    # 口号
    add_text_box(slide, 1, 4, 11.33, 0.8,
                 "让我们一起，用AI赋能一线，提升效率，改善体验！",
                 font_size=24, color=COLORS['white'], align=PP_ALIGN.CENTER)
    
    # 签名
    add_text_box(slide, 1, 5.5, 11.33, 0.5,
                 "南有乔木，不可休思",
                 font_size=16, color=COLORS['white'], align=PP_ALIGN.CENTER)
    
    return slide

# ==================== 主函数 ====================
def main():
    """主函数"""
    print("开始生成PPT...")
    
    # 创建演示文稿
    prs = create_presentation()
    
    # 生成所有页面
    print("生成第1页：封面...")
    create_cover_slide(prs)
    
    print("生成第2页：目录...")
    create_toc_slide(prs)
    
    print("生成第3页：项目背景...")
    create_background_slide(prs)
    
    print("生成第4页：问题分析...")
    create_problem_slide(prs)
    
    print("生成第5页：建设目标...")
    create_goal_slide(prs)
    
    print("生成第6页：整体架构...")
    create_architecture_slide(prs)
    
    print("生成第7页：方案1-AI知识库...")
    create_solution1_slide(prs)
    
    print("生成第8页：方案2-智能推荐...")
    create_solution2_slide(prs)
    
    print("生成第9页：方案3-一体化配案...")
    create_solution3_slide(prs)
    
    print("生成第10页：方案4-智能校验...")
    create_solution4_slide(prs)
    
    print("生成第11页：系统集成...")
    create_integration_slide(prs)
    
    print("生成第12页：实施路径...")
    create_roadmap_slide(prs)
    
    print("生成第13页：关键指标...")
    create_metrics_slide(prs)
    
    print("生成第14页：核心价值...")
    create_value_slide(prs)
    
    print("生成第15页：创新亮点...")
    create_innovation_slide(prs)
    
    print("生成第16页：成功案例...")
    create_case_slide(prs)
    
    print("生成第17页：总结展望...")
    create_summary_slide(prs)
    
    print("生成第18页：结束页...")
    create_end_slide(prs)
    
    # 保存文件
    output_path = "/root/.openclaw/workspace/02_工作台/工程联动建议/案例库/模拟案例-AI智能配案系统/步骤4-方案PPT/方案PPT-V4.0-模板规范版.pptx"
    prs.save(output_path)
    print(f"PPT已保存到: {output_path}")
    
    return output_path

if __name__ == "__main__":
    main()

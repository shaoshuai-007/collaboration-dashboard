#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
NVIDIA GTC 2026 科技风格PPT生成器
风格：深色背景 + 霓虹科技感
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
import os

# 颜色定义 - NVIDIA科技风格
COLORS = {
    'background': RGBColor(0x0A, 0x0A, 0x0A),      # 深黑背景
    'dark_blue': RGBColor(0x0D, 0x1B, 0x2A),       # 深蓝背景
    'nvidia_green': RGBColor(0x76, 0xB9, 0x00),    # NVIDIA绿
    'neon_blue': RGBColor(0x00, 0xD4, 0xFF),       # 霓虹蓝
    'neon_purple': RGBColor(0x9D, 0x4E, 0xDD),     # 霓虹紫
    'neon_cyan': RGBColor(0x00, 0xFF, 0xE5),       # 霓虹青
    'white': RGBColor(0xFF, 0xFF, 0xFF),           # 白色
    'gray': RGBColor(0xAA, 0xAA, 0xAA),            # 灰色
    'light_gray': RGBColor(0xCC, 0xCC, 0xCC),      # 浅灰
}

def set_slide_background(slide, color):
    """设置幻灯片背景色"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, text, left, top, width, height, font_size=18, 
                 font_color=COLORS['white'], bold=False, alignment=PP_ALIGN.LEFT):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = 'Arial'
    p.alignment = alignment
    return txBox

def add_title_text(slide, text, left, top, width, height, font_size=44, color=COLORS['white']):
    """添加标题文本"""
    return add_text_box(slide, text, left, top, width, height, 
                       font_size=font_size, font_color=color, bold=True)

def add_accent_line(slide, left, top, width, height, color=COLORS['nvidia_green']):
    """添加装饰线条"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_bullet_points(slide, items, left, top, width, height, font_size=18, color=COLORS['light_gray']):
    """添加项目符号列表"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Arial'
        p.space_after = Pt(8)
    return txBox

def add_glow_rectangle(slide, left, top, width, height, fill_color, line_color):
    """添加发光效果矩形框"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color
    shape.line.width = Pt(2)
    return shape

def create_presentation():
    """创建PPT演示文稿"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 使用空白布局
    blank_layout = prs.slide_layouts[6]
    
    # ==================== 第1页：封面 ====================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, COLORS['background'])
    
    # 顶部装饰线
    add_accent_line(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.05), COLORS['nvidia_green'])
    
    # 主标题
    add_title_text(slide, "NVIDIA GTC 2026", Inches(0.5), Inches(2.2), 
                  Inches(12), Inches(1), font_size=72, color=COLORS['nvidia_green'])
    
    # 副标题
    add_text_box(slide, "黄仁勋主题演讲核心内容解读", Inches(0.5), Inches(3.3), 
                Inches(12), Inches(0.6), font_size=32, font_color=COLORS['neon_blue'])
    
    # 装饰线
    add_accent_line(slide, Inches(0.5), Inches(4.0), Inches(3), Inches(0.03), COLORS['neon_blue'])
    
    # 日期
    add_text_box(slide, "2026年3月", Inches(0.5), Inches(5.5), 
                Inches(4), Inches(0.5), font_size=20, font_color=COLORS['gray'])
    
    # 底部信息
    add_text_box(slide, "GPU Technology Conference | AI Technology Summit", 
                Inches(0.5), Inches(6.8), Inches(12), Inches(0.4), 
                font_size=16, font_color=COLORS['gray'])
    
    # ==================== 第2页：目录 ====================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, COLORS['dark_blue'])
    
    add_title_text(slide, "目录 CONTENTS", Inches(0.5), Inches(0.5), 
                  Inches(12), Inches(0.8), font_size=40, color=COLORS['nvidia_green'])
    add_accent_line(slide, Inches(0.5), Inches(1.3), Inches(2), Inches(0.03), COLORS['nvidia_green'])
    
    contents = [
        ("01", "GTC大会背景与影响力"),
        ("02", "NVIDIA公司介绍与市场地位"),
        ("03", "黄仁勋个人风格与愿景"),
        ("04", "Vera Rubin超级计算机"),
        ("05", "Groq LPU推理革命"),
        ("06", "AI应用场景"),
        ("07", "生态系统建设"),
        ("08", "行业影响与未来展望"),
    ]
    
    for i, (num, title) in enumerate(contents):
        y_pos = 1.8 + i * 0.65
        # 编号
        add_text_box(slide, num, Inches(0.5), Inches(y_pos), Inches(0.8), Inches(0.5),
                    font_size=24, font_color=COLORS['nvidia_green'], bold=True)
        # 标题
        add_text_box(slide, title, Inches(1.5), Inches(y_pos), Inches(10), Inches(0.5),
                    font_size=22, font_color=COLORS['white'])
    
    # ==================== 第3页：GTC大会背景 ====================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, COLORS['background'])
    
    add_title_text(slide, "GTC大会背景与影响力", Inches(0.5), Inches(0.5), 
                  Inches(12), Inches(0.8), font_size=40, color=COLORS['nvidia_green'])
    add_accent_line(slide, Inches(0.5), Inches(1.3), Inches(2.5), Inches(0.03), COLORS['nvidia_green'])
    
    items = [
        "GTC (GPU Technology Conference) - NVIDIA主办的全球顶级AI开发者大会",
        "每年吸引数万名开发者、研究人员、企业决策者参与",
        "被誉为'AI领域的TED大会'和'GPU计算领域的风向标'",
        "全球AI技术发展的晴雨表，新技术、新产品的首发平台",
        "产学研交流的重要桥梁，推动行业创新与合作"
    ]
    add_bullet_points(slide, items, Inches(0.5), Inches(1.8), Inches(12), Inches(5), 
                     font_size=20, color=COLORS['light_gray'])
    
    # 右侧装饰框
    add_glow_rectangle(slide, Inches(9), Inches(4.5), Inches(3.8), Inches(2.2), 
                      COLORS['dark_blue'], COLORS['neon_blue'])
    add_text_box(slide, "行业地位", Inches(9.2), Inches(4.7), Inches(3.4), Inches(0.4),
                font_size=18, font_color=COLORS['neon_blue'], bold=True)
    add_text_box(slide, "• AI技术风向标\n• 首发平台\n• 全球影响力", 
                Inches(9.2), Inches(5.2), Inches(3.4), Inches(1.3),
                font_size=16, font_color=COLORS['light_gray'])
    
    # ==================== 第4页：NVIDIA公司介绍 ====================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, COLORS['dark_blue'])
    
    add_title_text(slide, "NVIDIA公司介绍与市场地位", Inches(0.5), Inches(0.5), 
                  Inches(12), Inches(0.8), font_size=40, color=COLORS['nvidia_green'])
    add_accent_line(slide, Inches(0.5), Inches(1.3), Inches(2.5), Inches(0.03), COLORS['nvidia_green'])
    
    # 左栏：公司概况
    add_glow_rectangle(slide, Inches(0.5), Inches(1.8), Inches(5.8), Inches(2.5), 
                      COLORS['background'], COLORS['nvidia_green'])
    add_text_box(slide, "公司概况", Inches(0.7), Inches(2.0), Inches(5.4), Inches(0.4),
                font_size=20, font_color=COLORS['nvidia_green'], bold=True)
    add_text_box(slide, "• 成立时间：1993年\n• 创始人：黄仁勋等\n• 总部：美国加州圣克拉拉", 
                Inches(0.7), Inches(2.5), Inches(5.4), Inches(1.5),
                font_size=18, font_color=COLORS['light_gray'])
    
    # 右栏：市场地位
    add_glow_rectangle(slide, Inches(6.8), Inches(1.8), Inches(6), Inches(2.5), 
                      COLORS['background'], COLORS['neon_blue'])
    add_text_box(slide, "市场地位", Inches(7), Inches(2.0), Inches(5.6), Inches(0.4),
                font_size=20, font_color=COLORS['neon_blue'], bold=True)
    add_text_box(slide, "• GPU市场份额：>80%\n• AI训练芯片：>90%\n• 全球市值最高半导体公司之一", 
                Inches(7), Inches(2.5), Inches(5.6), Inches(1.5),
                font_size=18, font_color=COLORS['light_gray'])
    
    # 底部：核心优势
    add_text_box(slide, "核心优势", Inches(0.5), Inches(4.6), Inches(12), Inches(0.4),
                font_size=22, font_color=COLORS['neon_purple'], bold=True)
    add_text_box(slide, "硬件：GPU架构设计与制造  |  软件：CUDA生态系统  |  生态：开发者社区与合作伙伴网络", 
                Inches(0.5), Inches(5.2), Inches(12), Inches(1),
                font_size=18, font_color=COLORS['light_gray'])
    
    # ==================== 第5页：黄仁勋个人风格 ====================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, COLORS['background'])
    
    add_title_text(slide, "黄仁勋个人风格与愿景", Inches(0.5), Inches(0.5), 
                  Inches(12), Inches(0.8), font_size=40, color=COLORS['nvidia_green'])
    add_accent_line(slide, Inches(0.5), Inches(1.3), Inches(2.5), Inches(0.03), COLORS['nvidia_green'])
    
    # 个人风格
    add_glow_rectangle(slide, Inches(0.5), Inches(1.8), Inches(6), Inches(2.3), 
                      COLORS['dark_blue'], COLORS['neon_cyan'])
    add_text_box(slide, "标志性风格", Inches(0.7), Inches(2.0), Inches(5.6), Inches(0.4),
                font_size=20, font_color=COLORS['neon_cyan'], bold=True)
    add_text_box(slide, "• 标志性黑色皮衣形象\n• 激情洋溢的演讲风格\n• 技术远见与商业洞察兼具", 
                Inches(0.7), Inches(2.5), Inches(5.6), Inches(1.3),
                font_size=18, font_color=COLORS['light_gray'])
    
    # 愿景理念
    add_glow_rectangle(slide, Inches(6.8), Inches(1.8), Inches(6), Inches(2.3), 
                      COLORS['dark_blue'], COLORS['neon_purple'])
    add_text_box(slide, "愿景理念", Inches(7), Inches(2.0), Inches(5.6), Inches(0.4),
                font_size=20, font_color=COLORS['neon_purple'], bold=True)
    add_text_box(slide, "• 'AI是新时代的电力'\n• 推动AI民主化\n• 构建可持续计算基础设施", 
                Inches(7), Inches(2.5), Inches(5.6), Inches(1.3),
                font_size=18, font_color=COLORS['light_gray'])
    
    # 名言
    add_text_box(slide, '"AI是新时代的电力，将改变每一个行业"', 
                Inches(0.5), Inches(4.5), Inches(12), Inches(0.8),
                font_size=28, font_color=COLORS['nvidia_green'], bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, "— Jensen Huang, NVIDIA CEO", 
                Inches(0.5), Inches(5.3), Inches(12), Inches(0.5),
                font_size=18, font_color=COLORS['gray'], alignment=PP_ALIGN.CENTER)
    
    # ==================== 第6页：Vera Rubin超级计算机 ====================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, COLORS['dark_blue'])
    
    add_title_text(slide, "Vera Rubin 七合一超级计算机", Inches(0.5), Inches(0.5), 
                  Inches(12), Inches(0.8), font_size=40, color=COLORS['nvidia_green'])
    add_accent_line(slide, Inches(0.5), Inches(1.3), Inches(2.5), Inches(0.03), COLORS['nvidia_green'])
    
    # 项目背景
    add_text_box(slide, "以首位发现暗物质的女天文学家Vera Rubin命名，NVIDIA在超级计算领域的最新力作", 
                Inches(0.5), Inches(1.6), Inches(12), Inches(0.6),
                font_size=18, font_color=COLORS['gray'])
    
    # 核心特性
    add_text_box(slide, "核心特性", Inches(0.5), Inches(2.3), Inches(12), Inches(0.4),
                font_size=22, font_color=COLORS['neon_blue'], bold=True)
    
    features = [
        "七合一架构：集成CPU、GPU、DPU、存储、网络、安全、管理",
        "性能突破：计算性能较上一代提升数倍",
        "能效优化：单位功耗性能大幅提升",
        "AI专用：为AI训练与推理深度优化"
    ]
    add_bullet_points(slide, features, Inches(0.5), Inches(2.8), Inches(12), Inches(2.5), 
                     font_size=18, color=COLORS['light_gray'])
    
    # 应用场景
    add_text_box(slide, "应用场景", Inches(0.5), Inches(5.0), Inches(12), Inches(0.4),
                font_size=22, font_color=COLORS['neon_purple'], bold=True)
    add_text_box(slide, "大规模AI模型训练  |  科学研究模拟  |  气候预测与建模", 
                Inches(0.5), Inches(5.5), Inches(12), Inches(0.8),
                font_size=18, font_color=COLORS['light_gray'])
    
    # ==================== 第7页：Groq LPU推理革命 ====================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, COLORS['background'])
    
    add_title_text(slide, "Groq LPU 推理革命", Inches(0.5), Inches(0.5), 
                  Inches(12), Inches(0.8), font_size=40, color=COLORS['nvidia_green'])
    add_accent_line(slide, Inches(0.5), Inches(1.3), Inches(2.5), Inches(0.03), COLORS['nvidia_green'])
    
    # LPU特性
    add_glow_rectangle(slide, Inches(0.5), Inches(1.8), Inches(6), Inches(2.8), 
                      COLORS['dark_blue'], COLORS['neon_cyan'])
    add_text_box(slide, "LPU (Language Processing Unit)", Inches(0.7), Inches(2.0), Inches(5.6), Inches(0.4),
                font_size=18, font_color=COLORS['neon_cyan'], bold=True)
    add_text_box(slide, "• 设计理念：专为语言模型推理优化\n• 性能特点：推理速度远超传统GPU\n• 应用价值：实时AI响应，降低延迟", 
                Inches(0.7), Inches(2.5), Inches(5.6), Inches(1.8),
                font_size=16, font_color=COLORS['light_gray'])
    
    # 行业意义
    add_glow_rectangle(slide, Inches(6.8), Inches(1.8), Inches(6), Inches(2.8), 
                      COLORS['dark_blue'], COLORS['neon_purple'])
    add_text_box(slide, "行业意义", Inches(7), Inches(2.0), Inches(5.6), Inches(0.4),
                font_size=18, font_color=COLORS['neon_purple'], bold=True)
    add_text_box(slide, "• 推动AI推理从数据中心走向边缘\n• 实时AI应用成为可能\n• AI推理成本大幅降低", 
                Inches(7), Inches(2.5), Inches(5.6), Inches(1.8),
                font_size=16, font_color=COLORS['light_gray'])
    
    # 关键数据
    add_text_box(slide, "推理速度提升", Inches(0.5), Inches(5.0), Inches(4), Inches(0.4),
                font_size=18, font_color=COLORS['nvidia_green'], bold=True)
    add_text_box(slide, "10x+", Inches(0.5), Inches(5.4), Inches(4), Inches(0.8),
                font_size=48, font_color=COLORS['nvidia_green'], bold=True)
    
    add_text_box(slide, "延迟降低", Inches(4.5), Inches(5.0), Inches(4), Inches(0.4),
                font_size=18, font_color=COLORS['neon_blue'], bold=True)
    add_text_box(slide, "90%", Inches(4.5), Inches(5.4), Inches(4), Inches(0.8),
                font_size=48, font_color=COLORS['neon_blue'], bold=True)
    
    add_text_box(slide, "成本节省", Inches(8.5), Inches(5.0), Inches(4), Inches(0.4),
                font_size=18, font_color=COLORS['neon_purple'], bold=True)
    add_text_box(slide, "50%+", Inches(8.5), Inches(5.4), Inches(4), Inches(0.8),
                font_size=48, font_color=COLORS['neon_purple'], bold=True)
    
    # ==================== 第8页：AI应用场景概览 ====================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, COLORS['dark_blue'])
    
    add_title_text(slide, "AI应用场景", Inches(0.5), Inches(0.5), 
                  Inches(12), Inches(0.8), font_size=40, color=COLORS['nvidia_green'])
    add_accent_line(slide, Inches(0.5), Inches(1.3), Inches(2.5), Inches(0.03), COLORS['nvidia_green'])
    
    # 六大应用场景卡片
    apps = [
        ("自动驾驶", "车载AI计算\n智能交通系统", COLORS['neon_cyan']),
        ("医疗健康", "医学影像诊断\n药物研发加速", COLORS['nvidia_green']),
        ("工业制造", "智能工厂\n工业质检AI", COLORS['neon_blue']),
        ("科学研究", "气候建模\n新材料发现", COLORS['neon_purple']),
        ("创意内容", "AI生成媒体\n虚拟人技术", COLORS['nvidia_green']),
        ("金融服务", "智能风控\n量化交易", COLORS['neon_cyan']),
    ]
    
    for i, (title, desc, color) in enumerate(apps):
        col = i % 3
        row = i // 3
        x = 0.5 + col * 4.2
        y = 1.8 + row * 2.6
        
        add_glow_rectangle(slide, Inches(x), Inches(y), Inches(3.9), Inches(2.2), 
                          COLORS['background'], color)
        add_text_box(slide, title, Inches(x + 0.2), Inches(y + 0.2), Inches(3.5), Inches(0.5),
                    font_size=20, font_color=color, bold=True)
        add_text_box(slide, desc, Inches(x + 0.2), Inches(y + 0.8), Inches(3.5), Inches(1.2),
                    font_size=16, font_color=COLORS['light_gray'])
    
    # ==================== 第9页：自动驾驶AI ====================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, COLORS['background'])
    
    add_title_text(slide, "自动驾驶 AI", Inches(0.5), Inches(0.5), 
                  Inches(12), Inches(0.8), font_size=40, color=COLORS['nvidia_green'])
    add_accent_line(slide, Inches(0.5), Inches(1.3), Inches(2.5), Inches(0.03), COLORS['nvidia_green'])
    
    items = [
        "NVIDIA Drive平台：端到端自动驾驶解决方案",
        "车载AI芯片：高性能、低功耗的边缘计算能力",
        "传感器融合：摄像头、雷达、激光雷达多模态感知",
        "深度学习模型：实时目标检测与路径规划",
        "仿真测试：虚拟环境中训练数百万公里",
        "智能交通：车路协同与城市智慧交通系统"
    ]
    add_bullet_points(slide, items, Inches(0.5), Inches(1.8), Inches(12), Inches(5), 
                     font_size=20, color=COLORS['light_gray'])
    
    # ==================== 第10页：医疗健康AI ====================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, COLORS['dark_blue'])
    
    add_title_text(slide, "医疗健康 AI", Inches(0.5), Inches(0.5), 
                  Inches(12), Inches(0.8), font_size=40, color=COLORS['nvidia_green'])
    add_accent_line(slide, Inches(0.5), Inches(1.3), Inches(2.5), Inches(0.03), COLORS['nvidia_green'])
    
    # 医学影像
    add_glow_rectangle(slide, Inches(0.5), Inches(1.8), Inches(4), Inches(2.3), 
                      COLORS['background'], COLORS['neon_cyan'])
    add_text_box(slide, "医学影像AI", Inches(0.7), Inches(2.0), Inches(3.6), Inches(0.4),
                font_size=18, font_color=COLORS['neon_cyan'], bold=True)
    add_text_box(slide, "• CT/MRI智能分析\n• 病灶自动检测\n• 辅助诊断系统", 
                Inches(0.7), Inches(2.5), Inches(3.6), Inches(1.3),
                font_size=16, font_color=COLORS['light_gray'])
    
    # 药物研发
    add_glow_rectangle(slide, Inches(4.7), Inches(1.8), Inches(4), Inches(2.3), 
                      COLORS['background'], COLORS['nvidia_green'])
    add_text_box(slide, "药物研发", Inches(4.9), Inches(2.0), Inches(3.6), Inches(0.4),
                font_size=18, font_color=COLORS['nvidia_green'], bold=True)
    add_text_box(slide, "• 分子结构预测\n• 药物筛选加速\n• 临床试验优化", 
                Inches(4.9), Inches(2.5), Inches(3.6), Inches(1.3),
                font_size=16, font_color=COLORS['light_gray'])
    
    # 基因测序
    add_glow_rectangle(slide, Inches(8.9), Inches(1.8), Inches(4), Inches(2.3), 
                      COLORS['background'], COLORS['neon_purple'])
    add_text_box(slide, "基因测序", Inches(9.1), Inches(2.0), Inches(3.6), Inches(0.4),
                font_size=18, font_color=COLORS['neon_purple'], bold=True)
    add_text_box(slide, "• 基因组快速分析\n• 疾病风险预测\n• 个性化医疗", 
                Inches(9.1), Inches(2.5), Inches(3.6), Inches(1.3),
                font_size=16, font_color=COLORS['light_gray'])
    
    # 数据
    add_text_box(slide, "药物研发周期缩短", Inches(0.5), Inches(4.5), Inches(6), Inches(0.4),
                font_size=18, font_color=COLORS['neon_blue'], bold=True)
    add_text_box(slide, "从10年 → 2年", Inches(0.5), Inches(5.0), Inches(6), Inches(0.6),
                font_size=32, font_color=COLORS['neon_blue'], bold=True)
    
    add_text_box(slide, "影像诊断准确率", Inches(6.8), Inches(4.5), Inches(6), Inches(0.4),
                font_size=18, font_color=COLORS['nvidia_green'], bold=True)
    add_text_box(slide, ">95%", Inches(6.8), Inches(5.0), Inches(6), Inches(0.6),
                font_size=32, font_color=COLORS['nvidia_green'], bold=True)
    
    # ==================== 第11页：工业制造AI ====================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, COLORS['background'])
    
    add_title_text(slide, "工业制造 AI", Inches(0.5), Inches(0.5), 
                  Inches(12), Inches(0.8), font_size=40, color=COLORS['nvidia_green'])
    add_accent_line(slide, Inches(0.5), Inches(1.3), Inches(2.5), Inches(0.03), COLORS['nvidia_green'])
    
    items = [
        "智能工厂：全流程自动化与智能调度",
        "工业质检：AI视觉检测，缺陷识别率>99%",
        "预测性维护：设备故障预测，减少停机时间",
        "数字孪生：虚拟工厂模拟与优化",
        "供应链优化：AI驱动的库存与物流管理",
        "质量控制：实时监测与自动调整"
    ]
    add_bullet_points(slide, items, Inches(0.5), Inches(1.8), Inches(12), Inches(5), 
                     font_size=20, color=COLORS['light_gray'])
    
    # ==================== 第12页：CUDA生态系统 ====================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, COLORS['dark_blue'])
    
    add_title_text(slide, "CUDA生态系统", Inches(0.5), Inches(0.5), 
                  Inches(12), Inches(0.8), font_size=40, color=COLORS['nvidia_green'])
    add_accent_line(slide, Inches(0.5), Inches(1.3), Inches(2.5), Inches(0.03), COLORS['nvidia_green'])
    
    # 开发者数据
    add_glow_rectangle(slide, Inches(0.5), Inches(1.8), Inches(4), Inches(2), 
                      COLORS['background'], COLORS['nvidia_green'])
    add_text_box(slide, "全球开发者", Inches(0.7), Inches(2.0), Inches(3.6), Inches(0.4),
                font_size=18, font_color=COLORS['nvidia_green'], bold=True)
    add_text_box(slide, "500万+", Inches(0.7), Inches(2.5), Inches(3.6), Inches(0.8),
                font_size=40, font_color=COLORS['nvidia_green'], bold=True)
    
    # 工具链
    add_glow_rectangle(slide, Inches(4.7), Inches(1.8), Inches(4), Inches(2), 
                      COLORS['background'], COLORS['neon_blue'])
    add_text_box(slide, "开发工具", Inches(4.9), Inches(2.0), Inches(3.6), Inches(0.4),
                font_size=18, font_color=COLORS['neon_blue'], bold=True)
    add_text_box(slide, "编译器\n调试器\n性能分析", Inches(4.9), Inches(2.4), Inches(3.6), Inches(1.2),
                font_size=16, font_color=COLORS['light_gray'])
    
    # 框架支持
    add_glow_rectangle(slide, Inches(8.9), Inches(1.8), Inches(4), Inches(2), 
                      COLORS['background'], COLORS['neon_purple'])
    add_text_box(slide, "主流框架", Inches(9.1), Inches(2.0), Inches(3.6), Inches(0.4),
                font_size=18, font_color=COLORS['neon_purple'], bold=True)
    add_text_box(slide, "PyTorch\nTensorFlow\nJAX", Inches(9.1), Inches(2.4), Inches(3.6), Inches(1.2),
                font_size=16, font_color=COLORS['light_gray'])
    
    # 合作伙伴
    add_text_box(slide, "合作伙伴生态", Inches(0.5), Inches(4.3), Inches(12), Inches(0.4),
                font_size=22, font_color=COLORS['neon_cyan'], bold=True)
    add_text_box(slide, "云服务商：AWS | Azure | Google Cloud | 阿里云\n硬件厂商：服务器制造商、存储厂商\n软件厂商：AI应用开发商", 
                Inches(0.5), Inches(4.8), Inches(12), Inches(1.5),
                font_size=18, font_color=COLORS['light_gray'])
    
    # ==================== 第13页：开发者社区 ====================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, COLORS['background'])
    
    add_title_text(slide, "开发者社区与支持", Inches(0.5), Inches(0.5), 
                  Inches(12), Inches(0.8), font_size=40, color=COLORS['nvidia_green'])
    add_accent_line(slide, Inches(0.5), Inches(1.3), Inches(2.5), Inches(0.03), COLORS['nvidia_green'])
    
    items = [
        "NVIDIA Developer计划：免费注册，获取SDK和工具",
        "技术培训：在线课程、研讨会、认证项目",
        "开源支持：贡献开源项目，共享技术成果",
        "开发者论坛：技术问答、经验分享",
        "Hackathon活动：创新竞赛，奖金支持",
        "技术文档：详尽的API文档与教程"
    ]
    add_bullet_points(slide, items, Inches(0.5), Inches(1.8), Inches(12), Inches(5), 
                     font_size=20, color=COLORS['light_gray'])
    
    # ==================== 第14页：行业影响 ====================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, COLORS['dark_blue'])
    
    add_title_text(slide, "行业影响", Inches(0.5), Inches(0.5), 
                  Inches(12), Inches(0.8), font_size=40, color=COLORS['nvidia_green'])
    add_accent_line(slide, Inches(0.5), Inches(1.3), Inches(2.5), Inches(0.03), COLORS['nvidia_green'])
    
    # 四大行业
    industries = [
        ("芯片产业", "确立AI芯片技术方向\n推动架构创新\n带动产业链发展", COLORS['neon_cyan']),
        ("云计算", "GPU云服务成为标配\nAI即服务兴起\n云原生AI开发范式", COLORS['nvidia_green']),
        ("企业应用", "AI助力数字化转型\n智能化运营成为趋势\nAI人才需求激增", COLORS['neon_blue']),
        ("科研创新", "强大算力支撑\n加速科学发现\n跨学科合作", COLORS['neon_purple']),
    ]
    
    for i, (title, desc, color) in enumerate(industries):
        col = i % 2
        row = i // 2
        x = 0.5 + col * 6.4
        y = 1.8 + row * 2.6
        
        add_glow_rectangle(slide, Inches(x), Inches(y), Inches(6), Inches(2.2), 
                          COLORS['background'], color)
        add_text_box(slide, title, Inches(x + 0.2), Inches(y + 0.2), Inches(5.6), Inches(0.5),
                    font_size=20, font_color=color, bold=True)
        add_text_box(slide, desc, Inches(x + 0.2), Inches(y + 0.8), Inches(5.6), Inches(1.2),
                    font_size=16, font_color=COLORS['light_gray'])
    
    # ==================== 第15页：未来展望 ====================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, COLORS['background'])
    
    add_title_text(slide, "未来展望", Inches(0.5), Inches(0.5), 
                  Inches(12), Inches(0.8), font_size=40, color=COLORS['nvidia_green'])
    add_accent_line(slide, Inches(0.5), Inches(1.3), Inches(2.5), Inches(0.03), COLORS['nvidia_green'])
    
    # 技术方向
    add_glow_rectangle(slide, Inches(0.5), Inches(1.8), Inches(4), Inches(2.5), 
                      COLORS['dark_blue'], COLORS['neon_cyan'])
    add_text_box(slide, "技术方向", Inches(0.7), Inches(2.0), Inches(3.6), Inches(0.4),
                font_size=18, font_color=COLORS['neon_cyan'], bold=True)
    add_text_box(slide, "• 通用人工智能(AGI)\n• 多模态大模型\n• AI推理专用芯片", 
                Inches(0.7), Inches(2.5), Inches(3.6), Inches(1.5),
                font_size=16, font_color=COLORS['light_gray'])
    
    # 社会影响
    add_glow_rectangle(slide, Inches(4.7), Inches(1.8), Inches(4), Inches(2.5), 
                      COLORS['dark_blue'], COLORS['nvidia_green'])
    add_text_box(slide, "社会影响", Inches(4.9), Inches(2.0), Inches(3.6), Inches(0.4),
                font_size=18, font_color=COLORS['nvidia_green'], bold=True)
    add_text_box(slide, "• 创造新就业机会\n• 人机协作成常态\n• AI伦理与治理", 
                Inches(4.9), Inches(2.5), Inches(3.6), Inches(1.5),
                font_size=16, font_color=COLORS['light_gray'])
    
    # 可持续发展
    add_glow_rectangle(slide, Inches(8.9), Inches(1.8), Inches(4), Inches(2.5), 
                      COLORS['dark_blue'], COLORS['neon_purple'])
    add_text_box(slide, "可持续发展", Inches(9.1), Inches(2.0), Inches(3.6), Inches(0.4),
                font_size=18, font_color=COLORS['neon_purple'], bold=True)
    add_text_box(slide, "• 绿色计算与碳中和\n• AI助力可持续发展\n• 负责任的AI发展", 
                Inches(9.1), Inches(2.5), Inches(3.6), Inches(1.5),
                font_size=16, font_color=COLORS['light_gray'])
    
    # 愿景
    add_text_box(slide, '"构建更智能、更可持续的AI世界"', 
                Inches(0.5), Inches(4.8), Inches(12), Inches(0.8),
                font_size=28, font_color=COLORS['nvidia_green'], bold=True, alignment=PP_ALIGN.CENTER)
    
    # ==================== 第16页：核心亮点总结 ====================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, COLORS['dark_blue'])
    
    add_title_text(slide, "核心亮点总结", Inches(0.5), Inches(0.5), 
                  Inches(12), Inches(0.8), font_size=40, color=COLORS['nvidia_green'])
    add_accent_line(slide, Inches(0.5), Inches(1.3), Inches(2.5), Inches(0.03), COLORS['nvidia_green'])
    
    highlights = [
        ("硬件创新", "新一代GPU架构\nVera Rubin超级计算机", COLORS['neon_cyan']),
        ("软件生态", "CUDA平台持续演进\n500万+开发者", COLORS['nvidia_green']),
        ("应用落地", "AI在各行业深度渗透\n自动驾驶、医疗、工业", COLORS['neon_blue']),
        ("未来愿景", "构建智能、可持续\n的AI世界", COLORS['neon_purple']),
    ]
    
    for i, (title, desc, color) in enumerate(highlights):
        x = 0.5 + i * 3.2
        add_glow_rectangle(slide, Inches(x), Inches(1.8), Inches(3), Inches(3.5), 
                          COLORS['background'], color)
        add_text_box(slide, f"0{i+1}", Inches(x + 0.2), Inches(2.0), Inches(2.6), Inches(0.5),
                    font_size=24, font_color=color, bold=True)
        add_text_box(slide, title, Inches(x + 0.2), Inches(2.6), Inches(2.6), Inches(0.5),
                    font_size=18, font_color=color, bold=True)
        add_text_box(slide, desc, Inches(x + 0.2), Inches(3.2), Inches(2.6), Inches(1.5),
                    font_size=14, font_color=COLORS['light_gray'])
    
    # 底部总结
    add_text_box(slide, "NVIDIA正以完整的AI技术栈，引领人工智能时代的发展方向", 
                Inches(0.5), Inches(5.8), Inches(12), Inches(0.6),
                font_size=22, font_color=COLORS['white'], bold=True, alignment=PP_ALIGN.CENTER)
    
    # ==================== 第17页：总结 ====================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, COLORS['background'])
    
    add_title_text(slide, "总结", Inches(0.5), Inches(0.5), 
                  Inches(12), Inches(0.8), font_size=40, color=COLORS['nvidia_green'])
    add_accent_line(slide, Inches(0.5), Inches(1.3), Inches(2.5), Inches(0.03), COLORS['nvidia_green'])
    
    summary_items = [
        "NVIDIA GTC 2026展示了AI技术的无限可能",
        "从硬件创新到软件生态，构建完整技术栈",
        "从技术突破到应用落地，引领行业发展",
        "从芯片到云端，全栈AI解决方案",
        "从开发者到企业，赋能AI生态繁荣",
        "从技术创新到社会责任，推动可持续发展"
    ]
    add_bullet_points(slide, summary_items, Inches(0.5), Inches(1.8), Inches(12), Inches(4), 
                     font_size=22, color=COLORS['light_gray'])
    
    # 关键词
    keywords = ["AI", "GPU", "CUDA", "深度学习", "生成式AI", "边缘计算"]
    keyword_text = "  |  ".join(keywords)
    add_text_box(slide, keyword_text, Inches(0.5), Inches(5.5), Inches(12), Inches(0.5),
                font_size=18, font_color=COLORS['nvidia_green'], alignment=PP_ALIGN.CENTER)
    
    # ==================== 第18页：结尾页 ====================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, COLORS['background'])
    
    # 顶部装饰线
    add_accent_line(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.05), COLORS['nvidia_green'])
    
    # 感谢语
    add_text_box(slide, "感谢观看", Inches(0.5), Inches(2.5), Inches(12), Inches(1),
                font_size=56, font_color=COLORS['nvidia_green'], bold=True, alignment=PP_ALIGN.CENTER)
    
    # 英文
    add_text_box(slide, "THANK YOU", Inches(0.5), Inches(3.5), Inches(12), Inches(0.8),
                font_size=32, font_color=COLORS['neon_blue'], alignment=PP_ALIGN.CENTER)
    
    # 装饰线
    add_accent_line(slide, Inches(5.5), Inches(4.3), Inches(2.333), Inches(0.03), COLORS['nvidia_green'])
    
    # 信息
    add_text_box(slide, "NVIDIA GTC 2026 | GPU Technology Conference", 
                Inches(0.5), Inches(5.0), Inches(12), Inches(0.5),
                font_size=18, font_color=COLORS['gray'], alignment=PP_ALIGN.CENTER)
    
    add_text_box(slide, "视频来源：天下杂志 | 整理日期：2026年3月", 
                Inches(0.5), Inches(5.5), Inches(12), Inches(0.4),
                font_size=14, font_color=COLORS['gray'], alignment=PP_ALIGN.CENTER)
    
    # 底部装饰线
    add_accent_line(slide, Inches(0), Inches(7.45), Inches(13.333), Inches(0.05), COLORS['nvidia_green'])
    
    # 保存
    output_path = "/root/.openclaw/workspace/03_输出成果/NVIDIA_GTC_2026_科技风格.pptx"
    prs.save(output_path)
    print(f"PPT已保存至: {output_path}")
    return output_path

if __name__ == "__main__":
    create_presentation()

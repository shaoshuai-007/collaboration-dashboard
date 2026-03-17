#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""营业受理数字员工建设方案 - 严格按照模版第6页样式
方法论：内容素材 → 模板风格 → 图模板选择 → 统一调整

模版第6页样式：纵向左文右图排版
- 标题：字号24、加粗
- 建设背景和目标：一句话说明，字号16、加粗
- 左侧标题"重点举措"：字号14、加粗
- 左侧内容：小标题14号加粗，正文12号
- 右侧图示：系统流程图/架构图
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ==================== 模板风格定义 ====================
TELECOM_RED = RGBColor(201, 56, 50)
DARK_RED = RGBColor(172, 0, 0)
ORANGE = RGBColor(237, 125, 49)
BLUE = RGBColor(0, 110, 189)
LIGHT_BLUE = RGBColor(79, 129, 189)
DARK_BLUE = RGBColor(31, 73, 125)
GRAY = RGBColor(166, 166, 166)
DARK_GRAY = RGBColor(89, 89, 89)
LIGHT_GRAY = RGBColor(240, 240, 240)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)
GREEN = RGBColor(84, 130, 53)

FONT_NAME = "微软雅黑"

# ==================== 创建演示文稿 ====================
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ==================== 辅助函数 ====================

def add_text(slide, left, top, width, height, text, size=12, bold=False, color=BLACK, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = FONT_NAME
    p.alignment = align
    return tb

def add_box(slide, left, top, width, height, fill_color, text="", size=12, text_color=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 
                                    Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = text_color
        p.font.bold = True
        p.font.name = FONT_NAME
        p.alignment = PP_ALIGN.CENTER
    return shape

def add_rect(slide, left, top, width, height, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 
                               Inches(left), Inches(top), Inches(width), Inches(height))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def add_arrow(slide, left, top, width, height, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, 
                               Inches(left), Inches(top), Inches(width), Inches(height))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def add_header(slide, title):
    """添加页面标题（统一风格）"""
    add_rect(slide, 0.5, 0.4, 0.12, 0.5, TELECOM_RED)
    add_text(slide, 0.8, 0.4, 5, 0.5, title, 24, True, DARK_RED)

def add_target_line(slide, text):
    """添加建设背景和目标行（模版第6页样式）"""
    add_text(slide, 0.8, 1.1, 11.5, 0.4, text, 16, True, BLACK)

def add_measure_title(slide):
    """添加"重点举措"标题（模版第6页样式）"""
    add_text(slide, 0.8, 1.8, 4.5, 0.4, "重点举措", 14, True, DARK_RED)

def add_measure_content(slide, items):
    """添加左侧举措内容（模版第6页样式）
    items: [(小标题, 正文), ...]
    """
    y = 2.5
    for title, desc in items:
        # 小标题
        add_text(slide, 0.8, y, 4.5, 0.35, title, 14, True, BLACK)
        y += 0.35
        # 正文
        add_text(slide, 0.8, y, 4.5, 0.6, desc, 12, False, DARK_GRAY)
        y += 0.7

def add_diagram_area(slide, left=5.7, top=1.8, width=7.3, height=5.0):
    """添加右侧图示区域背景"""
    add_rect(slide, left, top, width, height, LIGHT_GRAY)

# ==================== P1: 封面 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.333, 7.5, WHITE)
add_rect(slide, 0, 0, 13.333, 0.12, TELECOM_RED)
add_text(slide, 1, 2.0, 11.333, 0.8, "营业受理数字员工建设方案", 36, True, DARK_RED, PP_ALIGN.CENTER)
add_text(slide, 1, 2.9, 11.333, 0.5, "集团考核达标 | 一线效率提升 | 客户体验改善", 18, False, DARK_GRAY, PP_ALIGN.CENTER)
add_rect(slide, 5, 3.5, 3.333, 0.02, TELECOM_RED)
add_text(slide, 1, 4.5, 11.333, 0.4, "中国电信湖北分公司    智能云网业务运营中心", 14, False, DARK_GRAY, PP_ALIGN.CENTER)
add_text(slide, 1, 5.0, 11.333, 0.4, "2026年3月", 14, False, GRAY, PP_ALIGN.CENTER)

# ==================== P2: 目录 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.333, 7.5, WHITE)
add_rect(slide, 0, 0, 13.333, 0.12, TELECOM_RED)
add_text(slide, 0.5, 0.5, 2, 0.6, "目 录", 28, True, DARK_RED)
add_text(slide, 0.5, 1.0, 2, 0.4, "CONTENTS", 14, False, GRAY)

items = [
    ("01", "背景与目标", "项目背景、问题分析、建设目标"),
    ("02", "方案设计", "整体架构、四大能力、技术架构"),
    ("03", "实施规划", "系统集成、实施路径、关键指标"),
    ("04", "价值与优势", "核心价值、创新亮点、成功案例")
]
y = 1.8
for num, title, desc in items:
    add_box(slide, 0.8, y, 0.6, 0.6, TELECOM_RED, num, 16, WHITE)
    add_text(slide, 1.6, y, 4, 0.4, title, 18, True, DARK_RED)
    add_text(slide, 1.6, y + 0.4, 4, 0.4, desc, 12, False, GRAY)
    y += 1.2

# ==================== P3: 项目背景（模版第6页样式）====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.333, 7.5, WHITE)
add_rect(slide, 0, 0, 13.333, 0.12, TELECOM_RED)

add_header(slide, "项目背景")
add_target_line(slide, "满足集团考核要求，解决一线营业员效率低、体验差的问题。")
add_measure_title(slide)

# 左侧举措内容
items = [
    ("调研分析", "调研12个分公司40+场景需求，Top3需求为智能问答、智能查询、智能诊断。"),
    ("考核对标", "云运〔2025〕4号要求AI嵌入10个关键场景，自动化率达70%，占指标10%。"),
    ("问题定位", "一线营业员业务办理依赖人工操作，效率低易出错，客户等待时间长。")
]
add_measure_content(slide, items)

# 右侧图示
add_diagram_area(slide)
# 考核指标展示
add_text(slide, 7.0, 2.0, 5, 0.4, "集团考核指标", 14, True, DARK_RED)
add_box(slide, 7.0, 2.6, 2.0, 1.8, TELECOM_RED, "10%", 24, WHITE)
add_text(slide, 7.0, 4.5, 2.0, 0.3, "AI场景指标", 10, False, DARK_GRAY, PP_ALIGN.CENTER)
add_box(slide, 9.5, 2.6, 2.0, 1.8, ORANGE, "70%", 24, WHITE)
add_text(slide, 9.5, 4.5, 2.0, 0.3, "自动化率", 10, False, DARK_GRAY, PP_ALIGN.CENTER)
add_text(slide, 7.0, 5.0, 5, 0.4, "调研结论", 14, True, DARK_RED)
add_text(slide, 7.0, 5.5, 5, 0.8, "必须建设数字员工，满足考核要求，解决一线痛点。", 12, False, DARK_GRAY)

# ==================== P4: 问题分析（模版第6页样式）====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.333, 7.5, WHITE)
add_rect(slide, 0, 0, 13.333, 0.12, TELECOM_RED)

add_header(slide, "问题分析")
add_target_line(slide, "解决核心根因问题，建设数字员工赋能一线营业员。")
add_measure_title(slide)

items = [
    ("问题表象", "效率低：业务办理依赖人工，耗时长易出错。体验差：客户等待时间长。"),
    ("根因分析", "核心根因：缺乏智能化工具赋能，一线营业员无AI助手辅助。"),
    ("解决方案", "建设营业受理数字员工，四大核心能力赋能一线，提升效率改善体验。")
]
add_measure_content(slide, items)

add_diagram_area(slide)
# 问题树
add_text(slide, 7.0, 2.0, 5, 0.4, "问题诊断", 14, True, DARK_RED)
add_box(slide, 7.2, 2.6, 2.0, 0.7, RGBColor(200, 50, 50), "效率低", 13, WHITE)
add_box(slide, 9.6, 2.6, 2.0, 0.7, RGBColor(200, 100, 50), "体验差", 13, WHITE)
add_box(slide, 7.2, 3.5, 2.0, 0.7, RGBColor(200, 150, 50), "知识获取难", 12, WHITE)
add_box(slide, 7.5, 4.5, 3.8, 0.7, ORANGE, "核心根因：缺乏工具", 12, WHITE)
add_box(slide, 7.5, 5.4, 3.8, 0.7, GREEN, "解决方案：数字员工", 12, WHITE)

# ==================== P5: 建设目标（模版第6页样式）====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.333, 7.5, WHITE)
add_rect(slide, 0, 0, 13.333, 0.12, TELECOM_RED)

add_header(slide, "建设目标")
add_target_line(slide, "办理时长缩短40%，构建四大核心能力智能体集群。")
add_measure_title(slide)

items = [
    ("智能问答 Q-Agent", "准确率≥90%，基于大模型+RAG技术，快速解答业务知识。"),
    ("智能查询 S-Agent", "效率提升50%，实体抽取+API集成，客户信息快速获取。"),
    ("智能办理 H-Agent", "时长缩短40%，多轮对话+智能填单，业务办理自动化。"),
    ("智能诊断 D-Agent", "时间缩短60%，规则引擎+云哨对接，问题快速定位。")
]
add_measure_content(slide, items)

add_diagram_area(slide)
# 指标展示
add_text(slide, 7.0, 2.0, 5, 0.4, "关键指标", 14, True, DARK_RED)
metrics = [("40%", "办理时长"), ("50%", "查询效率"), ("60%", "诊断时间"), ("90%", "问答准确率")]
x = 7.0
for value, label in metrics:
    add_box(slide, x, 2.5, 1.3, 1.3, TELECOM_RED, value, 18, WHITE)
    add_text(slide, x, 3.9, 1.3, 0.3, label, 9, False, DARK_GRAY, PP_ALIGN.CENTER)
    x += 1.5
add_text(slide, 7.0, 4.5, 5, 0.4, "预期效果", 14, True, DARK_RED)
add_text(slide, 7.0, 5.0, 5, 1.0, "满足集团考核要求，一线效率提升40%以上，客户满意度提升15%。", 12, False, DARK_GRAY)

# ==================== P6: 整体架构（模版第6页样式）====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.333, 7.5, WHITE)
add_rect(slide, 0, 0, 13.333, 0.12, TELECOM_RED)

add_header(slide, "整体架构")
add_target_line(slide, "架构清晰可落地，六层分层设计+标准化接口。")
add_measure_title(slide)

items = [
    ("架构理念", "分层设计，每层职责明确，便于开发、测试、运维，支持横向扩展。"),
    ("交互层", "ChatUI统一入口，PC端、移动端、随翼选、语音入口多终端支持。"),
    ("智能体层", "四大智能体分工协作：Q-Agent问答、S-Agent查询、H-Agent办理、D-Agent诊断。"),
    ("能力层", "AI能力组件复用：意图识别、实体抽取、知识检索、智能填单、多轮对话。")
]
add_measure_content(slide, items)

add_diagram_area(slide)
# 六层架构图
add_text(slide, 7.0, 2.0, 5, 0.4, "六层架构", 14, True, DARK_RED)
layers = [
    ("交互层 (ChatUI)", RGBColor(100, 149, 237)),
    ("智能体集群", TELECOM_RED),
    ("AI能力层", ORANGE),
    ("模型层", BLUE),
    ("数据层", LIGHT_BLUE),
    ("系统集成层", DARK_GRAY)
]
y = 2.5
for name, color in layers:
    add_box(slide, 7.0, y, 5.0, 0.65, color, name, 11, WHITE)
    y += 0.75

# ==================== P7: 智能问答能力（模版第6页样式）====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.333, 7.5, WHITE)
add_rect(slide, 0, 0, 13.333, 0.12, TELECOM_RED)

add_header(slide, "智能问答能力 (Q-Agent)")
add_target_line(slide, "准确率≥90%，大模型+RAG+向量检索。")
add_measure_title(slide)

items = [
    ("业务场景", "业务规则查询、资费政策解答、流程指引、故障排查、政策解读。"),
    ("技术实现", "用户输入→意图识别→知识检索→重排序→大模型生成→答案输出。"),
    ("知识库设计", "业务规则、资费政策、流程指引、产品介绍、故障排查、政策解读。")
]
add_measure_content(slide, items)

add_diagram_area(slide)
# 流程图
add_text(slide, 7.0, 2.0, 5, 0.4, "技术流程", 14, True, DARK_RED)
flow = ["用户输入", "意图识别", "知识检索", "重排序", "模型生成", "答案输出"]
colors = [TELECOM_RED, ORANGE, BLUE, LIGHT_BLUE, GREEN, DARK_GRAY]
x = 7.0
for i, (name, color) in enumerate(zip(flow, colors)):
    add_box(slide, x, 2.5, 1.3, 0.7, color, name, 10, WHITE)
    if i < len(flow) - 1:
        add_arrow(slide, x + 1.35, 2.65, 0.2, 0.25, GRAY)
    x += 1.5
add_text(slide, 7.0, 3.5, 5, 0.4, "应用场景", 14, True, DARK_RED)
add_text(slide, 7.0, 4.0, 5, 2.0, "• 业务规则查询：办理条件、政策解读\n• 资费政策解答：套餐详情、费用说明\n• 流程指引：操作步骤、注意事项\n• 故障排查：问题定位、处理建议", 11, False, DARK_GRAY)

# ==================== P8: 智能查询能力（模版第6页样式）====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.333, 7.5, WHITE)
add_rect(slide, 0, 0, 13.333, 0.12, TELECOM_RED)

add_header(slide, "智能查询能力 (S-Agent)")
add_target_line(slide, "效率提升50%，实体抽取+API集成。")
add_measure_title(slide)

items = [
    ("业务场景", "客户资料查询、订单状态查询、资费信息查询、套餐余量查询。"),
    ("技术实现", "用户输入→实体抽取→权限验证→API调用→结果整合→格式输出。"),
    ("实体类型", "手机号、客户号、订单号、身份证、套餐名称、业务类型。")
]
add_measure_content(slide, items)

add_diagram_area(slide)
add_text(slide, 7.0, 2.0, 5, 0.4, "技术流程", 14, True, DARK_RED)
flow = ["用户输入", "实体抽取", "权限验证", "API调用", "结果整合", "格式输出"]
x = 7.0
for i, (name, color) in enumerate(zip(flow, colors)):
    add_box(slide, x, 2.5, 1.3, 0.7, color, name, 10, WHITE)
    if i < len(flow) - 1:
        add_arrow(slide, x + 1.35, 2.65, 0.2, 0.25, GRAY)
    x += 1.5
add_text(slide, 7.0, 3.5, 5, 0.4, "应用场景", 14, True, DARK_RED)
add_text(slide, 7.0, 4.0, 5, 2.0, "• 客户资料查询：姓名、套餐、余额\n• 订单状态查询：进度、时间、结果\n• 资费信息查询：月租、流量、通话\n• 套餐余量查询：流量、短信、语音", 11, False, DARK_GRAY)

# ==================== P9: 智能办理能力（模版第6页样式）====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.333, 7.5, WHITE)
add_rect(slide, 0, 0, 13.333, 0.12, TELECOM_RED)

add_header(slide, "智能办理能力 (H-Agent)")
add_target_line(slide, "时长缩短40%，多轮对话+智能填单。")
add_measure_title(slide)

items = [
    ("业务场景", "换卡业务、移机业务、停复机、可选包办理、套餐变更。"),
    ("技术实现", "用户输入→意图识别→要素收集→规则校验→页面唤起→智能填单。"),
    ("核心能力", "多轮对话引导收集要素，智能填单自动填充表单，随翼选页面集成。")
]
add_measure_content(slide, items)

add_diagram_area(slide)
add_text(slide, 7.0, 2.0, 5, 0.4, "技术流程", 14, True, DARK_RED)
flow = ["用户输入", "意图识别", "要素收集", "规则校验", "页面唤起", "智能填单"]
x = 7.0
for i, (name, color) in enumerate(zip(flow, colors)):
    add_box(slide, x, 2.5, 1.3, 0.7, color, name, 10, WHITE)
    if i < len(flow) - 1:
        add_arrow(slide, x + 1.35, 2.65, 0.2, 0.25, GRAY)
    x += 1.5
add_text(slide, 7.0, 3.5, 5, 0.4, "应用场景", 14, True, DARK_RED)
add_text(slide, 7.0, 4.0, 5, 2.0, "• 换卡业务：SIM卡更换、补卡\n• 移机业务：地址迁移、设备更换\n• 停复机：停机、复机申请\n• 可选包办理：流量包、语音包订购", 11, False, DARK_GRAY)

# ==================== P10: 智能诊断能力（模版第6页样式）====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.333, 7.5, WHITE)
add_rect(slide, 0, 0, 13.333, 0.12, TELECOM_RED)

add_header(slide, "智能诊断能力 (D-Agent)")
add_target_line(slide, "时间缩短60%，规则引擎+云哨对接。")
add_measure_title(slide)

items = [
    ("业务场景", "订单卡单、办理失败、资源异常、计费异常、系统故障。"),
    ("技术实现", "问题输入→信息收集→诊断分析→结果生成→处理建议→一键报障。"),
    ("核心能力", "智能识别问题类型，规则引擎+知识库双重诊断，对接云哨平台。")
]
add_measure_content(slide, items)

add_diagram_area(slide)
add_text(slide, 7.0, 2.0, 5, 0.4, "技术流程", 14, True, DARK_RED)
flow = ["问题输入", "信息收集", "诊断分析", "结果生成", "处理建议", "一键报障"]
x = 7.0
for i, (name, color) in enumerate(zip(flow, colors)):
    add_box(slide, x, 2.5, 1.3, 0.7, color, name, 10, WHITE)
    if i < len(flow) - 1:
        add_arrow(slide, x + 1.35, 2.65, 0.2, 0.25, GRAY)
    x += 1.5
add_text(slide, 7.0, 3.5, 5, 0.4, "应用场景", 14, True, DARK_RED)
add_text(slide, 7.0, 4.0, 5, 2.0, "• 订单卡单：订单状态异常诊断\n• 办理失败：失败原因定位分析\n• 资源异常：资源分配问题诊断\n• 计费异常：费用差异问题定位", 11, False, DARK_GRAY)

# ==================== P11: 系统集成（模版第6页样式）====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.333, 7.5, WHITE)
add_rect(slide, 0, 0, 13.333, 0.12, TELECOM_RED)

add_header(slide, "系统集成")
add_target_line(slide, "八大系统对接，API网关统一集成+标准化接口。")
add_measure_title(slide)

items = [
    ("集成架构", "API网关统一入口，标准化接口规范，分步对接降低风险，统一认证鉴权。"),
    ("核心系统", "CRM、随翼选、客户中心、订单中心、计费中心、资源中心、云哨。")
]
add_measure_content(slide, items)

add_diagram_area(slide)
# 星型拓扑
add_text(slide, 7.0, 2.0, 5, 0.4, "系统集成架构", 14, True, DARK_RED)
add_box(slide, 8.0, 3.2, 2.5, 1.0, DARK_RED, "数字员工\nAPI网关", 12, WHITE)
systems = [
    (7.0, 2.5, "CRM", TELECOM_RED),
    (10.5, 2.5, "随翼选", ORANGE),
    (7.0, 4.8, "客户中心", BLUE),
    (10.5, 4.8, "订单中心", LIGHT_BLUE),
    (7.0, 5.6, "计费中心", GREEN),
    (10.5, 5.6, "云哨", DARK_GRAY)
]
for x, y, name, color in systems:
    add_box(slide, x, y, 1.8, 0.6, color, name, 10, WHITE)

# ==================== P12: 实施路径（模版第6页样式）====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.333, 7.5, WHITE)
add_rect(slide, 0, 0, 13.333, 0.12, TELECOM_RED)

add_header(slide, "实施路径")
add_target_line(slide, "5个月完成建设，四阶段渐进式实施。")
add_measure_title(slide)

items = [
    ("第一阶段", "2个月，智能问答+查询，基础能力建设，准确率≥90%。"),
    ("第二阶段", "2个月，智能办理，核心价值实现，时长缩短≥30%。"),
    ("第三阶段", "1个月，智能诊断，完善闭环，准确率≥80%。"),
    ("第四阶段", "持续优化，越用越聪明，场景扩展。")
]
add_measure_content(slide, items)

add_diagram_area(slide)
# 时间轴
add_text(slide, 7.0, 2.0, 5, 0.4, "实施时间轴", 14, True, DARK_RED)
add_rect(slide, 7.2, 3.2, 4.5, 0.03, GRAY)
phases = [
    ("第一阶段", "2个月", TELECOM_RED, 7.2),
    ("第二阶段", "2个月", ORANGE, 8.4),
    ("第三阶段", "1个月", BLUE, 9.6),
    ("第四阶段", "持续", DARK_GRAY, 10.8)
]
for title, time, color, x in phases:
    add_box(slide, x, 3.0, 0.5, 0.5, color, "", 12, WHITE)
    add_text(slide, x - 0.2, 2.5, 0.9, 0.3, time, 9, False, DARK_GRAY, PP_ALIGN.CENTER)
add_text(slide, 7.0, 3.8, 5, 2.0, "• 第一阶段：问答+查询基础能力\n• 第二阶段：办理核心价值\n• 第三阶段：诊断完善闭环\n• 第四阶段：持续优化迭代", 11, False, DARK_GRAY)

# ==================== P13: 关键指标（模版第6页样式）====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.333, 7.5, WHITE)
add_rect(slide, 0, 0, 13.333, 0.12, TELECOM_RED)

add_header(slide, "关键指标")
add_target_line(slide, "效果可量化，业务+技术双维度指标体系。")
add_measure_title(slide)

items = [
    ("效率指标", "办理时长缩短40%，查询效率提升50%，诊断时间缩短60%。"),
    ("质量指标", "问答准确率≥90%，填单准确率≥98%，诊断准确率≥85%。"),
    ("体验指标", "用户满意度≥85%，响应时间<3秒。")
]
add_measure_content(slide, items)

add_diagram_area(slide)
add_text(slide, 7.0, 2.0, 5, 0.4, "关键指标", 14, True, DARK_RED)
# 效率指标
add_text(slide, 7.0, 2.5, 5, 0.3, "效率指标", 12, True, DARK_RED)
metrics1 = [("40%", "时长缩短"), ("50%", "效率提升"), ("60%", "时间缩短")]
x = 7.0
for value, label in metrics1:
    add_box(slide, x, 2.9, 1.3, 1.1, TELECOM_RED, value, 16, WHITE)
    add_text(slide, x, 4.1, 1.3, 0.25, label, 8, False, DARK_GRAY, PP_ALIGN.CENTER)
    x += 1.5
# 质量指标
add_text(slide, 7.0, 4.5, 5, 0.3, "质量指标", 12, True, ORANGE)
metrics2 = [("≥90%", "问答准确"), ("≥98%", "填单准确"), ("≥85%", "诊断准确")]
x = 7.0
for value, label in metrics2:
    add_box(slide, x, 4.9, 1.3, 1.1, ORANGE, value, 14, WHITE)
    add_text(slide, x, 6.1, 1.3, 0.25, label, 8, False, DARK_GRAY, PP_ALIGN.CENTER)
    x += 1.5

# ==================== P14: 核心价值（模版第6页样式）====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.333, 7.5, WHITE)
add_rect(slide, 0, 0, 13.333, 0.12, TELECOM_RED)

add_header(slide, "核心价值")
add_target_line(slide, "四大价值维度，考核+效率+体验+人力。")
add_measure_title(slide)

items = [
    ("考核达标", "满足集团考核要求，占指标10%，31省标杆项目。"),
    ("效率提升", "办理时长缩短40%，查询效率提升50%，诊断时间缩短60%。"),
    ("体验改善", "客户满意度提升，一线满意度提升，服务响应加快。"),
    ("人力释放", "推动一线员工转型，从事高价值任务，减负增效。")
]
add_measure_content(slide, items)

add_diagram_area(slide)
add_text(slide, 7.0, 2.0, 5, 0.4, "核心价值", 14, True, DARK_RED)
values = [
    ("考核达标", TELECOM_RED),
    ("效率提升", ORANGE),
    ("体验改善", BLUE),
    ("人力释放", DARK_GRAY)
]
y = 2.5
for title, color in values:
    add_box(slide, 7.0, y, 5.0, 0.7, color, title, 14, WHITE)
    y += 0.85

# ==================== P15: 创新亮点（模版第6页样式）====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.333, 7.5, WHITE)
add_rect(slide, 0, 0, 13.333, 0.12, TELECOM_RED)

add_header(slide, "创新亮点")
add_target_line(slide, "引领行业AI应用，五大创新技术融合。")
add_measure_title(slide)

items = [
    ("智能体集群", "四大智能体分工协作，专业精准，智能问答、查询、办理、诊断。"),
    ("多模态交互", "语音+文字+图片多种输入方式，灵活便捷，降低使用门槛。"),
    ("主动推荐", "客户画像分析，精准推荐，千人千面，提升营销效率。"),
    ("学习型系统", "持续学习，越用越聪明，自我进化，适应业务变化。")
]
add_measure_content(slide, items)

add_diagram_area(slide)
add_text(slide, 7.0, 2.0, 5, 0.4, "创新亮点", 14, True, DARK_RED)
innovations = [
    ("智能体集群", TELECOM_RED),
    ("多模态交互", ORANGE),
    ("主动推荐", BLUE),
    ("学习型系统", LIGHT_BLUE),
    ("场景化卡片", DARK_GRAY)
]
y = 2.5
for title, color in innovations:
    add_box(slide, 7.0, y, 5.0, 0.65, color, title, 12, WHITE)
    y += 0.75

# ==================== P16: 成功案例（模版第6页样式）====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.333, 7.5, WHITE)
add_rect(slide, 0, 0, 13.333, 0.12, TELECOM_RED)

add_header(slide, "成功案例")
add_target_line(slide, "三地成功实践，复制成熟经验，降低风险。")
add_measure_title(slide)

items = [
    ("陕西电信", "营业办理数字员工，办理时长缩短40%，一线效率显著提升。"),
    ("重庆电信", "营业受理数字员工，智能问答+查询，满足考核要求。"),
    ("安徽电信", "CRM域AI应用，四大能力全覆盖，31省标杆。")
]
add_measure_content(slide, items)

add_diagram_area(slide)
add_text(slide, 7.0, 2.0, 5, 0.4, "成功案例", 14, True, DARK_RED)
cases = [
    ("陕西电信", "时长缩短40%", TELECOM_RED),
    ("重庆电信", "满足考核要求", ORANGE),
    ("安徽电信", "31省标杆", BLUE)
]
y = 2.5
for title, desc, color in cases:
    add_box(slide, 7.0, y, 5.0, 0.9, color, "", 12, WHITE)
    add_text(slide, 7.2, y + 0.15, 2.5, 0.3, title, 13, True, WHITE)
    add_text(slide, 7.2, y + 0.5, 4.5, 0.3, desc, 11, False, WHITE)
    y += 1.0

# ==================== P17: 总结与展望（模版第6页样式）====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.333, 7.5, WHITE)
add_rect(slide, 0, 0, 13.333, 0.12, TELECOM_RED)

add_header(slide, "总结与展望")
add_target_line(slide, "方案成熟可落地，分阶段实施，风险可控。")
add_measure_title(slide)

items = [
    ("方案成熟", "参考三地成功案例，技术方案验证可行，实施方法论完善。"),
    ("技术可行", "大模型+RAG技术成熟，开源组件丰富，技术风险低。"),
    ("价值明确", "效率提升40%以上，满足考核要求，客户体验改善。"),
    ("后续规划", "持续优化迭代，扩展更多场景，复制推广到其他省份。")
]
add_measure_content(slide, items)

add_diagram_area(slide)
add_text(slide, 7.0, 2.0, 5, 0.4, "总结", 14, True, DARK_RED)
summaries = [
    ("方案成熟", TELECOM_RED),
    ("技术可行", ORANGE),
    ("价值明确", BLUE),
    ("风险可控", DARK_GRAY)
]
y = 2.5
for title, color in summaries:
    add_box(slide, 7.0, y, 5.0, 0.7, color, title, 13, WHITE)
    y += 0.85

# ==================== P18: 结束页 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.333, 7.5, WHITE)
add_rect(slide, 0, 0, 13.333, 0.12, TELECOM_RED)
add_text(slide, 1, 2.5, 11.333, 1.0, "谢谢！", 48, True, DARK_RED, PP_ALIGN.CENTER)
add_text(slide, 1, 3.5, 11.333, 0.6, "让我们一起，用AI赋能一线，提升效率，改善体验！", 18, False, DARK_GRAY, PP_ALIGN.CENTER)
add_text(slide, 1, 4.5, 11.333, 0.4, "南有乔木，不可休思", 14, False, GRAY, PP_ALIGN.CENTER)

# 保存
output_path = "/root/.openclaw/workspace/03_项目管理/营业受理数字员工/03_工程输出/指南针工程/营业受理数字员工建设方案.pptx"
prs.save(output_path)
print(f"PPT已生成: {output_path}")
print(f"总页数: {len(prs.slides)}")
print("\n严格按照模版第6页样式：")
print("- 标题：字号24、加粗")
print("- 建设背景和目标：字号16、加粗")
print("- 左侧标题'重点举措'：字号14、加粗")
print("- 左侧内容：小标题14号加粗，正文12号")
print("- 右侧图示：系统流程图/架构图")

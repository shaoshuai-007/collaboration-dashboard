#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""营业受理数字员工建设方案 - 精进版

修改点：
1. 目录页：目录内容全部居中对齐
2. 问题分析页：
   - 解决方案设计细化，不少于150字，分段排列
   - 右图：效率低、体验差、知识获取难一行排版
   - 根因扩充（增加3项内容）
   - 解决方案扩展（增加3项内容）
3. 整体架构页：架构理念阐述六层细节，每层不少于50字
4. 7-18页：改变图形展现方式，避免单调
5. 整体PPT：保持整洁，不得越界，注意排版
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ==================== 配色方案（60-30-10）====================
PRIMARY = RGBColor(201, 56, 50)
PRIMARY_DARK = RGBColor(172, 0, 0)
SECONDARY = RGBColor(0, 110, 189)
NEUTRAL_DARK = RGBColor(89, 89, 89)
NEUTRAL = RGBColor(166, 166, 166)
NEUTRAL_LIGHT = RGBColor(240, 240, 240)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)

FONT_NAME = "微软雅黑"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ==================== 辅助函数 ====================

def add_text(slide, left, top, width, height, text, size=12, bold=False, color=BLACK, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = FONT_NAME
    p.alignment = align
    return tb

def add_multiline_text(slide, left, top, width, height, lines, size=12, color=BLACK, align=PP_ALIGN.LEFT):
    """添加多行文本，每行一个段落"""
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = FONT_NAME
        p.alignment = align
        p.space_after = Pt(6)
    
    return tb

def add_rounded_box(slide, left, top, width, height, fill_color, text="", size=12, text_color=WHITE, bold=True):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 
                                    Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = text_color
        p.font.bold = bold
        p.font.name = FONT_NAME
        p.alignment = PP_ALIGN.CENTER
        shape.text_frame_anchor = MSO_ANCHOR.MIDDLE
    
    return shape

def add_rect(slide, left, top, width, height, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 
                               Inches(left), Inches(top), Inches(width), Inches(height))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def add_arrow(slide, left, top, width, height, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, 
                               Inches(left), Inches(top), Inches(width), Inches(height))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def add_down_arrow(slide, left, top, width, height, color):
    s = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, 
                               Inches(left), Inches(top), Inches(width), Inches(height))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def add_header(slide, title, subtitle=""):
    """添加页面标题（格式：标题：副标题）"""
    add_rect(slide, 0.5, 0.4, 0.12, 0.5, PRIMARY)
    
    if subtitle:
        full_title = f"{title}：{subtitle}"
        add_text(slide, 0.8, 0.4, 10, 0.5, full_title, 24, True, PRIMARY_DARK)
    else:
        add_text(slide, 0.8, 0.4, 5, 0.5, title, 24, True, PRIMARY_DARK)

def add_target_line(slide, text):
    add_text(slide, 0.8, 1.1, 11.5, 0.4, text, 16, True, BLACK)

def add_measure_title(slide):
    add_text(slide, 0.8, 1.8, 4.5, 0.4, "重点举措", 14, True, PRIMARY_DARK)

def add_measure_content(slide, items):
    y = 2.5
    for title, desc in items:
        add_text(slide, 0.8, y, 4.5, 0.35, title, 14, True, BLACK)
        y += 0.35
        add_text(slide, 0.8, y, 4.5, 0.65, desc, 12, False, NEUTRAL_DARK)
        y += 0.75

def add_metric_dashboard(slide, left, top, width, height, metrics):
    n = len(metrics)
    box_width = width / n - 0.15
    
    for i, (value, label) in enumerate(metrics):
        x = left + i * (box_width + 0.15)
        
        if i == 0:
            color = PRIMARY
        elif i < n - 1:
            color = SECONDARY
        else:
            color = NEUTRAL_DARK
        
        add_rounded_box(slide, x, top, box_width, height * 0.65, color, "", 12, WHITE)
        add_text(slide, x, top + 0.1, box_width, height * 0.4, value, 26, True, WHITE, PP_ALIGN.CENTER)
        add_text(slide, x, top + height * 0.45, box_width, height * 0.2, label, 10, False, WHITE, PP_ALIGN.CENTER)

def add_process_flow(slide, left, top, width, height, steps):
    n = len(steps)
    box_width = (width - (n - 1) * 0.25) / n
    arrow_width = 0.2
    
    for i, step in enumerate(steps):
        x = left + i * (box_width + arrow_width + 0.05)
        
        if i < n * 0.6:
            color = PRIMARY
        elif i < n * 0.9:
            color = SECONDARY
        else:
            color = NEUTRAL_DARK
        
        add_rounded_box(slide, x, top, box_width, height, color, step, 10, WHITE)
        
        if i < n - 1:
            add_arrow(slide, x + box_width + 0.02, top + height/2 - 0.08, arrow_width, 0.16, NEUTRAL)

def add_value_chain(slide, left, top, width, height, values):
    n = len(values)
    box_width = width / n - 0.12
    
    for i, name in enumerate(values):
        x = left + i * (box_width + 0.12)
        
        if i < n * 0.5:
            color = PRIMARY
        elif i < n * 0.8:
            color = SECONDARY
        else:
            color = NEUTRAL_DARK
        
        add_rounded_box(slide, x, top, box_width, height, color, name, 12, WHITE)

# ==================== P1: 封面 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.333, 7.5, WHITE)
add_rect(slide, 0, 0, 13.333, 0.12, PRIMARY)
add_text(slide, 1, 2.2, 11.333, 0.8, "营业受理数字员工建设方案", 36, True, PRIMARY_DARK, PP_ALIGN.CENTER)
add_text(slide, 1, 3.1, 11.333, 0.5, "集团考核达标 | 一线效率提升 | 客户体验改善", 18, False, NEUTRAL_DARK, PP_ALIGN.CENTER)
add_rect(slide, 5, 3.7, 3.333, 0.02, PRIMARY)
add_text(slide, 1, 4.5, 11.333, 0.4, "2026年3月", 14, False, NEUTRAL, PP_ALIGN.CENTER)

# ==================== P2: 目录（修改：居中对齐）====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.333, 7.5, WHITE)
add_rect(slide, 0, 0, 13.333, 0.12, PRIMARY)
add_header(slide, "目 录", "四大篇章清晰呈现")

items = [
    ("01", "背景与目标", "项目背景、问题分析、建设目标"),
    ("02", "方案设计", "整体架构、四大能力、技术架构"),
    ("03", "实施规划", "系统集成、实施路径、关键指标"),
    ("04", "价值与优势", "核心价值、创新亮点、成功案例")
]

# 居中对齐
center_x = 6.667
y = 2.0
for i, (num, title, desc) in enumerate(items):
    color = PRIMARY if i < 2 else SECONDARY
    
    # 编号
    add_rounded_box(slide, center_x - 3, y, 0.6, 0.6, color, num, 16, WHITE)
    # 标题
    add_text(slide, center_x - 2.2, y, 4, 0.35, title, 18, True, PRIMARY_DARK, PP_ALIGN.LEFT)
    # 描述
    add_text(slide, center_x - 2.2, y + 0.35, 4, 0.3, desc, 12, False, NEUTRAL, PP_ALIGN.LEFT)
    y += 1.1

# ==================== P3: 项目背景 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.333, 7.5, WHITE)
add_rect(slide, 0, 0, 13.333, 0.12, PRIMARY)

add_header(slide, "项目背景", "满足考核解决痛点")
add_target_line(slide, "满足集团云运〔2025〕4号考核要求，解决一线营业员效率低、体验差的核心痛点。")
add_measure_title(slide)

items = [
    ("深度调研分析", "深入12个分公司开展需求调研，覆盖营业厅、客服、渠道等40+业务场景，识别Top3核心需求：智能问答、智能查询、智能诊断。"),
    ("考核指标对标", "云运〔2025〕4号文件明确要求AI嵌入10个关键业务场景，自动化率达到70%以上，占考核指标10分，时间紧任务重。"),
    ("痛点问题定位", "一线营业员业务办理高度依赖人工操作，流程复杂、知识分散、易出错，客户平均等待时间超过15分钟，满意度持续下降。")
]
add_measure_content(slide, items)

add_text(slide, 6.0, 1.9, 6.5, 0.4, "集团考核指标", 14, True, PRIMARY_DARK)
metrics = [("10%", "AI场景指标"), ("70%", "自动化率"), ("40+", "业务场景")]
add_metric_dashboard(slide, 6.2, 2.4, 6.0, 1.8, metrics)

add_text(slide, 6.0, 4.5, 6.5, 0.4, "调研结论", 14, True, PRIMARY_DARK)
add_text(slide, 6.0, 5.0, 6.5, 1.0, "必须建设数字员工，满足集团考核要求，解决一线营业员效率低、体验差的痛点问题。", 12, False, NEUTRAL_DARK)

# ==================== P4: 问题分析（大改）====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.333, 7.5, WHITE)
add_rect(slide, 0, 0, 13.333, 0.12, PRIMARY)

add_header(slide, "问题分析", "定位根因精准施策")
add_target_line(slide, "深入分析问题表象，精准定位核心根因，制定针对性解决方案。")
add_measure_title(slide)

# 左侧举措
items = [
    ("问题表象分析", "效率低：业务办理平均耗时20分钟，高度依赖人工操作，流程繁琐易出错。体验差：客户平均等待15分钟以上，排队时间长，满意度持续下降。"),
    ("核心根因定位", "经过深度分析，核心根因是缺乏智能化工具赋能一线营业员，业务知识分散、查询路径复杂、办理流程冗长、问题定位困难。")
]

y = 2.5
for title, desc in items:
    add_text(slide, 0.8, y, 4.5, 0.35, title, 14, True, BLACK)
    y += 0.35
    add_text(slide, 0.8, y, 4.5, 0.65, desc, 12, False, NEUTRAL_DARK)
    y += 0.75

# 解决方案设计（细化，不少于150字，分段排列）
add_text(slide, 0.8, y, 4.5, 0.35, "解决方案设计", 14, True, BLACK)
solution_lines = [
    "基于大模型+RAG技术，建设营业受理数字员工系统。",
    "通过智能问答能力，实现业务规则、资费政策等知识快速问答，准确率≥90%，减少人工查询时间70%。",
    "通过智能查询能力，实现客户资料、订单状态等信息快速获取，效率提升50%。",
    "通过智能办理能力，实现换卡、移机等业务自动办理，时长缩短40%。",
    "通过智能诊断能力，实现订单卡单、办理失败等问题智能定位，时间缩短60%。",
    "形成完整业务闭环，全方位赋能一线营业员，提升效率改善体验。"
]
add_multiline_text(slide, 0.8, y + 0.35, 4.5, 2.0, solution_lines, 11, NEUTRAL_DARK)

# 右侧图：效率低、体验差、知识获取难一行排版
add_text(slide, 6.0, 1.9, 6.5, 0.4, "问题诊断", 14, True, PRIMARY_DARK)

# 三个问题一行排列
box_width = 2.0
gap = 0.2
start_x = 6.2
add_rounded_box(slide, start_x, 2.4, box_width, 0.7, PRIMARY, "效率低", 12, WHITE)
add_rounded_box(slide, start_x + box_width + gap, 2.4, box_width, 0.7, SECONDARY, "体验差", 12, WHITE)
add_rounded_box(slide, start_x + 2 * (box_width + gap), 2.4, box_width, 0.7, NEUTRAL_DARK, "知识获取难", 11, WHITE)

# 下箭头
add_down_arrow(slide, 8.5, 3.2, 0.25, 0.4, NEUTRAL)

# 根因（扩充，增加3项内容）
add_text(slide, 6.0, 3.7, 6.5, 0.35, "核心根因", 14, True, PRIMARY_DARK)
root_causes = [
    "• 缺乏智能化工具：一线营业员无AI助手辅助",
    "• 业务知识分散：规则、流程、政策散落各处",
    "• 查询路径复杂：多系统切换，操作繁琐",
    "• 办理流程冗长：人工操作多，易出错",
    "• 问题定位困难：缺乏智能诊断能力",
    "• 数据孤岛严重：系统间数据不互通"
]
add_multiline_text(slide, 6.0, 4.1, 6.5, 1.5, root_causes, 10, NEUTRAL_DARK)

# 解决方案（扩展，增加3项内容）
add_text(slide, 6.0, 5.6, 6.5, 0.35, "解决方案", 14, True, PRIMARY_DARK)
solutions = [
    "• 建设数字员工：打造四大核心能力智能体",
    "• 知识库整合：构建统一业务知识库",
    "• API集成：打通八大核心系统",
    "• 智能填单：自动化业务办理流程",
    "• 规则引擎：智能诊断问题根因",
    "• 数据打通：消除数据孤岛，信息共享"
]
add_multiline_text(slide, 6.0, 6.0, 6.5, 1.2, solutions, 10, NEUTRAL_DARK)

# ==================== P5: 建设目标 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.333, 7.5, WHITE)
add_rect(slide, 0, 0, 13.333, 0.12, PRIMARY)

add_header(slide, "建设目标", "四大能力集群作战")
add_target_line(slide, "办理时长缩短40%以上，构建智能问答、查询、办理、诊断四大核心能力智能体集群。")
add_measure_title(slide)

items = [
    ("智能问答 Q-Agent", "基于大模型+RAG技术，实现业务规则、资费政策、流程指引等知识快速问答，准确率≥90%，响应时间<3秒，大幅减少人工查询时间。"),
    ("智能查询 S-Agent", "通过实体抽取+API集成，实现客户资料、订单状态、资费信息等快速查询，效率提升50%以上，一键获取客户360度视图。"),
    ("智能办理 H-Agent", "基于多轮对话+智能填单，实现换卡、移机、停复机等业务自动办理，时长缩短40%，错误率降低80%。"),
    ("智能诊断 D-Agent", "结合规则引擎+云哨平台，实现订单卡单、办理失败等问题智能诊断，时间缩短60%，一键报障闭环处理。")
]
add_measure_content(slide, items)

add_text(slide, 6.0, 1.9, 6.5, 0.4, "关键指标", 14, True, PRIMARY_DARK)
metrics = [("40%", "办理时长"), ("50%", "查询效率"), ("60%", "诊断时间"), ("90%", "准确率")]
add_metric_dashboard(slide, 6.0, 2.4, 6.5, 1.8, metrics)

add_text(slide, 6.0, 4.5, 6.5, 0.4, "预期效果", 14, True, PRIMARY_DARK)
add_text(slide, 6.0, 5.0, 6.5, 1.0, "满足集团考核要求，一线效率提升40%以上，客户满意度提升15%，打造31省标杆项目。", 12, False, NEUTRAL_DARK)

# ==================== P6: 整体架构（大改：架构理念详细描述）====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.333, 7.5, WHITE)
add_rect(slide, 0, 0, 13.333, 0.12, PRIMARY)

add_header(slide, "整体架构", "六层分层清晰可落地")
add_target_line(slide, "架构清晰可落地，六层分层设计，每层职责明确，内容充实，便于开发运维。")

# 架构理念（阐述六层架构实现细节，每层不少于50字）
add_text(slide, 0.8, 1.6, 4.5, 0.35, "架构理念", 14, True, PRIMARY_DARK)

arch_concepts = [
    "【交互层】提供统一对话入口，支持PC端、移动端、语音等多种交互方式，降低使用门槛，一线营业员无需培训即可上手使用，提升用户体验。",
    "【智能体层】部署四大智能体（Q问答、S查询、H办理、D诊断），分工协作，专业精准，形成完整业务闭环，实现端到端智能服务。",
    "【能力层】构建意图识别、知识检索、智能填单三大核心能力，通过NLU引擎理解用户意图，RAG检索精准获取知识，自动化填充表单。",
    "【模型层】集成大语言模型（GPT/文心/通义）、RAG检索引擎、向量数据库（Milvus/ES），提供强大的AI能力支撑，确保准确率≥90%。",
    "【数据层】整合业务知识库（规则/资费/流程/产品）、客户数据（档案/套餐/余额/订单）、系统数据（配置/日志/接口/报文），消除数据孤岛。",
    "【集成层】对接CRM、随翼选、云哨等八大核心系统，通过API网关统一集成，标准化接口规范，实现系统间数据互通。"
]
add_multiline_text(slide, 0.8, 2.0, 4.5, 4.5, arch_concepts, 9, NEUTRAL_DARK)

# 右侧六层架构图
add_text(slide, 6.0, 1.6, 6.5, 0.35, "六层架构", 14, True, PRIMARY_DARK)

layers_data = [
    {"name": "交互层", "color": PRIMARY, "items": [("ChatUI入口", "统一对话界面"), ("PC/移动端", "多终端支持"), ("语音入口", "语音交互")]},
    {"name": "智能体层", "color": PRIMARY, "items": [("Q-Agent", "智能问答"), ("S-Agent", "智能查询"), ("H/D-Agent", "办理/诊断")]},
    {"name": "能力层", "color": SECONDARY, "items": [("意图识别", "NLU引擎"), ("知识检索", "RAG检索"), ("智能填单", "表单自动填充")]},
    {"name": "模型层", "color": SECONDARY, "items": [("大语言模型", "GPT/文心/通义"), ("RAG检索", "向量检索+重排序"), ("向量数据库", "Milvus/ES")]},
    {"name": "数据层", "color": NEUTRAL_DARK, "items": [("业务知识库", "规则/资费/流程/产品"), ("客户数据", "档案/套餐/余额/订单"), ("系统数据", "配置/日志/接口/报文")]},
    {"name": "集成层", "color": NEUTRAL_DARK, "items": [("CRM系统", "客户关系管理"), ("随翼选", "营业办理平台"), ("云哨平台", "智能诊断平台")]}
]

layer_height = 0.7
layer_gap = 0.05
start_y = 2.0

for i, layer in enumerate(layers_data):
    y = start_y + i * (layer_height + layer_gap)
    name_width = 1.1
    add_rounded_box(slide, 6.2, y, name_width, layer_height, layer["color"], layer["name"], 9, WHITE)
    
    item_width = 1.55
    item_gap = 0.06
    item_start_x = 6.2 + name_width + 0.1
    
    for j, (item_name, item_desc) in enumerate(layer["items"]):
        x = item_start_x + j * (item_width + item_gap)
        item_color = PRIMARY if i < 2 else (SECONDARY if i < 4 else NEUTRAL_DARK)
        add_rounded_box(slide, x, y, item_width, layer_height * 0.52, item_color, item_name, 8, WHITE)
        add_text(slide, x, y + layer_height * 0.55, item_width, layer_height * 0.35, item_desc, 6, False, NEUTRAL_DARK, PP_ALIGN.CENTER)

# ==================== P7: 智能问答能力（改变展现方式）====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.333, 7.5, WHITE)
add_rect(slide, 0, 0, 13.333, 0.12, PRIMARY)

add_header(slide, "智能问答能力 (Q-Agent)", "大模型加RAG技术")
add_target_line(slide, "准确率≥90%，基于大模型+RAG+向量检索，实现业务知识快速问答。")
add_measure_title(slide)

items = [
    ("核心业务场景", "覆盖业务规则查询、资费政策解答、流程指引、故障排查、政策解读等高频场景，日均咨询量超过5000次，人工客服压力巨大。"),
    ("技术实现路径", "用户自然语言输入→意图智能识别→知识向量检索→结果重排序→大模型生成答案→结构化输出，端到端响应时间<3秒。"),
    ("知识库体系设计", "构建业务规则库（办理条件、审批流程）、资费政策库（套餐详情、费用说明）、流程指引库（操作步骤、注意事项）、故障排查库（问题定位、处理建议）。")
]
add_measure_content(slide, items)

# 右侧：改变展现方式（垂直流程图）
add_text(slide, 6.0, 1.9, 6.5, 0.4, "技术流程", 14, True, PRIMARY_DARK)

steps = ["用户输入", "意图识别", "知识检索", "重排序", "模型生成", "答案输出"]
y = 2.4
for i, step in enumerate(steps):
    color = PRIMARY if i < 3 else SECONDARY
    add_rounded_box(slide, 6.5, y, 2.0, 0.5, color, step, 10, WHITE)
    if i < len(steps) - 1:
        add_down_arrow(slide, 7.4, y + 0.52, 0.2, 0.2, NEUTRAL)
    y += 0.75

add_text(slide, 9.0, 2.4, 3.3, 0.4, "应用场景", 14, True, PRIMARY_DARK)
add_text(slide, 9.0, 2.9, 3.3, 3.5, "• 业务规则查询\n• 资费政策解答\n• 流程指引\n• 故障排查\n• 政策解读", 11, False, NEUTRAL_DARK)

# ==================== P8: 智能查询能力（改变展现方式）====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.333, 7.5, WHITE)
add_rect(slide, 0, 0, 13.333, 0.12, PRIMARY)

add_header(slide, "智能查询能力 (S-Agent)", "实体抽取API集成")
add_target_line(slide, "效率提升50%，通过实体抽取+API集成，实现客户信息快速查询。")
add_measure_title(slide)

items = [
    ("核心业务场景", "覆盖客户资料查询（姓名、套餐、余额）、订单状态查询（进度、时间、结果）、资费信息查询（月租、流量、通话）、套餐余量查询（流量、短信、语音）等高频查询场景。"),
    ("技术实现路径", "用户自然语言输入→实体智能抽取→权限实时验证→API并行调用→结果整合格式化→结构化输出，支持手机号、客户号、订单号、身份证等多种实体识别。"),
    ("实体识别能力", "支持识别手机号（11位数字）、客户号（CRM编号）、订单号（业务流水号）、身份证（18位号码）、套餐名称、业务类型等6大类实体，准确率≥95%。")
]
add_measure_content(slide, items)

# 右侧：改变展现方式（实体识别展示）
add_text(slide, 6.0, 1.9, 6.5, 0.4, "实体识别能力", 14, True, PRIMARY_DARK)

entities = ["手机号", "客户号", "订单号", "身份证", "套餐名称", "业务类型"]
x = 6.0
for i, entity in enumerate(entities):
    color = PRIMARY if i < 2 else (SECONDARY if i < 4 else NEUTRAL_DARK)
    if i < 3:
        add_rounded_box(slide, x, 2.4, 2.0, 0.6, color, entity, 11, WHITE)
        x += 2.15
    else:
        if i == 3:
            x = 6.0
        add_rounded_box(slide, x, 3.2, 2.0, 0.6, color, entity, 11, WHITE)
        x += 2.15

add_text(slide, 6.0, 4.1, 6.5, 0.4, "查询类型", 14, True, PRIMARY_DARK)
add_text(slide, 6.0, 4.6, 6.5, 2.5, "• 客户资料查询\n• 订单状态查询\n• 资费信息查询\n• 套餐余量查询", 11, False, NEUTRAL_DARK)

# ==================== P9: 智能办理能力（改变展现方式）===================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.333, 7.5, WHITE)
add_rect(slide, 0, 0, 13.333, 0.12, PRIMARY)

add_header(slide, "智能办理能力 (H-Agent)", "多轮对话智能填单")
add_target_line(slide, "时长缩短40%，通过多轮对话+智能填单，实现业务自动办理。")
add_measure_title(slide)

items = [
    ("核心业务场景", "覆盖换卡业务（SIM卡更换、补卡）、移机业务（地址迁移、设备更换）、停复机（停机、复机申请）、可选包办理（流量包、语音包订购）、套餐变更等高频办理场景。"),
    ("技术实现路径", "用户自然语言输入→意图智能识别→多轮对话要素收集→业务规则校验→随翼选页面唤起→智能填单自动填充，实现端到端自动化办理。"),
    ("核心能力亮点", "多轮对话智能引导收集业务要素，支持上下文理解、意图澄清、要素补全；智能填单自动填充表单字段，填单准确率≥98%，大幅减少人工录入错误。")
]
add_measure_content(slide, items)

# 右侧：改变展现方式（业务场景卡片）
add_text(slide, 6.0, 1.9, 6.5, 0.4, "业务场景", 14, True, PRIMARY_DARK)

scenarios = [
    ("换卡业务", "SIM卡更换、补卡"),
    ("移机业务", "地址迁移、设备更换"),
    ("停复机", "停机、复机申请"),
    ("可选包办理", "流量包、语音包订购")
]

x = 6.0
y = 2.4
for i, (name, desc) in enumerate(scenarios):
    color = PRIMARY if i < 2 else SECONDARY
    add_rounded_box(slide, x, y, 3.1, 0.55, color, name, 11, WHITE)
    add_text(slide, x, y + 0.58, 3.1, 0.4, desc, 9, False, NEUTRAL_DARK, PP_ALIGN.CENTER)
    x += 3.25
    if i == 1:
        x = 6.0
        y = 3.3

add_text(slide, 6.0, 4.5, 6.5, 0.4, "核心能力", 14, True, PRIMARY_DARK)
add_text(slide, 6.0, 5.0, 6.5, 2.0, "• 多轮对话智能引导\n• 上下文理解\n• 意图澄清\n• 要素补全\n• 智能填单", 11, False, NEUTRAL_DARK)

# ==================== P10: 智能诊断能力（改变展现方式）====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.333, 7.5, WHITE)
add_rect(slide, 0, 0, 13.333, 0.12, PRIMARY)

add_header(slide, "智能诊断能力 (D-Agent)", "规则引擎云哨对接")
add_target_line(slide, "时间缩短60%，通过规则引擎+云哨对接，实现问题智能诊断。")
add_measure_title(slide)

items = [
    ("核心业务场景", "覆盖订单卡单（订单状态异常）、办理失败（失败原因定位）、资源异常（资源分配问题）、计费异常（费用差异问题）、系统故障（系统报错诊断）等高频诊断场景。"),
    ("技术实现路径", "问题描述自然语言输入→关键信息智能收集→规则引擎知识库双重诊断→诊断结果生成→处理建议输出→一键报障闭环处理，诊断准确率≥85%。"),
    ("核心能力亮点", "智能识别问题类型，结合规则引擎（业务规则匹配）+知识库（历史案例检索）双重诊断，对接云哨平台实现自动化诊断，诊断时间从平均30分钟缩短至10分钟。")
]
add_measure_content(slide, items)

# 右侧：改变展现方式（诊断类型展示）
add_text(slide, 6.0, 1.9, 6.5, 0.4, "诊断类型", 14, True, PRIMARY_DARK)

diag_types = [
    ("订单卡单", "订单状态异常诊断"),
    ("办理失败", "失败原因定位"),
    ("资源异常", "资源分配问题"),
    ("计费异常", "费用差异问题"),
    ("系统故障", "系统报错诊断")
]

y = 2.4
for i, (name, desc) in enumerate(diag_types):
    color = PRIMARY if i < 2 else SECONDARY
    add_rounded_box(slide, 6.0, y, 2.0, 0.55, color, name, 10, WHITE)
    add_text(slide, 8.1, y + 0.1, 4.0, 0.4, desc, 10, False, NEUTRAL_DARK, PP_ALIGN.LEFT)
    y += 0.65

add_text(slide, 6.0, 5.8, 6.5, 0.4, "核心能力", 14, True, PRIMARY_DARK)
add_text(slide, 6.0, 6.2, 6.5, 0.8, "规则引擎+知识库双重诊断，一键报障闭环处理", 10, False, NEUTRAL_DARK)

# ==================== P11: 系统集成（改变展现方式）====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.333, 7.5, WHITE)
add_rect(slide, 0, 0, 13.333, 0.12, PRIMARY)

add_header(slide, "系统集成", "API网关统一集成")
add_target_line(slide, "八大系统对接，通过API网关统一集成+标准化接口规范，降低集成风险。")
add_measure_title(slide)

items = [
    ("集成架构设计", "采用API网关统一入口，标准化接口规范，分步对接降低风险，统一认证鉴权确保安全。支持RESTful API、WebService等多种接口协议，适配不同系统架构。"),
    ("核心系统对接", "CRM系统（客户关系管理）、随翼选（营业办理平台）、客户中心（客户档案）、订单中心（订单流转）、计费中心（费用结算）、资源中心（资源分配）、云哨平台（智能诊断）七大核心系统全覆盖。")
]
add_measure_content(slide, items)

# 右侧：改变展现方式（系统集成架构图）
add_text(slide, 6.0, 1.9, 6.5, 0.4, "系统集成架构", 14, True, PRIMARY_DARK)

# API网关居中
add_rounded_box(slide, 7.5, 2.4, 2.0, 0.7, PRIMARY, "API网关", 12, WHITE)

# 上游系统
upstream = ["ChatUI", "PC端", "移动端"]
x = 6.2
for i, name in enumerate(upstream):
    add_rounded_box(slide, x, 3.4, 1.3, 0.5, SECONDARY, name, 9, WHITE)
    x += 1.45

# 下游系统
downstream = ["CRM", "随翼选", "客户中心", "订单中心", "计费中心", "云哨"]
x = 6.0
for i, name in enumerate(downstream):
    add_rounded_box(slide, x, 4.2, 1.0, 0.5, NEUTRAL_DARK, name, 8, WHITE)
    x += 1.1

# 连接线（简化）
add_text(slide, 6.0, 4.9, 6.5, 0.4, "集成方式", 14, True, PRIMARY_DARK)
add_text(slide, 6.0, 5.3, 6.5, 1.5, "• API网关统一入口\n• 标准化接口规范\n• 统一认证鉴权\n• 分步对接降低风险", 11, False, NEUTRAL_DARK)

# ==================== P12: 实施路径（改变展现方式）====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.333, 7.5, WHITE)
add_rect(slide, 0, 0, 13.333, 0.12, PRIMARY)

add_header(slide, "实施路径", "五个月四阶段落地")
add_target_line(slide, "5个月完成建设，四阶段渐进式实施，风险可控、价值快速释放。")
add_measure_title(slide)

items = [
    ("第一阶段：基础能力建设", "历时2个月，完成智能问答+智能查询两大基础能力建设，准确率≥90%，覆盖Top10高频场景，快速释放价值。"),
    ("第二阶段：核心价值实现", "历时2个月，完成智能办理核心能力建设，实现换卡、移机、停复机等高频业务自动办理，时长缩短≥30%，一线反馈良好。"),
    ("第三阶段：完善闭环构建", "历时1个月，完成智能诊断能力建设，实现订单卡单、办理失败等问题智能诊断，准确率≥80%，形成完整业务闭环。"),
    ("第四阶段：持续优化迭代", "长期持续，越用越聪明，扩展更多业务场景，优化用户体验，打造31省标杆项目。")
]
add_measure_content(slide, items)

# 右侧：改变展现方式（时间轴展示）
add_text(slide, 6.0, 1.9, 6.5, 0.4, "实施时间轴", 14, True, PRIMARY_DARK)

phases = [
    ("第一阶段", "2个月", "问答+查询"),
    ("第二阶段", "2个月", "智能办理"),
    ("第三阶段", "1个月", "智能诊断"),
    ("第四阶段", "持续", "优化迭代")
]

y = 2.4
for i, (name, time, desc) in enumerate(phases):
    color = PRIMARY if i < 2 else SECONDARY
    add_rounded_box(slide, 6.0, y, 1.8, 0.6, color, name, 10, WHITE)
    add_rounded_box(slide, 8.0, y, 1.2, 0.6, SECONDARY, time, 10, WHITE)
    add_text(slide, 9.4, y + 0.15, 3.0, 0.4, desc, 10, False, NEUTRAL_DARK, PP_ALIGN.LEFT)
    if i < len(phases) - 1:
        add_down_arrow(slide, 7.0, y + 0.65, 0.15, 0.15, NEUTRAL)
    y += 0.85

# ==================== P13: 关键指标（改变展现方式）====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.333, 7.5, WHITE)
add_rect(slide, 0, 0, 13.333, 0.12, PRIMARY)

add_header(slide, "关键指标", "效果可量化可追踪")
add_target_line(slide, "效果可量化，建立业务+技术双维度指标体系，实时追踪优化。")
add_measure_title(slide)

items = [
    ("效率提升指标", "办理时长缩短40%（从平均20分钟降至12分钟），查询效率提升50%（从平均5分钟降至2.5分钟），诊断时间缩短60%（从平均30分钟降至10分钟）。"),
    ("质量保障指标", "问答准确率≥90%（基于大模型+RAG技术），填单准确率≥98%（智能填单自动校验），诊断准确率≥85%（规则引擎+知识库双重诊断）。"),
    ("用户体验指标", "用户满意度≥85%（一线营业员反馈），客户满意度提升15%（客户等待时间缩短），响应时间<3秒（端到端延迟控制）。")
]
add_measure_content(slide, items)

# 右侧：改变展现方式（指标卡片）
add_text(slide, 6.0, 1.9, 6.5, 0.4, "效率指标", 14, True, PRIMARY_DARK)
metrics1 = [("40%", "时长缩短"), ("50%", "效率提升"), ("60%", "时间缩短")]
add_metric_dashboard(slide, 6.0, 2.4, 6.5, 1.5, metrics1)

add_text(slide, 6.0, 4.2, 6.5, 0.4, "质量指标", 14, True, PRIMARY_DARK)
metrics2 = [("≥90%", "问答准确"), ("≥98%", "填单准确"), ("≥85%", "诊断准确")]
add_metric_dashboard(slide, 6.0, 4.7, 6.5, 1.5, metrics2)

# ==================== P14: 核心价值（改变展现方式）====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.333, 7.5, WHITE)
add_rect(slide, 0, 0, 13.333, 0.12, PRIMARY)

add_header(slide, "核心价值", "四大价值全面释放")
add_target_line(slide, "四大价值维度：考核达标、效率提升、体验改善、人力释放，全面赋能一线。")
add_measure_title(slide)

items = [
    ("考核达标价值", "满足集团云运〔2025〕4号考核要求，AI嵌入10个关键场景，自动化率达70%，占考核指标10分，打造31省标杆项目。"),
    ("效率提升价值", "办理时长缩短40%，查询效率提升50%，诊断时间缩短60%，一线营业员日均可多服务8-10名客户，整体效率显著提升。"),
    ("体验改善价值", "客户平均等待时间从15分钟降至8分钟，客户满意度提升15%，一线营业员工作压力减轻，服务体验明显改善。"),
    ("人力释放价值", "推动一线营业员从简单重复劳动转向高价值客户服务，从事营销推广、客户维系等核心业务，实现减负增效。")
]
add_measure_content(slide, items)

# 右侧：改变展现方式（价值卡片）
add_text(slide, 6.0, 1.9, 6.5, 0.4, "核心价值", 14, True, PRIMARY_DARK)

values = [
    ("考核达标", "满足集团要求"),
    ("效率提升", "效率提升40%"),
    ("体验改善", "满意度提升15%"),
    ("人力释放", "减负增效")
]

x = 6.0
y = 2.4
for i, (name, desc) in enumerate(values):
    color = PRIMARY if i < 2 else SECONDARY
    add_rounded_box(slide, x, y, 3.1, 0.6, color, name, 11, WHITE)
    add_text(slide, x, y + 0.65, 3.1, 0.4, desc, 9, False, NEUTRAL_DARK, PP_ALIGN.CENTER)
    x += 3.25
    if i == 1:
        x = 6.0
        y = 3.4

# ==================== P15: 创新亮点（改变展现方式）====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.333, 7.5, WHITE)
add_rect(slide, 0, 0, 13.333, 0.12, PRIMARY)

add_header(slide, "创新亮点", "五大创新引领行业")
add_target_line(slide, "引领行业AI应用，五大创新技术融合：智能体集群、多模态交互、主动推荐、学习型系统、场景化卡片。")
add_measure_title(slide)

items = [
    ("智能体集群创新", "四大智能体（Q-Agent问答、S-Agent查询、H-Agent办理、D-Agent诊断）分工协作，专业精准，形成完整业务闭环。"),
    ("多模态交互创新", "支持语音+文字+图片多种输入方式，灵活便捷，降低使用门槛，一线营业员无需培训即可上手使用。"),
    ("主动推荐创新", "基于客户画像分析，实现千人千面精准推荐，智能推送适合客户的套餐、产品、服务，提升营销效率。"),
    ("学习型系统创新", "持续学习用户反馈和业务数据，越用越聪明，自我进化，适应业务变化，准确率持续提升。")
]
add_measure_content(slide, items)

# 右侧：改变展现方式（创新展示）
add_text(slide, 6.0, 1.9, 6.5, 0.4, "五大创新", 14, True, PRIMARY_DARK)

innovations = [
    ("智能体集群", "四大Agent协作"),
    ("多模态交互", "语音+文字+图片"),
    ("主动推荐", "千人千面"),
    ("学习型系统", "越用越聪明"),
    ("场景化卡片", "快速响应")
]

y = 2.4
for i, (name, desc) in enumerate(innovations):
    color = PRIMARY if i < 2 else SECONDARY
    add_rounded_box(slide, 6.0, y, 2.2, 0.55, color, name, 10, WHITE)
    add_text(slide, 8.4, y + 0.1, 4.0, 0.4, desc, 10, False, NEUTRAL_DARK, PP_ALIGN.LEFT)
    y += 0.65

# ==================== P16: 成功案例（改变展现方式）====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.333, 7.5, WHITE)
add_rect(slide, 0, 0, 13.333, 0.12, PRIMARY)

add_header(slide, "成功案例", "三地成功可复制")
add_target_line(slide, "三地成功实践，复制成熟经验，降低实施风险，加速落地见效。")
add_measure_title(slide)

items = [
    ("陕西电信案例", "营业办理数字员工项目，覆盖智能问答、智能查询、智能办理三大能力，办理时长缩短40%，一线营业员效率显著提升，客户满意度提升12%。"),
    ("重庆电信案例", "营业受理数字员工项目，聚焦智能问答+智能查询两大基础能力，满足集团考核要求，自动化率达75%，成为省级标杆项目。"),
    ("安徽电信案例", "CRM域AI应用项目，四大能力（问答、查询、办理、诊断）全覆盖，打造31省标杆，经验可复制推广。")
]
add_measure_content(slide, items)

# 右侧：改变展现方式（案例卡片）
add_text(slide, 6.0, 1.9, 6.5, 0.4, "成功案例", 14, True, PRIMARY_DARK)

cases = [
    ("陕西电信", "时长缩短40%"),
    ("重庆电信", "自动化率75%"),
    ("安徽电信", "31省标杆")
]

y = 2.4
for i, (name, desc) in enumerate(cases):
    color = PRIMARY if i == 0 else SECONDARY
    add_rounded_box(slide, 6.0, y, 2.5, 0.7, color, name, 12, WHITE)
    add_text(slide, 8.7, y + 0.2, 3.5, 0.4, desc, 11, False, NEUTRAL_DARK, PP_ALIGN.LEFT)
    y += 0.9

# ==================== P17: 总结与展望（改变展现方式）====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.333, 7.5, WHITE)
add_rect(slide, 0, 0, 13.333, 0.12, PRIMARY)

add_header(slide, "总结与展望", "方案成熟风险可控")
add_target_line(slide, "方案成熟可落地，分阶段实施，风险可控，价值明确，打造31省标杆。")
add_measure_title(slide)

items = [
    ("方案成熟可靠", "参考陕西、重庆、安徽三地成功案例，技术方案验证可行，实施方法论完善，风险可控。"),
    ("技术先进可行", "大模型+RAG技术成熟，开源组件丰富（GPT、文心、通义、Milvus等），技术风险低，可快速落地。"),
    ("价值明确可期", "效率提升40%以上，满足集团考核要求，客户体验明显改善，一线营业员工作压力减轻。"),
    ("后续规划清晰", "持续优化迭代，扩展更多业务场景，复制推广到其他省份，打造31省标杆项目，引领行业AI应用。")
]
add_measure_content(slide, items)

# 右侧：改变展现方式（总结卡片）
add_text(slide, 6.0, 1.9, 6.5, 0.4, "总结", 14, True, PRIMARY_DARK)

summaries = [
    ("方案成熟", "三地验证"),
    ("技术可行", "技术风险低"),
    ("价值明确", "效率提升40%"),
    ("风险可控", "分阶段实施")
]

y = 2.4
for i, (name, desc) in enumerate(summaries):
    color = PRIMARY if i < 2 else SECONDARY
    add_rounded_box(slide, 6.0, y, 2.5, 0.6, color, name, 11, WHITE)
    add_text(slide, 8.7, y + 0.15, 3.5, 0.4, desc, 10, False, NEUTRAL_DARK, PP_ALIGN.LEFT)
    y += 0.75

# ==================== P18: 结束页 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.333, 7.5, WHITE)
add_rect(slide, 0, 0, 13.333, 0.12, PRIMARY)
add_text(slide, 1, 2.5, 11.333, 1.0, "谢谢！", 48, True, PRIMARY_DARK, PP_ALIGN.CENTER)
add_text(slide, 1, 3.5, 11.333, 0.6, "让我们一起，用AI赋能一线，提升效率，改善体验！", 18, False, NEUTRAL_DARK, PP_ALIGN.CENTER)

# 保存
output_path = "/root/.openclaw/workspace/03_项目管理/营业受理数字员工/03_工程输出/指南针工程/营业受理数字员工建设方案_精进版.pptx"
prs.save(output_path)
print(f"\n✅ 精进版PPT已生成: {output_path}")
print(f"总页数: {len(prs.slides)}")
print("\n修改点：")
print("1. 目录页：目录内容全部居中对齐")
print("2. 问题分析页：解决方案设计细化，右图一行排版，根因扩充，解决方案扩展")
print("3. 整体架构页：架构理念阐述六层细节，每层不少于50字")
print("4. 7-18页：改变图形展现方式，避免单调")
print("5. 整体PPT：保持整洁，不得越界，注意排版")
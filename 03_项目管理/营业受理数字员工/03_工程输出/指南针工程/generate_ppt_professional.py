#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""营业受理数字员工建设方案 - 专业版
使用素材库中的专业图标和架构图组件

核心改进：
1. 嵌入专业图标（166个PNG/SVG）
2. 嵌入专业架构图组件（323个）
3. 严格遵循设计规范库
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# ==================== 素材库路径 ====================
ASSETS_BASE = "/root/.openclaw/workspace/02_工作台/美学工坊/素材库/PPT素材库/素材索引"
ICONS_DIR = f"{ASSETS_BASE}/icons/ppt/media"
LOGICS_DIR = f"{ASSETS_BASE}/logics/ppt/media"

# ==================== 设计规范（从设计规范库导入）====================
# 配色
TELECOM_RED = RGBColor(201, 56, 50)
DARK_RED = RGBColor(172, 0, 0)
ORANGE = RGBColor(237, 125, 49)
BLUE = RGBColor(0, 110, 189)
LIGHT_BLUE = RGBColor(79, 129, 189)
DARK_BLUE = RGBColor(31, 73, 125)
GRAY = RGBColor(166, 166, 166)
DARK_GRAY = RGBColor(89, 89, 89)
LIGHT_GRAY = RGBColor(240, 240, 240)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)
GREEN = RGBColor(84, 130, 53)

# 字体
FONT_NAME = "微软雅黑"

# ==================== 创建演示文稿 ====================
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ==================== 辅助函数 ====================

def add_text(slide, left, top, width, height, text, size=12, bold=False, color=BLACK, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = FONT_NAME
    p.alignment = align
    return tb

def add_box(slide, left, top, width, height, fill_color, text="", size=12, text_color=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 
                                    Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = text_color
        p.font.bold = True
        p.font.name = FONT_NAME
        p.alignment = PP_ALIGN.CENTER
    return shape

def add_rect(slide, left, top, width, height, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 
                               Inches(left), Inches(top), Inches(width), Inches(height))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def add_arrow(slide, left, top, width, height, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, 
                               Inches(left), Inches(top), Inches(width), Inches(height))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def add_icon(slide, left, top, width, height, icon_path):
    """嵌入专业图标"""
    if os.path.exists(icon_path):
        try:
            slide.shapes.add_picture(icon_path, Inches(left), Inches(top), Inches(width), Inches(height))
            return True
        except:
            return False
    return False

def add_header(slide, title):
    """添加页面标题（统一风格）"""
    add_rect(slide, 0.5, 0.4, 0.12, 0.5, TELECOM_RED)
    add_text(slide, 0.8, 0.4, 5, 0.5, title, 24, True, DARK_RED)

def add_target_line(slide, text):
    """添加建设背景和目标行"""
    add_text(slide, 0.8, 1.1, 11.5, 0.4, text, 16, True, BLACK)

def add_measure_title(slide):
    """添加"重点举措"标题"""
    add_text(slide, 0.8, 1.8, 4.5, 0.4, "重点举措", 14, True, DARK_RED)

def add_measure_content(slide, items):
    """添加左侧举措内容"""
    y = 2.5
    for title, desc in items:
        add_text(slide, 0.8, y, 4.5, 0.35, title, 14, True, BLACK)
        y += 0.35
        add_text(slide, 0.8, y, 4.5, 0.6, desc, 12, False, DARK_GRAY)
        y += 0.7

def add_diagram_area(slide, left=5.7, top=1.8, width=7.3, height=5.0):
    """添加右侧图示区域背景"""
    add_rect(slide, left, top, width, height, LIGHT_GRAY)

# 获取可用的图标文件
def get_available_icons():
    """获取可用的图标文件列表"""
    icons = []
    if os.path.exists(ICONS_DIR):
        for f in os.listdir(ICONS_DIR):
            if f.endswith('.png'):
                icons.append(os.path.join(ICONS_DIR, f))
    return icons

def get_available_logics():
    """获取可用的架构图组件列表"""
    logics = []
    if os.path.exists(LOGICS_DIR):
        for f in os.listdir(LOGICS_DIR):
            if f.endswith('.png'):
                logics.append(os.path.join(LOGICS_DIR, f))
    return logics

# 获取素材
icons = get_available_icons()
logics = get_available_logics()
print(f"可用图标: {len(icons)}个")
print(f"可用架构图组件: {len(logics)}个")

# ==================== P1: 封面 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.333, 7.5, WHITE)
add_rect(slide, 0, 0, 13.333, 0.12, TELECOM_RED)
add_text(slide, 1, 2.0, 11.333, 0.8, "营业受理数字员工建设方案", 36, True, DARK_RED, PP_ALIGN.CENTER)
add_text(slide, 1, 2.9, 11.333, 0.5, "集团考核达标 | 一线效率提升 | 客户体验改善", 18, False, DARK_GRAY, PP_ALIGN.CENTER)
add_rect(slide, 5, 3.5, 3.333, 0.02, TELECOM_RED)
add_text(slide, 1, 4.5, 11.333, 0.4, "中国电信湖北分公司    智能云网业务运营中心", 14, False, DARK_GRAY, PP_ALIGN.CENTER)
add_text(slide, 1, 5.0, 11.333, 0.4, "2026年3月", 14, False, GRAY, PP_ALIGN.CENTER)

# 尝试嵌入装饰图标
if len(icons) > 0:
    add_icon(slide, 11.5, 0.5, 1.2, 1.2, icons[0])

# ==================== P2: 目录 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.333, 7.5, WHITE)
add_rect(slide, 0, 0, 13.333, 0.12, TELECOM_RED)
add_text(slide, 0.5, 0.5, 2, 0.6, "目 录", 28, True, DARK_RED)
add_text(slide, 0.5, 1.0, 2, 0.4, "CONTENTS", 14, False, GRAY)

items = [
    ("01", "背景与目标", "项目背景、问题分析、建设目标"),
    ("02", "方案设计", "整体架构、四大能力、技术架构"),
    ("03", "实施规划", "系统集成、实施路径、关键指标"),
    ("04", "价值与优势", "核心价值、创新亮点、成功案例")
]
y = 1.8
for i, (num, title, desc) in enumerate(items):
    add_box(slide, 0.8, y, 0.6, 0.6, TELECOM_RED, num, 16, WHITE)
    add_text(slide, 1.6, y, 4, 0.4, title, 18, True, DARK_RED)
    add_text(slide, 1.6, y + 0.4, 4, 0.4, desc, 12, False, GRAY)
    # 嵌入装饰图标
    if len(icons) > i + 1:
        add_icon(slide, 5.5, y + 0.1, 0.5, 0.5, icons[i + 1])
    y += 1.2

# ==================== P3: 项目背景（嵌入专业图标）====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.333, 7.5, WHITE)
add_rect(slide, 0, 0, 13.333, 0.12, TELECOM_RED)

add_header(slide, "项目背景")
add_target_line(slide, "满足集团考核要求，解决一线营业员效率低、体验差的问题。")
add_measure_title(slide)

items = [
    ("调研分析", "调研12个分公司40+场景需求，Top3需求为智能问答、智能查询、智能诊断。"),
    ("考核对标", "云运〔2025〕4号要求AI嵌入10个关键场景，自动化率达70%，占指标10%。"),
    ("问题定位", "一线营业员业务办理依赖人工操作，效率低易出错，客户等待时间长。")
]
add_measure_content(slide, items)

add_diagram_area(slide)

# 嵌入专业指标展示图标
add_text(slide, 7.0, 2.0, 5, 0.4, "集团考核指标", 14, True, DARK_RED)

# 使用圆角矩形+嵌入图标的方式展示指标
if len(icons) > 10:
    # 电信红指标卡片
    add_box(slide, 7.0, 2.6, 2.0, 2.0, TELECOM_RED, "", 12, WHITE)
    add_text(slide, 7.0, 3.0, 2.0, 0.6, "10%", 28, True, WHITE, PP_ALIGN.CENTER)
    add_text(slide, 7.0, 3.6, 2.0, 0.4, "AI场景指标", 11, False, WHITE, PP_ALIGN.CENTER)
    # 嵌入图标
    add_icon(slide, 7.7, 4.7, 0.5, 0.5, icons[10])

if len(icons) > 11:
    # 橙色指标卡片
    add_box(slide, 9.5, 2.6, 2.0, 2.0, ORANGE, "", 12, WHITE)
    add_text(slide, 9.5, 3.0, 2.0, 0.6, "70%", 28, True, WHITE, PP_ALIGN.CENTER)
    add_text(slide, 9.5, 3.6, 2.0, 0.4, "自动化率", 11, False, WHITE, PP_ALIGN.CENTER)
    add_icon(slide, 10.2, 4.7, 0.5, 0.5, icons[11])

add_text(slide, 7.0, 5.3, 5, 0.4, "调研结论", 14, True, DARK_RED)
add_text(slide, 7.0, 5.8, 5, 0.8, "必须建设数字员工，满足考核要求，解决一线痛点。", 12, False, DARK_GRAY)

# ==================== P4: 问题分析（嵌入专业图标）====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.333, 7.5, WHITE)
add_rect(slide, 0, 0, 13.333, 0.12, TELECOM_RED)

add_header(slide, "问题分析")
add_target_line(slide, "解决核心根因问题，建设数字员工赋能一线营业员。")
add_measure_title(slide)

items = [
    ("问题表象", "效率低：业务办理依赖人工，耗时长易出错。体验差：客户等待时间长。"),
    ("根因分析", "核心根因：缺乏智能化工具赋能，一线营业员无AI助手辅助。"),
    ("解决方案", "建设营业受理数字员工，四大核心能力赋能一线，提升效率改善体验。")
]
add_measure_content(slide, items)

add_diagram_area(slide)
add_text(slide, 7.0, 2.0, 5, 0.4, "问题诊断", 14, True, DARK_RED)

# 问题卡片+图标
if len(icons) > 20:
    add_box(slide, 7.2, 2.6, 2.0, 0.9, RGBColor(200, 50, 50), "效率低", 13, WHITE)
    add_icon(slide, 7.3, 2.7, 0.3, 0.3, icons[20])
    
if len(icons) > 21:
    add_box(slide, 9.6, 2.6, 2.0, 0.9, RGBColor(200, 100, 50), "体验差", 13, WHITE)
    add_icon(slide, 9.7, 2.7, 0.3, 0.3, icons[21])

if len(icons) > 22:
    add_box(slide, 7.2, 3.7, 2.0, 0.9, RGBColor(200, 150, 50), "获取难", 13, WHITE)
    add_icon(slide, 7.3, 3.8, 0.3, 0.3, icons[22])

add_box(slide, 7.5, 4.8, 3.8, 0.8, ORANGE, "核心根因：缺乏工具", 12, WHITE)
add_box(slide, 7.5, 5.8, 3.8, 0.8, GREEN, "解决方案：数字员工", 12, WHITE)

# ==================== P5: 建设目标（嵌入专业图标）====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.333, 7.5, WHITE)
add_rect(slide, 0, 0, 13.333, 0.12, TELECOM_RED)

add_header(slide, "建设目标")
add_target_line(slide, "办理时长缩短40%，构建四大核心能力智能体集群。")
add_measure_title(slide)

items = [
    ("智能问答 Q-Agent", "准确率≥90%，基于大模型+RAG技术，快速解答业务知识。"),
    ("智能查询 S-Agent", "效率提升50%，实体抽取+API集成，客户信息快速获取。"),
    ("智能办理 H-Agent", "时长缩短40%，多轮对话+智能填单，业务办理自动化。"),
    ("智能诊断 D-Agent", "时间缩短60%，规则引擎+云哨对接，问题快速定位。")
]
add_measure_content(slide, items)

add_diagram_area(slide)
add_text(slide, 7.0, 2.0, 5, 0.4, "关键指标", 14, True, DARK_RED)

# 指标仪表盘+图标
metrics = [("40%", "办理时长"), ("50%", "查询效率"), ("60%", "诊断时间"), ("90%", "问答准确率")]
x = 7.0
for i, (value, label) in enumerate(metrics):
    add_box(slide, x, 2.5, 1.3, 1.5, TELECOM_RED, "", 12, WHITE)
    add_text(slide, x, 2.7, 1.3, 0.5, value, 20, True, WHITE, PP_ALIGN.CENTER)
    add_text(slide, x, 3.3, 1.3, 0.3, label, 9, False, WHITE, PP_ALIGN.CENTER)
    # 嵌入图标
    if len(icons) > 30 + i:
        add_icon(slide, x + 0.4, 3.7, 0.5, 0.5, icons[30 + i])
    x += 1.5

add_text(slide, 7.0, 4.5, 5, 0.4, "预期效果", 14, True, DARK_RED)
add_text(slide, 7.0, 5.0, 5, 1.0, "满足集团考核要求，一线效率提升40%以上，客户满意度提升15%。", 12, False, DARK_GRAY)

# ==================== P6: 整体架构（嵌入专业架构图组件）====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.333, 7.5, WHITE)
add_rect(slide, 0, 0, 13.333, 0.12, TELECOM_RED)

add_header(slide, "整体架构")
add_target_line(slide, "架构清晰可落地，六层分层设计+标准化接口。")
add_measure_title(slide)

items = [
    ("架构理念", "分层设计，每层职责明确，便于开发、测试、运维，支持横向扩展。"),
    ("交互层", "ChatUI统一入口，PC端、移动端、随翼选、语音入口多终端支持。"),
    ("智能体层", "四大智能体分工协作：Q-Agent问答、S-Agent查询、H-Agent办理、D-Agent诊断。"),
    ("能力层", "AI能力组件复用：意图识别、实体抽取、知识检索、智能填单、多轮对话。")
]
add_measure_content(slide, items)

add_diagram_area(slide)
add_text(slide, 7.0, 2.0, 5, 0.4, "六层架构", 14, True, DARK_RED)

# 尝试嵌入专业架构图组件
if len(logics) > 0:
    # 嵌入第一个可用的架构图组件
    add_icon(slide, 7.0, 2.5, 5.5, 4.0, logics[0])
else:
    # 备用：用形状绘制
    layers = [
        ("交互层 (ChatUI)", RGBColor(100, 149, 237)),
        ("智能体集群", TELECOM_RED),
        ("AI能力层", ORANGE),
        ("模型层", BLUE),
        ("数据层", LIGHT_BLUE),
        ("系统集成层", DARK_GRAY)
    ]
    y = 2.5
    for name, color in layers:
        add_box(slide, 7.0, y, 5.0, 0.65, color, name, 11, WHITE)
        y += 0.75

# ==================== P7-10: 四大能力（嵌入专业图标）====================
# 为了节省篇幅，这里只展示一个能力页的完整代码，其他类似

capabilities = [
    ("智能问答能力 (Q-Agent)", "准确率≥90%", "大模型+RAG+向量检索",
     "业务规则查询、资费政策解答、流程指引、故障排查、政策解读。",
     "用户输入→意图识别→知识检索→重排序→大模型生成→答案输出。",
     "业务规则、资费政策、流程指引、产品介绍、故障排查、政策解读。"),
    
    ("智能查询能力 (S-Agent)", "效率提升50%", "实体抽取+API集成",
     "客户资料查询、订单状态查询、资费信息查询、套餐余量查询。",
     "用户输入→实体抽取→权限验证→API调用→结果整合→格式输出。",
     "手机号、客户号、订单号、身份证、套餐名称、业务类型。"),
    
    ("智能办理能力 (H-Agent)", "时长缩短40%", "多轮对话+智能填单",
     "换卡业务、移机业务、停复机、可选包办理、套餐变更。",
     "用户输入→意图识别→要素收集→规则校验→页面唤起→智能填单。",
     "多轮对话引导收集要素，智能填单自动填充表单，随翼选页面集成。"),
    
    ("智能诊断能力 (D-Agent)", "时间缩短60%", "规则引擎+云哨对接",
     "订单卡单、办理失败、资源异常、计费异常、系统故障。",
     "问题输入→信息收集→诊断分析→结果生成→处理建议→一键报障。",
     "智能识别问题类型，规则引擎+知识库双重诊断，对接云哨平台。")
]

for idx, (title, target, measure, scene, tech, capability) in enumerate(capabilities):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, 13.333, 7.5, WHITE)
    add_rect(slide, 0, 0, 13.333, 0.12, TELECOM_RED)
    
    add_header(slide, title)
    add_target_line(slide, f"{target}，{measure}。")
    add_measure_title(slide)
    
    items = [
        ("业务场景", scene),
        ("技术实现", tech),
        ("核心能力", capability)
    ]
    add_measure_content(slide, items)
    
    add_diagram_area(slide)
    add_text(slide, 7.0, 2.0, 5, 0.4, "技术流程", 14, True, DARK_RED)
    
    # 流程图（嵌入图标）
    flow_colors = [TELECOM_RED, ORANGE, BLUE, LIGHT_BLUE, GREEN, DARK_GRAY]
    flow_x = 7.0
    for i, color in enumerate(flow_colors):
        add_box(slide, flow_x, 2.5, 1.3, 0.7, color, f"步骤{i+1}", 9, WHITE)
        if len(icons) > 40 + idx * 6 + i:
            add_icon(slide, flow_x + 0.4, 2.55, 0.3, 0.3, icons[40 + idx * 6 + i])
        if i < 5:
            add_arrow(slide, flow_x + 1.35, 2.65, 0.15, 0.2, GRAY)
        flow_x += 1.5
    
    add_text(slide, 7.0, 3.5, 5, 0.4, "应用场景", 14, True, DARK_RED)
    add_text(slide, 7.0, 4.0, 5, 2.0, f"• {scene.replace('、', '\n• ')}", 11, False, DARK_GRAY)

# ==================== P11: 系统集成（嵌入专业拓扑图组件）====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.333, 7.5, WHITE)
add_rect(slide, 0, 0, 13.333, 0.12, TELECOM_RED)

add_header(slide, "系统集成")
add_target_line(slide, "八大系统对接，API网关统一集成+标准化接口。")
add_measure_title(slide)

items = [
    ("集成架构", "API网关统一入口，标准化接口规范，分步对接降低风险，统一认证鉴权。"),
    ("核心系统", "CRM、随翼选、客户中心、订单中心、计费中心、资源中心、云哨。")
]
add_measure_content(slide, items)

add_diagram_area(slide)
add_text(slide, 7.0, 2.0, 5, 0.4, "系统集成架构", 14, True, DARK_RED)

# 尝试嵌入专业拓扑图组件
if len(logics) > 100:
    add_icon(slide, 7.0, 2.5, 5.5, 4.0, logics[100])
else:
    # 备用：用形状绘制
    add_box(slide, 8.0, 3.2, 2.5, 1.0, DARK_RED, "数字员工\nAPI网关", 12, WHITE)
    systems = [
        (7.0, 2.5, "CRM", TELECOM_RED),
        (10.5, 2.5, "随翼选", ORANGE),
        (7.0, 4.8, "客户中心", BLUE),
        (10.5, 4.8, "订单中心", LIGHT_BLUE),
        (7.0, 5.6, "计费中心", GREEN),
        (10.5, 5.6, "云哨", DARK_GRAY)
    ]
    for x, y, name, color in systems:
        add_box(slide, x, y, 1.8, 0.6, color, name, 10, WHITE)

# ==================== P12-18: 其余页面（简化版，保持一致性）====================
# P12: 实施路径
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.333, 7.5, WHITE)
add_rect(slide, 0, 0, 13.333, 0.12, TELECOM_RED)
add_header(slide, "实施路径")
add_target_line(slide, "5个月完成建设，四阶段渐进式实施。")
add_measure_title(slide)
items = [
    ("第一阶段", "2个月，智能问答+查询，基础能力建设，准确率≥90%。"),
    ("第二阶段", "2个月，智能办理，核心价值实现，时长缩短≥30%。"),
    ("第三阶段", "1个月，智能诊断，完善闭环，准确率≥80%。"),
    ("第四阶段", "持续优化，越用越聪明，场景扩展。")
]
add_measure_content(slide, items)
add_diagram_area(slide)
add_text(slide, 7.0, 2.0, 5, 0.4, "实施时间轴", 14, True, DARK_RED)
# 时间轴
phases = ["第一阶段", "第二阶段", "第三阶段", "第四阶段"]
colors = [TELECOM_RED, ORANGE, BLUE, DARK_GRAY]
x = 7.0
for i, (phase, color) in enumerate(zip(phases, colors)):
    add_box(slide, x, 2.5, 1.3, 0.6, color, phase, 9, WHITE)
    if len(icons) > 80 + i:
        add_icon(slide, x + 0.4, 3.2, 0.5, 0.5, icons[80 + i])
    x += 1.5

# P13: 关键指标
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.333, 7.5, WHITE)
add_rect(slide, 0, 0, 13.333, 0.12, TELECOM_RED)
add_header(slide, "关键指标")
add_target_line(slide, "效果可量化，业务+技术双维度指标体系。")
add_measure_title(slide)
items = [
    ("效率指标", "办理时长缩短40%，查询效率提升50%，诊断时间缩短60%。"),
    ("质量指标", "问答准确率≥90%，填单准确率≥98%，诊断准确率≥85%。"),
    ("体验指标", "用户满意度≥85%，响应时间<3秒。")
]
add_measure_content(slide, items)
add_diagram_area(slide)
add_text(slide, 7.0, 2.0, 5, 0.4, "关键指标", 14, True, DARK_RED)
# 指标展示
add_text(slide, 7.0, 2.5, 5, 0.3, "效率指标", 12, True, DARK_RED)
metrics1 = [("40%", "时长缩短"), ("50%", "效率提升"), ("60%", "时间缩短")]
x = 7.0
for i, (value, label) in enumerate(metrics1):
    add_box(slide, x, 2.9, 1.3, 1.1, TELECOM_RED, value, 16, WHITE)
    add_text(slide, x, 3.6, 1.3, 0.25, label, 8, False, DARK_GRAY, PP_ALIGN.CENTER)
    if len(icons) > 90 + i:
        add_icon(slide, x + 0.4, 4.1, 0.4, 0.4, icons[90 + i])
    x += 1.5

# P14: 核心价值
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.333, 7.5, WHITE)
add_rect(slide, 0, 0, 13.333, 0.12, TELECOM_RED)
add_header(slide, "核心价值")
add_target_line(slide, "四大价值维度，考核+效率+体验+人力。")
add_measure_title(slide)
items = [
    ("考核达标", "满足集团考核要求，占指标10%，31省标杆项目。"),
    ("效率提升", "办理时长缩短40%，查询效率提升50%，诊断时间缩短60%。"),
    ("体验改善", "客户满意度提升，一线满意度提升，服务响应加快。"),
    ("人力释放", "推动一线员工转型，从事高价值任务，减负增效。")
]
add_measure_content(slide, items)
add_diagram_area(slide)
add_text(slide, 7.0, 2.0, 5, 0.4, "核心价值", 14, True, DARK_RED)
values = ["考核达标", "效率提升", "体验改善", "人力释放"]
colors = [TELECOM_RED, ORANGE, BLUE, DARK_GRAY]
y = 2.5
for i, (title, color) in enumerate(zip(values, colors)):
    add_box(slide, 7.0, y, 5.0, 0.7, color, title, 13, WHITE)
    if len(icons) > 100 + i:
        add_icon(slide, 11.5, y + 0.1, 0.5, 0.5, icons[100 + i])
    y += 0.85

# P15: 创新亮点
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.333, 7.5, WHITE)
add_rect(slide, 0, 0, 13.333, 0.12, TELECOM_RED)
add_header(slide, "创新亮点")
add_target_line(slide, "引领行业AI应用，五大创新技术融合。")
add_measure_title(slide)
items = [
    ("智能体集群", "四大智能体分工协作，专业精准，智能问答、查询、办理、诊断。"),
    ("多模态交互", "语音+文字+图片多种输入方式，灵活便捷，降低使用门槛。"),
    ("主动推荐", "客户画像分析，精准推荐，千人千面，提升营销效率。"),
    ("学习型系统", "持续学习，越用越聪明，自我进化，适应业务变化。")
]
add_measure_content(slide, items)
add_diagram_area(slide)
add_text(slide, 7.0, 2.0, 5, 0.4, "创新亮点", 14, True, DARK_RED)
innovations = ["智能体集群", "多模态交互", "主动推荐", "学习型系统", "场景化卡片"]
colors = [TELECOM_RED, ORANGE, BLUE, LIGHT_BLUE, DARK_GRAY]
y = 2.5
for i, (title, color) in enumerate(zip(innovations, colors)):
    add_box(slide, 7.0, y, 5.0, 0.65, color, title, 11, WHITE)
    if len(icons) > 110 + i:
        add_icon(slide, 11.5, y + 0.1, 0.4, 0.4, icons[110 + i])
    y += 0.75

# P16: 成功案例
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.333, 7.5, WHITE)
add_rect(slide, 0, 0, 13.333, 0.12, TELECOM_RED)
add_header(slide, "成功案例")
add_target_line(slide, "三地成功实践，复制成熟经验，降低风险。")
add_measure_title(slide)
items = [
    ("陕西电信", "营业办理数字员工，办理时长缩短40%，一线效率显著提升。"),
    ("重庆电信", "营业受理数字员工，智能问答+查询，满足考核要求。"),
    ("安徽电信", "CRM域AI应用，四大能力全覆盖，31省标杆。")
]
add_measure_content(slide, items)
add_diagram_area(slide)
add_text(slide, 7.0, 2.0, 5, 0.4, "成功案例", 14, True, DARK_RED)
cases = [("陕西电信", "时长缩短40%"), ("重庆电信", "满足考核要求"), ("安徽电信", "31省标杆")]
colors = [TELECOM_RED, ORANGE, BLUE]
y = 2.5
for i, ((title, desc), color) in enumerate(zip(cases, colors)):
    add_box(slide, 7.0, y, 5.0, 0.9, color, "", 12, WHITE)
    add_text(slide, 7.2, y + 0.15, 2.5, 0.3, title, 13, True, WHITE)
    add_text(slide, 7.2, y + 0.5, 4.5, 0.3, desc, 11, False, WHITE)
    if len(icons) > 120 + i:
        add_icon(slide, 11.5, y + 0.2, 0.5, 0.5, icons[120 + i])
    y += 1.0

# P17: 总结与展望
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.333, 7.5, WHITE)
add_rect(slide, 0, 0, 13.333, 0.12, TELECOM_RED)
add_header(slide, "总结与展望")
add_target_line(slide, "方案成熟可落地，分阶段实施，风险可控。")
add_measure_title(slide)
items = [
    ("方案成熟", "参考三地成功案例，技术方案验证可行，实施方法论完善。"),
    ("技术可行", "大模型+RAG技术成熟，开源组件丰富，技术风险低。"),
    ("价值明确", "效率提升40%以上，满足考核要求，客户体验改善。"),
    ("后续规划", "持续优化迭代，扩展更多场景，复制推广到其他省份。")
]
add_measure_content(slide, items)
add_diagram_area(slide)
add_text(slide, 7.0, 2.0, 5, 0.4, "总结", 14, True, DARK_RED)
summaries = ["方案成熟", "技术可行", "价值明确", "风险可控"]
colors = [TELECOM_RED, ORANGE, BLUE, DARK_GRAY]
y = 2.5
for i, (title, color) in enumerate(zip(summaries, colors)):
    add_box(slide, 7.0, y, 5.0, 0.7, color, title, 13, WHITE)
    if len(icons) > 130 + i:
        add_icon(slide, 11.5, y + 0.1, 0.5, 0.5, icons[130 + i])
    y += 0.85

# P18: 结束页
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.333, 7.5, WHITE)
add_rect(slide, 0, 0, 13.333, 0.12, TELECOM_RED)
add_text(slide, 1, 2.5, 11.333, 1.0, "谢谢！", 48, True, DARK_RED, PP_ALIGN.CENTER)
add_text(slide, 1, 3.5, 11.333, 0.6, "让我们一起，用AI赋能一线，提升效率，改善体验！", 18, False, DARK_GRAY, PP_ALIGN.CENTER)
add_text(slide, 1, 4.5, 11.333, 0.4, "南有乔木，不可休思", 14, False, GRAY, PP_ALIGN.CENTER)

# 保存
output_path = "/root/.openclaw/workspace/03_项目管理/营业受理数字员工/03_工程输出/指南针工程/营业受理数字员工建设方案_专业版.pptx"
prs.save(output_path)
print(f"\n✅ 专业版PPT已生成: {output_path}")
print(f"总页数: {len(prs.slides)}")
print(f"嵌入图标: {len(icons)}个可用")
print(f"嵌入架构图组件: {len(logics)}个可用")
